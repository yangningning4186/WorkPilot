from pathlib import Path
from zipfile import ZipFile

import pymupdf
import pytest
from docx import Document
from docx.oxml.ns import qn
from openpyxl import load_workbook
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.util import Inches
from pydantic import ValidationError

from app.cowork.artifact_manifest import ArtifactSkillRef, bind_claim_evidence
from app.cowork.artifact_pipeline import render_validate_commit
from app.cowork.artifact_renderers import render_candidate
from app.cowork.artifact_renderers.contracts import (
    ClaimRecord,
    DocumentBlock,
    DocumentSection,
    DocumentSpec,
    HtmlReportSpec,
    HtmlSection,
    PdfSpec,
    PresentationSpec,
    SlideSpec,
    WorkbookCell,
    WorkbookSpec,
    WorksheetSpec,
)
from app.cowork.artifact_renderers.image_assets import ArtifactImageError, sanitized_svg_bytes
from app.cowork.artifact_validation import validate_artifact
from app.cowork.files import CoworkFileError
from app.cowork.skills.builtin.pptx.scripts.pptx2image import (
    _font,
    _missing_characters,
    render_presentation_pages,
)
from app.cowork.skills.builtin.pptx.scripts.render_pptx import render_presentation
from app.cowork.skills.builtin.pptx.scripts.visual_kits import (
    verify_visual_kit_asset,
    visual_kit_ids,
)
from eval.artifact_bench import score_artifact


def _skill(name: str) -> ArtifactSkillRef:
    return ArtifactSkillRef(name=name, origin="builtin", kind="artifact", sha256="a" * 64)


def test_pptx_visual_kit_catalog_locks_all_copied_template_assets() -> None:
    kit_ids = visual_kit_ids()

    assert len(kit_ids) == 17
    assert "workpilot-clean" in kit_ids
    assert sum(verify_visual_kit_asset(name) is not None for name in kit_ids) == 16


def test_format_specific_renderers_produce_reopenable_artifacts(tmp_path: Path) -> None:
    presentation = PresentationSpec(
        title="可信交付",
        slides=[
            SlideSpec(id="s1", role="hero", layout="title", title="生成不是终点"),
            SlideSpec(
                id="s2",
                layout="two_column",
                title="生产线",
                left_title="模型",
                left_items=["内容", "结构"],
                right_title="Runtime",
                right_items=["Renderer", "Validator"],
                notes="来源：WorkPilot 固定 Renderer 设计说明。",
            ),
        ],
    )
    document = DocumentSpec(
        title="研究报告",
        sections=[
            DocumentSection(
                id="summary",
                heading="摘要",
                blocks=[DocumentBlock(id="p1", type="paragraph", text="有据可查。")],
            )
        ],
    )
    workbook = WorkbookSpec(
        title="指标",
        sheets=[
            WorksheetSpec(
                name="Summary",
                cells=[
                    WorkbookCell(address="A1", value="指标", style="header"),
                    WorkbookCell(address="A2", value=42, style="metric"),
                ],
            )
        ],
    )
    html = HtmlReportSpec(
        title="离线报告",
        sections=[
            HtmlSection(
                id="summary",
                heading="摘要",
                blocks=[DocumentBlock(type="paragraph", text="无远程依赖。")],
            )
        ],
    )
    pdf = PdfSpec(
        title="归档报告",
        summary="固定渲染的 PDF。",
        sections=[
            HtmlSection(
                id="summary",
                heading="摘要",
                blocks=[DocumentBlock(type="paragraph", text="文件可重新打开。")],
            )
        ],
    )
    cases = [
        (presentation, tmp_path / "deck.pptx"),
        (document, tmp_path / "report.docx"),
        (workbook, tmp_path / "metrics.xlsx"),
        (html, tmp_path / "report.html"),
        (pdf, tmp_path / "report.pdf"),
    ]

    for spec, path in cases:
        render_candidate(spec, path)
        report = validate_artifact(
            path,
            spec=spec,
            render_visual=path.suffix.casefold() == ".pptx",
        )
        assert report.deliverable is True
        assert report.structural.status == "passed"

    rendered_presentation = Presentation(str(tmp_path / "deck.pptx"))
    assert len(rendered_presentation.slides) == 2
    assert "固定 Renderer" in rendered_presentation.slides[1].notes_slide.notes_text_frame.text
    assert Document(str(tmp_path / "report.docx")).paragraphs[0].text == "研究报告"


def test_required_evidence_failure_never_commits_candidate(tmp_path: Path) -> None:
    target = tmp_path / "report.docx"
    spec = DocumentSpec(
        title="证据报告",
        evidence_policy="required",
        claims=[
            ClaimRecord(
                claim_id="c1",
                text="关键结论",
                evidence_ids=["S1"],
                target_type="docx_paragraph",
                target_id="p1",
            )
        ],
        sections=[
            DocumentSection(
                id="summary",
                heading="摘要",
                blocks=[DocumentBlock(id="p1", type="paragraph", text="关键结论")],
            )
        ],
    )
    bindings = bind_claim_evidence(spec, [])

    with pytest.raises(CoworkFileError, match="验证失败"):
        render_validate_commit(
            spec=spec,
            target=target,
            baseline_sha256=None,
            skill=_skill("docx"),
            evidence_bindings=bindings,
            max_bytes=5 * 1024 * 1024,
            backup_versions=3,
        )

    assert not target.exists()


def test_manifest_binds_ppt_slide_to_existing_evidence(tmp_path: Path) -> None:
    target = tmp_path / "deck.pptx"
    spec = PresentationSpec(
        title="证据演示",
        evidence_policy="required",
        claims=[
            ClaimRecord(
                claim_id="c1",
                text="证据驱动交付",
                evidence_ids=["S1"],
                target_type="pptx_slide",
                target_id="s1",
            )
        ],
        slides=[SlideSpec(id="s1", layout="title", title="证据驱动交付")],
    )
    bindings = bind_claim_evidence(
        spec,
        [
            {
                "citation_id": "S1",
                "kind": "knowledge",
                "title": "论文",
                "source_uri": "/workspace/paper.pdf",
                "quote": "原文证据",
                "locator": 7,
                "locations": [],
            }
        ],
    )

    result = render_validate_commit(
        spec=spec,
        target=target,
        baseline_sha256=None,
        skill=_skill("pptx"),
        evidence_bindings=bindings,
        max_bytes=5 * 1024 * 1024,
        backup_versions=3,
    )

    assert result.manifest.status == "validated"
    assert result.manifest.evidence_bindings[0].target_id == "s1"
    assert result.manifest.evidence_bindings[0].evidence[0].locator == 7
    _, metrics = score_artifact(target, spec=spec, evidence_bindings=bindings)
    assert metrics.file_open_rate == 1.0
    assert metrics.evidence_coverage == 1.0
    assert metrics.citation_accuracy == 1.0


def test_pptx_visual_validation_is_required_and_renders_every_slide(tmp_path: Path) -> None:
    target = tmp_path / "visual.pptx"
    spec = PresentationSpec(
        title="逐页验证",
        slides=[
            SlideSpec(id="s1", layout="title", title="第一页"),
            SlideSpec(id="s2", layout="statement", title="第二页", body="一个清楚的结论"),
        ],
    )
    render_candidate(spec, target)

    without_render = validate_artifact(target, spec=spec, render_visual=False)
    assert without_render.deliverable is False
    assert without_render.visual.status == "not_run"

    pages = render_presentation_pages(target, tmp_path / "pages")
    assert len(pages.pages) == 2
    assert all(page.is_file() and page.stat().st_size > 0 for page in pages.pages)
    assert not pages.overflow_shapes
    assert not pages.unsupported_shapes
    with_render = validate_artifact(target, spec=spec, render_visual=True)
    assert with_render.deliverable is True
    assert with_render.visual.status == "passed"


def test_pptx_rasterizer_selects_a_font_with_bold_chinese_glyphs() -> None:
    text = "月满中秋 · 情满校园"
    resolved = _font(text, 35, bold=True, pixels_per_emu=1600 / 12_191_695)

    assert not _missing_characters(resolved.font, text)


def test_pptx_runtime_is_owned_by_the_builtin_skill_package() -> None:
    from app.cowork.artifact_renderers.pptx_renderer import (
        render_presentation as compatibility_renderer,
    )
    from app.cowork.pptx_rasterizer import (
        render_presentation_pages as compatibility_rasterizer,
    )

    assert compatibility_renderer is render_presentation
    assert compatibility_rasterizer is render_presentation_pages
    assert render_presentation.__module__.endswith(
        "skills.builtin.pptx.scripts.render_pptx"
    )
    assert render_presentation_pages.__module__.endswith(
        "skills.builtin.pptx.scripts.pptx2image"
    )


def test_safe_svg_is_embedded_in_pptx_and_docx_with_alt_text(tmp_path: Path) -> None:
    svg = tmp_path / "architecture.svg"
    svg.write_text(
        """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 800 480">
        <rect width="800" height="480" rx="32" fill="#edf5f1"/>
        <path d="M120 240 H680" stroke="#167a5b" stroke-width="18"/>
        <circle cx="210" cy="240" r="70" fill="#ffffff" stroke="#167a5b" stroke-width="12"/>
        <circle cx="590" cy="240" r="70" fill="#ffffff" stroke="#167a5b" stroke-width="12"/>
        </svg>""",
        encoding="utf-8",
    )
    pptx_target = tmp_path / "visual.pptx"
    pptx_spec = PresentationSpec(
        title="视觉素材",
        slides=[
            SlideSpec(
                id="s1",
                layout="image_text",
                title="安全 SVG",
                image_path=str(svg),
                image_alt="两个节点通过一条绿色连线连接",
                image_caption="系统关系示意",
                bullets=["保留源 SVG", "交付时嵌入兼容图片"],
            )
        ],
    )
    render_candidate(pptx_spec, pptx_target)
    presentation = Presentation(pptx_target)
    pictures = [shape for shape in presentation.slides[0].shapes if shape.shape_type == 13]
    assert len(pictures) == 1
    assert pictures[0]._pic.nvPicPr.cNvPr.get("descr") == "两个节点通过一条绿色连线连接"
    with ZipFile(pptx_target) as archive:
        assert any(
            name.startswith("ppt/media/") and name.endswith(".png") for name in archive.namelist()
        )
    pages = render_presentation_pages(pptx_target, tmp_path / "svg-pages")
    assert not pages.overflow_shapes
    assert not pages.unsupported_shapes

    docx_target = tmp_path / "visual.docx"
    docx_spec = DocumentSpec(
        title="图文报告",
        sections=[
            DocumentSection(
                id="diagram",
                heading="系统关系",
                blocks=[
                    DocumentBlock(
                        type="image",
                        image_path=str(svg),
                        image_alt="两个节点通过一条绿色连线连接",
                        image_caption="图 1 系统关系示意",
                    )
                ],
            )
        ],
    )
    render_candidate(docx_spec, docx_target)
    document = Document(docx_target)
    assert len(document.inline_shapes) == 1
    properties = document.element.body.xpath(".//wp:docPr")
    assert properties[0].get("descr") == "两个节点通过一条绿色连线连接"
    assert any(paragraph.text == "图 1 系统关系示意" for paragraph in document.paragraphs)
    with ZipFile(docx_target) as archive:
        assert any(
            name.startswith("word/media/") and name.endswith(".png") for name in archive.namelist()
        )


def test_pptx_cover_visual_cards_and_activity_render_as_native_layouts(tmp_path: Path) -> None:
    hero = tmp_path / "mid-autumn-hero.svg"
    hero.write_text(
        """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 600 900">
        <rect width="600" height="900" fill="#192746"/>
        <circle cx="300" cy="300" r="170" fill="#f2cf78"/>
        <rect x="170" y="570" width="260" height="170" rx="34" fill="#c9892d"/>
        <path d="M300 470 V820" stroke="#f2cf78" stroke-width="14"/>
        </svg>""",
        encoding="utf-8",
    )
    target = tmp_path / "facilitation.pptx"
    spec = PresentationSpec.model_validate(
        {
            "title": "中秋主题班会",
            "theme": {
                "name": "mid-autumn",
                "background": "F8F1E5",
                "surface": "FFFDF8",
                "text_primary": "2B241D",
                "text_secondary": "665A4A",
                "accent": "C9892D",
                "positive": "567A61",
                "warning": "A64D3D",
                "title_font": "Arial",
                "body_font": "Arial",
                "east_asia_font": "Microsoft YaHei",
            },
            "slides": [
                {
                    "id": "cover",
                    "role": "hero",
                    "rhythm": "peak",
                    "layout": "title",
                    "title": "月满中秋 · 情满班级",
                    "subtitle": "从文化理解走向共同参与",
                    "image_path": str(hero),
                    "image_alt": "金色圆月与传统花灯的国风矢量主视觉",
                },
                {
                    "id": "culture",
                    "layout": "cards",
                    "title": "四种习俗都在回答同一个问题：如何表达团圆",
                    "cards": [
                        {"kicker": "观察", "title": "赏月", "detail": "从圆月意象理解团圆与思念。"},
                        {"kicker": "品味", "title": "月饼", "detail": "从分享食物理解共同体与祝福。"},
                        {"kicker": "合作", "title": "灯谜", "detail": "用协作猜谜让文化知识进入体验。"},
                        {"kicker": "表达", "title": "花灯", "detail": "把祝愿写成可展示、可分享的作品。"},
                    ],
                },
                {
                    "id": "activity",
                    "layout": "activity",
                    "title": "三分钟共创：把祝福变成班级行动",
                    "activity_prompt": "如果“团圆”不只发生在家里，我们能为班级做一件什么事？",
                    "activity_steps": [
                        "四人一组，每人先写下一个具体行动。",
                        "合并相似想法，选出最容易本周完成的一项。",
                        "由一位同学用 20 秒说明行动、负责人和完成时间。",
                    ],
                    "activity_timebox": "3 分钟",
                    "activity_debrief": "哪一个行动最能让平时容易被忽略的同学感到被看见？",
                },
            ],
        }
    )

    render_candidate(spec, target)
    report = validate_artifact(target, spec=spec, render_visual=True)
    presentation = Presentation(target)

    assert report.deliverable is True, report.model_dump(mode="json")
    assert any(shape.shape_type == 13 for shape in presentation.slides[0].shapes)
    assert "赏月" in "\n".join(shape.text for shape in presentation.slides[1].shapes if getattr(shape, "has_text_frame", False))
    assert "3 分钟" in "\n".join(shape.text for shape in presentation.slides[2].shapes if getattr(shape, "has_text_frame", False))


def test_pptx_chart_layout_is_an_editable_native_chart(tmp_path: Path) -> None:
    target = tmp_path / "native-chart.pptx"
    spec = PresentationSpec.model_validate(
        {
            "title": "季度趋势",
            "visual_kit": "consulting-02",
            "slides": [
                {
                    "id": "cover",
                    "role": "hero",
                    "rhythm": "peak",
                    "layout": "title",
                    "title": "收入增长正在恢复",
                    "subtitle": "2026 年季度复盘",
                },
                {
                    "id": "trend",
                    "layout": "chart",
                    "title": "Q3 收入增速转正，Q4 仍需验证续约质量",
                    "chart": {
                        "chart_type": "line",
                        "categories": ["Q1", "Q2", "Q3", "Q4E"],
                        "series": [
                            {"name": "收入（百万元）", "values": [82, 79, 91, 104]},
                            {"name": "目标（百万元）", "values": [86, 88, 94, 108]},
                        ],
                        "unit": "百万元",
                    },
                    "body": "趋势已反转，但预计值仍低于目标，下一步优先关闭续约缺口。",
                    "notes": "数据为 Renderer 回归样例，不代表真实经营数据。",
                },
            ],
        }
    )

    render_candidate(spec, target)
    report = validate_artifact(target, spec=spec, render_visual=True)
    presentation = Presentation(target)

    assert report.deliverable is True, report.model_dump(mode="json")
    assert any(
        shape.shape_type == MSO_SHAPE_TYPE.CHART
        for shape in presentation.slides[1].shapes
    )
    assert (
        "Renderer 回归样例"
        in presentation.slides[1].notes_slide.notes_text_frame.text
    )


def test_pptx_diagrams_and_constrained_canvas_render_as_editable_native_shapes(
    tmp_path: Path,
) -> None:
    target = tmp_path / "diagrams-and-canvas.pptx"
    spec = PresentationSpec.model_validate(
        {
            "title": "图示组件与安全画布",
            "slides": [
                {
                    "id": "process",
                    "layout": "diagram",
                    "title": "从问题到落地形成闭环",
                    "diagram": {
                        "kind": "process",
                        "nodes": [
                            {"id": "observe", "title": "识别问题", "detail": "统一事实", "emphasis": "primary"},
                            {"id": "design", "title": "形成方案", "detail": "比较路径"},
                            {"id": "pilot", "title": "小步试点", "detail": "验证假设"},
                            {"id": "scale", "title": "规模推广", "detail": "沉淀标准"},
                        ],
                    },
                },
                {
                    "id": "cycle",
                    "layout": "diagram",
                    "title": "持续改进依靠四步循环",
                    "diagram": {
                        "kind": "cycle",
                        "center_label": "持续改进",
                        "nodes": [
                            {"id": "observe", "title": "观察", "detail": "收集反馈"},
                            {"id": "judge", "title": "判断", "detail": "识别偏差"},
                            {"id": "act", "title": "行动", "detail": "实施调整", "emphasis": "primary"},
                            {"id": "review", "title": "复盘", "detail": "沉淀方法"},
                        ],
                    },
                },
                {
                    "id": "hierarchy",
                    "layout": "diagram",
                    "title": "三级责任结构让决策快速落地",
                    "diagram": {
                        "kind": "hierarchy",
                        "nodes": [
                            {"id": "root", "title": "项目委员会", "detail": "方向与资源", "emphasis": "primary"},
                            {"id": "product", "title": "产品组", "detail": "方案与范围"},
                            {"id": "delivery", "title": "交付组", "detail": "计划与质量"},
                            {"id": "research", "title": "用户研究"},
                            {"id": "engineering", "title": "工程实施"},
                            {"id": "qa", "title": "质量保障"},
                        ],
                        "edges": [
                            {"source": "root", "target": "product"},
                            {"source": "root", "target": "delivery"},
                            {"source": "product", "target": "research"},
                            {"source": "delivery", "target": "engineering"},
                            {"source": "delivery", "target": "qa"},
                        ],
                    },
                },
                {
                    "id": "funnel",
                    "layout": "diagram",
                    "title": "线索经过四层筛选转化为客户",
                    "diagram": {
                        "kind": "funnel",
                        "nodes": [
                            {"id": "lead", "title": "全部线索", "detail": "2,400 条"},
                            {"id": "qualified", "title": "有效需求", "detail": "860 条"},
                            {"id": "review", "title": "方案评审", "detail": "260 条"},
                            {"id": "won", "title": "签约客户", "detail": "72 家", "emphasis": "primary"},
                        ],
                    },
                },
                {
                    "id": "pyramid",
                    "layout": "diagram",
                    "title": "能力建设从基础走向规模化",
                    "diagram": {
                        "kind": "pyramid",
                        "nodes": [
                            {"id": "strategy", "title": "战略牵引", "detail": "明确优先级", "emphasis": "primary"},
                            {"id": "mechanism", "title": "机制协同", "detail": "打通职责"},
                            {"id": "data", "title": "数据闭环", "detail": "持续反馈"},
                            {"id": "foundation", "title": "基础设施", "detail": "平台与标准"},
                        ],
                    },
                },
                {
                    "id": "canvas",
                    "layout": "canvas",
                    "title": "受约束画布表达非标准关系",
                    "canvas": {
                        "elements": [
                            {"type": "shape", "id": "input", "x": 2, "y": 24, "width": 24, "height": 42, "title": "输入", "detail": "事实与需求"},
                            {"type": "shape", "id": "engine", "x": 38, "y": 12, "width": 24, "height": 64, "shape": "hexagon", "title": "决策引擎", "detail": "规则 + 判断", "fill_role": "accent", "fill_style": "solid"},
                            {"type": "shape", "id": "output", "x": 74, "y": 24, "width": 24, "height": 42, "title": "输出", "detail": "行动与结果", "fill_role": "positive"},
                            {"type": "connector", "id": "edge-in", "source_id": "input", "target_id": "engine", "label": "校验"},
                            {"type": "connector", "id": "edge-out", "source_id": "engine", "target_id": "output", "label": "执行", "color_role": "positive"},
                        ]
                    },
                },
            ],
        }
    )

    render_candidate(spec, target)
    pages = render_presentation_pages(target, tmp_path / "diagram-pages")
    report = validate_artifact(target, spec=spec, render_visual=True)
    presentation = Presentation(target)

    assert report.deliverable is True, report.model_dump(mode="json")
    assert not pages.overflow_shapes
    assert not pages.unsupported_shapes
    assert any(
        shape.shape_type == MSO_SHAPE_TYPE.LINE
        or (
            getattr(getattr(shape._element, "spPr", None), "prstGeom", None) is not None
            and shape._element.spPr.prstGeom.get("prst", "").startswith(
                ("line", "bentConnector")
            )
        )
        for slide in presentation.slides
        for shape in slide.shapes
    )
    assert any(
        shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and "决策引擎" in shape.text
        for shape in presentation.slides[-1].shapes
    )
    visual_story = next(
        check for check in report.semantic.checks if check.name == "visual_story"
    )
    assert visual_story.status == "passed"


@pytest.mark.parametrize(
    ("slide", "message"),
    [
        (
            {
                "id": "tree",
                "layout": "diagram",
                "title": "错误层级",
                "diagram": {
                    "kind": "hierarchy",
                    "nodes": [
                        {"id": "a", "title": "A"},
                        {"id": "b", "title": "B"},
                        {"id": "c", "title": "C"},
                    ],
                    "edges": [
                        {"source": "a", "target": "b"},
                        {"source": "b", "target": "a"},
                    ],
                },
            },
            "必须从根节点可达",
        ),
        (
            {
                "id": "outside",
                "layout": "canvas",
                "title": "越界画布",
                "canvas": {
                    "elements": [
                        {"type": "shape", "id": "a", "x": 80, "y": 10, "width": 30, "height": 20, "title": "越界"},
                        {"type": "shape", "id": "b", "x": 5, "y": 50, "width": 30, "height": 20, "title": "正常"},
                    ]
                },
            },
            "完整位于",
        ),
        (
            {
                "id": "overlap",
                "layout": "canvas",
                "title": "重叠画布",
                "canvas": {
                    "elements": [
                        {"type": "shape", "id": "a", "x": 10, "y": 10, "width": 45, "height": 45, "title": "A"},
                        {"type": "shape", "id": "b", "x": 30, "y": 25, "width": 45, "height": 45, "title": "B"},
                    ]
                },
            },
            "大面积重叠",
        ),
        (
            {
                "id": "dense-text",
                "layout": "canvas",
                "title": "文字容量",
                "canvas": {
                    "elements": [
                        {"type": "text", "id": "a", "x": 5, "y": 5, "width": 12, "height": 12, "font_size": 36, "text": "这是一段远远超过小边界盒容量的文字" * 4},
                        {"type": "shape", "id": "b", "x": 55, "y": 40, "width": 30, "height": 30, "title": "另一区块"},
                    ]
                },
            },
            "文字超过当前边界容量",
        ),
        (
            {
                "id": "dangling-edge",
                "layout": "canvas",
                "title": "悬空连接线",
                "canvas": {
                    "elements": [
                        {"type": "shape", "id": "a", "x": 5, "y": 10, "width": 30, "height": 30, "title": "A"},
                        {"type": "shape", "id": "b", "x": 65, "y": 10, "width": 30, "height": 30, "title": "B"},
                        {"type": "connector", "id": "edge", "source_id": "a", "target_id": "missing"},
                    ]
                },
            },
            "必须连接两个不同的现有可见元素",
        ),
        (
            {
                "id": "crossing-edge",
                "layout": "canvas",
                "title": "连接线穿过节点",
                "canvas": {
                    "elements": [
                        {"type": "shape", "id": "a", "x": 2, "y": 38, "width": 20, "height": 20, "title": "A"},
                        {"type": "shape", "id": "blocker", "x": 40, "y": 34, "width": 20, "height": 28, "title": "中间节点"},
                        {"type": "shape", "id": "b", "x": 78, "y": 38, "width": 20, "height": 20, "title": "B"},
                        {"type": "connector", "id": "edge", "source_id": "a", "target_id": "b"},
                    ]
                },
            },
            "穿过 element blocker",
        ),
    ],
)
def test_pptx_advanced_layout_contracts_reject_unsafe_geometry(
    slide: dict[str, object],
    message: str,
) -> None:
    with pytest.raises(ValidationError, match=message):
        PresentationSpec.model_validate(
            {"artifact_type": "pptx", "title": "安全边界", "slides": [slide]}
        )


def test_svg_with_active_or_external_content_is_rejected(tmp_path: Path) -> None:
    unsafe = tmp_path / "unsafe.svg"
    unsafe.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg"><script>alert(1)</script></svg>',
        encoding="utf-8",
    )

    with pytest.raises(ArtifactImageError, match="脚本"):
        sanitized_svg_bytes(unsafe)

    remote = tmp_path / "remote.svg"
    remote.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg"><rect fill="url(https://evil.example/a)"/></svg>',
        encoding="utf-8",
    )
    with pytest.raises(ArtifactImageError, match="外部资源"):
        sanitized_svg_bytes(remote)


def test_pptx_text_overflow_never_commits(tmp_path: Path) -> None:
    target = tmp_path / "overflow.pptx"
    spec = PresentationSpec(
        title="溢出门禁",
        slides=[
            SlideSpec(
                id="s1",
                layout="two_column",
                title="这是一个故意写得非常长、会在普通页标题安全区域里发生换行并溢出的标题" * 4,
                left_title="左侧",
                left_items=["内容"],
                right_title="右侧",
                right_items=["内容"],
            )
        ],
    )

    with pytest.raises(CoworkFileError, match="文本溢出"):
        render_validate_commit(
            spec=spec,
            target=target,
            baseline_sha256=None,
            skill=_skill("pptx"),
            evidence_bindings=[],
            max_bytes=5 * 1024 * 1024,
            backup_versions=3,
        )

    assert not target.exists()


def test_presentation_contract_rejects_matrix_beyond_fixed_layout_capacity() -> None:
    with pytest.raises(ValidationError):
        SlideSpec.model_validate(
            {
                "id": "matrix",
                "layout": "matrix",
                "title": "容量",
                "matrix": [
                    {"x": str(index), "y": str(index), "label": str(index)} for index in range(5)
                ],
            }
        )


@pytest.mark.parametrize(
    ("payload", "message"),
    [
        (
            {
                "id": "columns",
                "layout": "two_column",
                "title": "两栏",
                "left_title": "左",
                "right_title": "右",
            },
            "左右非空条目",
        ),
        ({"id": "metric", "layout": "big_number", "title": "指标"}, "metrics"),
        ({"id": "chart", "layout": "chart", "title": "趋势"}, "必须提供 chart"),
        (
            {"id": "image", "layout": "image_text", "title": "图片", "body": "解读"},
            "image_path",
        ),
        ({"id": "quote", "layout": "quote", "title": "引文"}, "必须提供 body"),
        (
            {
                "id": "timeline",
                "layout": "timeline",
                "title": "计划",
                "timeline": [{"label": "现在", "title": "开始"}],
            },
            "至少包含 2 个节点",
        ),
        (
            {
                "id": "matrix",
                "layout": "matrix",
                "title": "矩阵",
                "matrix": [{"x": "高", "y": "高", "label": "优先"}],
            },
            "至少包含 2 个对象",
        ),
        (
            {
                "id": "hidden",
                "layout": "quote",
                "title": "引文",
                "body": "可见引文",
                "bullets": ["不会显示"],
            },
            "不消费字段",
        ),
        (
            {
                "id": "hidden-title-body",
                "layout": "title",
                "title": "封面",
                "subtitle": "可见副标题",
                "body": "不会显示",
            },
            "title 不能同时提供 subtitle 与 body",
        ),
        (
            {
                "id": "hidden-image-body",
                "layout": "image_text",
                "title": "图片解读",
                "image_path": "/tmp/example.png",
                "bullets": ["可见解释"],
                "body": "不会显示",
            },
            "image_text 不能同时提供 bullets 与 body",
        ),
    ],
)
def test_presentation_contract_rejects_empty_or_hidden_layout_content(
    payload: dict[str, object], message: str
) -> None:
    with pytest.raises(ValidationError, match=message):
        SlideSpec.model_validate(payload)


def test_pptx_rich_native_layouts_render_without_sparse_supporting_pages(
    tmp_path: Path,
) -> None:
    target = tmp_path / "rich-layouts.pptx"
    spec = PresentationSpec(
        title="经营复盘",
        visual_kit="consulting-02",
        slides=[
            SlideSpec(
                id="cover",
                role="hero",
                rhythm="peak",
                layout="title",
                title="增长恢复，但续约仍是下一阶段约束",
                subtitle="2026 年 8 月经营复盘",
            ),
            SlideSpec(
                id="drivers",
                layout="two_column",
                title="增长由新签拉动，续约缺口集中在华南",
                left_title="增长驱动",
                left_items=["新签企业客户 7 家", "工单积压环比下降 26.9%", "重点客户上线按期"],
                right_title="主要约束",
                right_items=["续约率低于目标 2 个百分点", "3 家客户待安全评审", "华南流失原因待复盘"],
            ),
            SlideSpec.model_validate(
                {
                    "id": "metrics",
                    "layout": "big_number",
                    "title": "收入接近目标，风险尚未解除",
                    "metrics": [
                        {"value": "95.4%", "label": "收入目标达成率", "detail": "截至 2026-08-31"},
                        {"value": "91%", "label": "续约率", "detail": "目标 93%"},
                        {"value": "38", "label": "积压工单", "detail": "上期 52 单"},
                    ],
                    "body": "扩张条件基本具备，但需先关闭续约与安全评审风险。",
                }
            ),
            SlideSpec.model_validate(
                {
                    "id": "plan",
                    "layout": "timeline",
                    "title": "未来 30 天按三个检查点关闭风险",
                    "timeline": [
                        {"label": "09-10", "title": "上线计划", "detail": "逐户确认 3 家客户"},
                        {"label": "09-15", "title": "续约复盘", "detail": "完成华南流失分析"},
                        {"label": "09-30", "title": "经营复核", "detail": "检查续约与上线结果"},
                    ],
                }
            ),
            SlideSpec.model_validate(
                {
                    "id": "matrix",
                    "layout": "matrix",
                    "title": "资源优先投向高影响、可快速关闭的事项",
                    "matrix": [
                        {"x": "高影响", "y": "高可控", "label": "安全评审"},
                        {"x": "高影响", "y": "低可控", "label": "华南续约"},
                        {"x": "低影响", "y": "高可控", "label": "工单积压"},
                        {"x": "低影响", "y": "低可控", "label": "长尾需求"},
                    ],
                }
            ),
        ],
    )

    render_candidate(spec, target)
    report = validate_artifact(target, spec=spec, render_visual=True)
    presentation = Presentation(target)
    cover_title = next(
        shape
        for shape in presentation.slides[0].shapes
        if getattr(shape, "has_text_frame", False)
        and shape.text == "增长恢复，但续约仍是下一阶段约束"
    )

    assert report.deliverable is True
    assert cover_title.text_frame.paragraphs[0].runs[0].font.size.pt == 39
    assert next(
        check for check in report.semantic.checks if check.name == "supporting_density"
    ).status == "passed"
    assert next(
        check for check in report.semantic.checks if check.name == "visual_story"
    ).status == "passed"


def test_pptx_focus_layouts_render_titles_and_messages_without_overlap(
    tmp_path: Path,
) -> None:
    target = tmp_path / "focus-layouts.pptx"
    spec = PresentationSpec(
        title="焦点页字段可见性",
        slides=[
            SlideSpec(
                id="statement",
                role="hero",
                rhythm="peak",
                layout="statement",
                title="本轮判断",
                body="增长已经恢复，但续约仍是下一阶段的首要约束。",
                subtitle="结论来自 2026 年 8 月经营数据",
            ),
            SlideSpec(
                id="section",
                role="transition",
                layout="section",
                title="第二部分",
                body="从结果转向约束与行动",
                subtitle="先关闭续约缺口，再扩大增长投入",
            ),
            SlideSpec(
                id="quote",
                layout="quote",
                title="客户原声",
                body="上线速度不错，但我们仍需要更明确的安全评审计划。",
                quote_attribution="重点客户访谈",
            ),
        ],
    )

    render_candidate(spec, target)
    report = validate_artifact(target, spec=spec, render_visual=True)
    presentation = Presentation(target)
    slide_text = [
        "\n".join(shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False))
        for slide in presentation.slides
    ]

    assert report.deliverable is True
    assert "本轮判断" in slide_text[0]
    assert "增长已经恢复" in slide_text[0]
    assert "第二部分" in slide_text[1]
    assert "从结果转向约束与行动" in slide_text[1]
    assert "客户原声" in slide_text[2]
    assert "安全评审计划" in slide_text[2]


def test_pptx_sparse_supporting_page_fails_delivery_even_when_not_blank(
    tmp_path: Path,
) -> None:
    target = tmp_path / "sparse.pptx"
    spec = PresentationSpec(
        title="空洞门禁",
        slides=[
            SlideSpec(
                id="thin",
                layout="two_column",
                title="只有两个短词不构成完整内容页",
                left_title="左",
                left_items=["一条"],
                right_title="右",
                right_items=["一条"],
            )
        ],
    )

    render_candidate(spec, target)
    report = validate_artifact(target, spec=spec, render_visual=True)

    assert report.deliverable is False
    density = next(
        check for check in report.semantic.checks if check.name == "supporting_density"
    )
    assert density.status == "failed"
    assert "thin" in density.message


def test_pptx_output_density_catches_legacy_title_and_body_page(tmp_path: Path) -> None:
    target = tmp_path / "legacy-sparse.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.5), Inches(0.8))
    title.text_frame.text = "中秋习俗一览"
    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.5), Inches(4.5))
    body.text_frame.text = "赏月、吃月饼、猜灯谜、玩花灯"
    presentation.save(target)

    report = validate_artifact(target, render_visual=False)

    density = next(
        check for check in report.visual.checks if check.name == "content_density"
    )
    assert density.status == "failed"
    assert density.value == 1
    assert report.deliverable is False


def test_pptx_output_density_catches_large_label_only_containers(tmp_path: Path) -> None:
    target = tmp_path / "hollow-cards.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    presentation.slides.add_slide(presentation.slide_layouts[6])
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(8), Inches(0.7))
    title.text_frame.text = "中秋习俗一览"
    positions = ((0.8, 1.4), (6.8, 1.4), (0.8, 4.2), (6.8, 4.2))
    for (left, top), label in zip(positions, ("赏月", "月饼", "灯谜", "花灯"), strict=True):
        slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(5.35),
            Inches(2.25),
        )
        text = slide.shapes.add_textbox(
            Inches(left + 0.35),
            Inches(top + 0.75),
            Inches(4.65),
            Inches(0.5),
        )
        text.text_frame.text = label
    presentation.save(target)

    report = validate_artifact(target, render_visual=False)
    hollow = next(check for check in report.visual.checks if check.name == "hollow_containers")

    assert hollow.status == "failed"
    assert hollow.value == 1


def test_pptx_output_density_catches_title_plus_four_short_labels(tmp_path: Path) -> None:
    target = tmp_path / "label-grid.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    presentation.slides.add_slide(presentation.slide_layouts[6])
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(8), Inches(0.7))
    title.text_frame.text = "中秋习俗：在家与在外、吃与玩"
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.2),
        Inches(1.7),
        Inches(10.7),
        Inches(4.7),
    )
    positions = ((1.6, 2.1), (7.1, 2.1), (1.6, 4.5), (7.1, 4.5))
    for (left, top), label in zip(
        positions,
        ("吃月饼", "赏月·拜月", "饮桂花酒", "玩花灯·猜灯谜"),
        strict=True,
    ):
        text = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(4.2), Inches(0.7)
        )
        text.text_frame.text = label
    presentation.save(target)

    report = validate_artifact(target, render_visual=False)
    density = next(
        check for check in report.visual.checks if check.name == "content_density"
    )

    assert density.status == "failed"
    assert density.value == 1


def test_pptx_multiple_metric_pages_are_not_misclassified_as_focus(tmp_path: Path) -> None:
    target = tmp_path / "metric-evidence.pptx"
    slides = [
        SlideSpec(
            id="cover",
            role="hero",
            rhythm="peak",
            layout="title",
            title="指标证据页",
            subtitle="经营复盘",
        )
    ]
    for index in range(4):
        slides.append(
            SlideSpec.model_validate(
                {
                    "id": f"metric-{index}",
                    "layout": "big_number",
                    "title": f"第 {index + 1} 组经营指标",
                    "metrics": [
                        {"value": f"{90 + index}%", "label": "目标达成率"},
                        {"value": str(30 + index), "label": "待办数量"},
                    ],
                }
            )
        )
    spec = PresentationSpec(title="指标证据页", slides=slides)

    render_candidate(spec, target)
    report = validate_artifact(target, spec=spec, render_visual=False)

    inferred = next(
        check for check in report.visual.checks if check.name == "inferred_focus_ratio"
    )
    density = next(
        check for check in report.visual.checks if check.name == "content_density"
    )
    assert inferred.status == "passed"
    assert inferred.value == 1
    assert density.status == "passed"


def test_docx_uses_safe_page_geometry_and_fixed_table_grid(tmp_path: Path) -> None:
    target = tmp_path / "table.docx"
    spec = DocumentSpec(
        title="表格报告",
        sections=[
            DocumentSection(
                id="summary",
                heading="摘要",
                blocks=[
                    DocumentBlock(
                        type="table",
                        headers=["项目", "说明"],
                        rows=[["A", "包含较长中文说明的内容"], ["B", "可自动分配列宽"]],
                    )
                ],
            )
        ],
    )
    render_candidate(spec, target)
    document = Document(target)
    section = document.sections[0]
    assert section.page_width.inches == pytest.approx(8.5, abs=0.01)
    assert section.left_margin.inches == pytest.approx(1.0, abs=0.01)
    layout = document.tables[0]._tbl.tblPr.find(qn("w:tblLayout"))
    assert layout is not None and layout.get(qn("w:type")) == "fixed"
    assert not document.tables[0]._tbl.xpath(".//w:trHeight")
    assert validate_artifact(target, spec=spec).structural.status == "passed"


def test_xlsx_wraps_long_cjk_text_and_avoids_clipping(tmp_path: Path) -> None:
    target = tmp_path / "wrapped.xlsx"
    spec = WorkbookSpec(
        title="长文本",
        sheets=[
            WorksheetSpec(
                name="Summary",
                cells=[WorkbookCell(address="A1", value="很长的中文内容" * 20)],
            )
        ],
    )
    render_candidate(spec, target)
    workbook = load_workbook(target)
    try:
        sheet = workbook["Summary"]
        assert sheet.column_dimensions["A"].width == pytest.approx(48.0)
        assert sheet["A1"].alignment.wrap_text is True
        assert float(sheet.row_dimensions[1].height or 0) > 15
    finally:
        workbook.close()
    report = validate_artifact(target, spec=spec)
    assert (
        next(check for check in report.semantic.checks if check.name == "potential_clipping").status
        == "passed"
    )


def test_pdf_renderer_paginates_and_rasterizes_every_page(tmp_path: Path) -> None:
    target = tmp_path / "multipage.pdf"
    spec = PdfSpec(
        title="多页报告",
        sections=[
            HtmlSection(
                id=f"section-{index}",
                heading=f"章节 {index}",
                blocks=[DocumentBlock(type="paragraph", text="用于验证自然分页的正文。" * 80)],
            )
            for index in range(12)
        ],
    )
    render_candidate(spec, target)
    document = pymupdf.open(target)
    try:
        assert document.page_count > 1
    finally:
        document.close()
    report = validate_artifact(target, spec=spec)
    assert report.deliverable is True
    assert report.visual.status == "passed"
