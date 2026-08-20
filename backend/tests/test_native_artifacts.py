import hashlib
import stat
from pathlib import Path

import fitz  # type: ignore[import-untyped]
import pytest
from docx import Document
from openpyxl import load_workbook  # type: ignore[import-untyped]
from pptx import Presentation  # type: ignore[import-untyped]

from app.cowork.native_artifacts import create_native_artifact


def test_native_artifact_generation_docx_xlsx_pdf(tmp_path: Path) -> None:
    docx_path = tmp_path / "report.docx"
    xlsx_path = tmp_path / "table.xlsx"
    pdf_path = tmp_path / "brief.pdf"

    docx_result = create_native_artifact(
        docx_path,
        format="docx",
        title="季度报告",
        content="# 结论\n- 收入增长\n- 风险可控",
        sheets=[],
        baseline_sha256=None,
    )
    create_native_artifact(
        xlsx_path,
        format="xlsx",
        title="数据",
        content="",
        sheets=[{"name": "汇总", "rows": [["项目", "金额"], ["收入", 100]]}],
        baseline_sha256=None,
    )
    create_native_artifact(
        pdf_path,
        format="pdf",
        title="一页摘要",
        content="这是可直接预览的原生 PDF。",
        sheets=[],
        baseline_sha256=None,
    )

    assert docx_result.sha256
    assert Document(str(docx_path)).core_properties.title == "季度报告"
    workbook = load_workbook(xlsx_path, read_only=True)
    try:
        assert workbook["汇总"]["B2"].value == 100
    finally:
        workbook.close()
    with fitz.open(pdf_path) as pdf:
        assert pdf.page_count == 1
        assert pdf.metadata["title"] == "一页摘要"

    baseline = hashlib.sha256(docx_path.read_bytes()).hexdigest()
    overwritten = create_native_artifact(
        docx_path,
        format="docx",
        title="修订报告",
        content="新正文",
        sheets=[],
        baseline_sha256=baseline,
        backup_versions=1,
    )
    assert overwritten.backup_path is not None
    assert overwritten.backup_path.exists()


def test_native_pptx_is_editable_widescreen_deck(tmp_path: Path) -> None:
    pptx_path = tmp_path / "儿童节.pptx"
    result = create_native_artifact(
        pptx_path,
        format="pptx",
        title="快乐六一",
        content="",
        sheets=[],
        slides=[
            {"title": "节日由来", "bullets": ["关爱儿童", "快乐成长"]},
            {"title": "活动安排", "body": "游戏、表演与分享"},
        ],
        baseline_sha256=None,
    )

    presentation = Presentation(str(pptx_path))
    assert result.mime_type.endswith("presentationml.presentation")
    assert len(presentation.slides) == 3
    assert presentation.slide_width / presentation.slide_height == pytest.approx(16 / 9)
    assert "快乐六一" in " ".join(
        shape.text for shape in presentation.slides[0].shapes if hasattr(shape, "text")
    )


def test_native_pdf_paginates_cjk_and_overwrite_preserves_mode(tmp_path: Path) -> None:
    pdf_path = tmp_path / "long.pdf"
    tail = "末尾不可丢失"
    create_native_artifact(
        pdf_path,
        format="pdf",
        title="中文长文",
        content="正文" * 4_000 + tail,
        sheets=[],
        baseline_sha256=None,
    )
    with fitz.open(pdf_path) as pdf:
        extracted = "".join(page.get_text() for page in pdf)
        assert pdf.page_count > 1
        assert tail in extracted

    pdf_path.chmod(0o640)
    baseline = hashlib.sha256(pdf_path.read_bytes()).hexdigest()
    create_native_artifact(
        pdf_path,
        format="pdf",
        title="修订",
        content="修订正文",
        sheets=[],
        baseline_sha256=baseline,
    )
    assert stat.S_IMODE(pdf_path.stat().st_mode) == 0o640

    with pytest.raises(ValueError, match="必须省略"):
        create_native_artifact(
            tmp_path / "missing.pdf",
            format="pdf",
            title="不存在",
            content="",
            sheets=[],
            baseline_sha256="0" * 64,
        )


def test_native_artifact_rejects_oversized_existing_baseline(tmp_path: Path) -> None:
    path = tmp_path / "large.docx"
    with path.open("wb") as stream:
        stream.truncate(2 * 1024 * 1024)

    with pytest.raises(ValueError, match="超过读取上限"):
        create_native_artifact(
            path,
            format="docx",
            title="Large",
            content="replacement",
            sheets=[],
            baseline_sha256="0" * 64,
            max_existing_bytes=1024,
        )


def test_native_pdf_renders_markdown_as_document_layout(tmp_path: Path) -> None:
    pdf_path = tmp_path / "markdown.pdf"
    create_native_artifact(
        pdf_path,
        format="pdf",
        title="MCP 总结",
        content=(
            "# MCP 内容标题\n\n"
            "## 核心结论\n\n"
            "**MCP** 用于连接模型与工具。\n\n"
            "- 支持资源\n- 支持工具\n\n"
            "| 角色 | 收益 |\n| --- | --- |\n| 开发者 | 降低集成成本 |\n\n"
            "```text\nHost -> Client -> Server\n```"
        ),
        sheets=[],
        baseline_sha256=None,
    )

    with fitz.open(pdf_path) as pdf:
        extracted = "\n".join(page.get_text() for page in pdf)
        assert pdf.metadata["title"] == "MCP 总结"
        assert extracted.count("MCP 总结") == 1
        assert "MCP 内容标题" not in extracted
        assert "核心结论" in extracted
        assert "降低集成成本" in extracted
        assert "Host -> Client -> Server" in extracted
        assert "**MCP**" not in extracted
        assert "| --- |" not in extracted
