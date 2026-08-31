from __future__ import annotations

import asyncio
import hashlib
import json
from pathlib import Path

import pymupdf
import pytest
from docx import Document
from openpyxl import Workbook  # type: ignore[import-untyped]
from openpyxl.chart import BarChart, Reference  # type: ignore[import-untyped]
from PIL import Image
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.presentation import Presentation as PresentationDocument
from pptx.util import Inches, Pt
from reportlab.pdfbase import pdfmetrics  # type: ignore[import-untyped]
from reportlab.pdfbase.cidfonts import UnicodeCIDFont  # type: ignore[import-untyped]
from reportlab.pdfgen import canvas  # type: ignore[import-untyped]

from app.cowork.artifact_validation import validate_artifact
from eval.office_content_suite import (
    DEFAULT_SUITE,
    OfficeContentItem,
    OfficeContentSuite,
    OfficeContentSuiteError,
    PptxChartDataCheck,
    ReviewFile,
    ScopeSpec,
    TextClaimContainsCheck,
    TextClaimNotContainsCheck,
    TextContainsCheck,
    TextLabeledValuesCheck,
    XlsxCellsValueCheck,
    XlsxFormulaValueCheck,
    _main,
    evaluate_check,
    evaluate_suite,
    extract_artifact,
    load_suite,
    prepare_suite,
    suite_summary,
)
from eval.office_eval import _configured_judge, _parse_args, run_one_click_evaluation
from eval.office_model_judge import (
    RenderedArtifact,
    parse_model_response,
    render_artifact_for_review,
    run_model_reviews,
)
from workpilot_ai.providers.gemini import GeminiProvider
from workpilot_ai.types import CacheRetention, CompletionResult, Message, Usage


def test_seed_suite_has_balanced_formats_and_content_contracts() -> None:
    suite = load_suite(DEFAULT_SUITE)

    assert suite_summary(suite) == {
        "name": "office-content-seed",
        "version": "1.1.0",
        "origin": "synthetic",
        "data_classification": "synthetic",
        "review_status": "pending_human_review",
        "items": 12,
        "splits": {"dev": 8, "test": 4},
        "formats": {"docx": 3, "pdf": 3, "pptx": 3, "xlsx": 3},
        "task_types": {
            "analyze": 6,
            "conflict_handling": 2,
            "create": 1,
            "transform": 3,
        },
        "difficulties": {"hard": 7, "medium": 5},
        "automatic_checks": 99,
        "penalties": 12,
        "review_criteria": 36,
    }
    assert all(
        {rubric.dimension for rubric in item.rubric}
        == {"fundamentals", "completeness", "correctness", "fidelity", "usability"}
        for item in suite.items
    )
    assert all(rubric.source_refs for item in suite.items for rubric in item.rubric)
    assert all(
        {criterion.dimension for criterion in item.review_criteria}
        == {"coherence", "decision_utility", "visual_quality"}
        for item in suite.items
    )
    assert all(
        criterion.minimum_score > 0 for item in suite.items for criterion in item.review_criteria
    )
    assert all(
        item.gate.render_visual
        and {"structural", "visual", "security"}.issubset(item.gate.fail_on_dimensions)
        and {"structural", "visual", "security"}.issubset(item.gate.require_measured_dimensions)
        for item in suite.items
        if item.artifact_type == "pptx"
    )


def test_prepare_defaults_to_dev_and_audits_test_access(tmp_path: Path) -> None:
    suite = load_suite(DEFAULT_SUITE)
    dev_root = tmp_path / "dev"

    manifest = prepare_suite(suite, dev_root)

    assert len(manifest["items"]) == 8
    assert len(manifest["suite_sha256"]) == 64
    assert all(
        item["id"] not in {value.id for value in suite.items if value.split == "test"}
        for item in manifest["items"]
    )
    first = suite.items[0]
    assert (dev_root / first.id / "TASK.md").is_file()
    assert (dev_root / first.id / "inputs" / first.fixtures[0].path).is_file()
    assert (dev_root / first.id / "submission").is_dir()
    with pytest.raises(OfficeContentSuiteError, match="include_test"):
        prepare_suite(suite, tmp_path / "test-without-audit", split="test")
    with pytest.raises(OfficeContentSuiteError, match="test_access_note"):
        prepare_suite(
            suite,
            tmp_path / "test-without-note",
            split="test",
            include_test=True,
        )

    test_manifest = prepare_suite(
        suite,
        tmp_path / "test",
        split="test",
        include_test=True,
        test_access_note="冻结验收",
    )

    assert len(test_manifest["items"]) == 4
    assert test_manifest["test_access_note"] == "冻结验收"


def _save_formula_workbook(path: Path) -> None:
    workbook = Workbook()
    data = workbook.active
    data.title = "明细"
    data["A1"] = "value"
    data["A2"] = 10
    data["A3"] = 15
    data["B2"] = "华东"
    data["B3"] = "华南"
    summary = workbook.create_sheet("汇总")
    summary["A1"] = "total"
    summary["B1"] = "=SUM(明细!A2:A3)"
    summary["A2"] = "direct"
    summary["B2"] = "=明细!A2+明细!A3"
    summary["A3"] = "east"
    summary["B3"] = '=SUMIF(明细!B2:B3,"华东",明细!A2:A3)'
    summary["B1"].number_format = "¥#,##0"
    workbook.save(path)


def test_xlsx_formula_checks_accept_semantically_equivalent_calculations(tmp_path: Path) -> None:
    path = tmp_path / "formula.xlsx"
    _save_formula_workbook(path)
    view = extract_artifact(path, "xlsx")
    validation = validate_artifact(path, render_visual=False)

    sum_outcome = evaluate_check(
        XlsxFormulaValueCheck(
            type="xlsx_formula_value",
            address="汇总!B1",
            expected=25,
            tolerance=0,
        ),
        view,
        validation,
    )
    direct_outcome = evaluate_check(
        XlsxFormulaValueCheck(
            type="xlsx_formula_value",
            address="汇总!B2",
            expected=25,
            tolerance=0,
        ),
        view,
        validation,
    )
    batch_outcome = evaluate_check(
        XlsxCellsValueCheck(
            type="xlsx_cells_value",
            cells={"汇总!B1": 25, "汇总!B2": 25},
            tolerance=0,
        ),
        view,
        validation,
    )
    sumif_outcome = evaluate_check(
        XlsxFormulaValueCheck(
            type="xlsx_formula_value",
            address="汇总!B3",
            expected=10,
            tolerance=0,
        ),
        view,
        validation,
    )

    assert sum_outcome.passed
    assert direct_outcome.passed
    assert batch_outcome.passed
    assert sumif_outcome.passed
    assert "→ 25" in sum_outcome.evidence[0]


def _save_scoped_presentation(path: Path) -> None:
    presentation = Presentation()
    for title, body in (("经营结果", "ARR 860，计划 900"), ("主要风险", "API 准备度黄灯")):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = title
        text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
        text_box.text_frame.text = body
    presentation.save(str(path))


def test_slide_scope_prevents_global_keyword_stuffing(tmp_path: Path) -> None:
    path = tmp_path / "scoped.pptx"
    _save_scoped_presentation(path)
    view = extract_artifact(path, "pptx")
    validation = validate_artifact(path, render_visual=False)

    global_check = TextContainsCheck(
        type="text_contains",
        values=["860"],
        scope=ScopeSpec(kind="artifact"),
    )
    wrong_slide_check = TextContainsCheck(
        type="text_contains",
        values=["860"],
        scope=ScopeSpec(kind="pptx_slide", selector="风险"),
    )

    assert evaluate_check(global_check, view, validation).passed
    assert not evaluate_check(wrong_slide_check, view, validation).passed


def test_slide_scope_excludes_speaker_notes(tmp_path: Path) -> None:
    path = tmp_path / "notes-only.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "经营结果"
    notes_frame = slide.notes_slide.notes_text_frame
    assert notes_frame is not None
    notes_frame.text = "ARR 860，计划 900"
    presentation.save(str(path))
    view = extract_artifact(path, "pptx")
    validation = validate_artifact(path, render_visual=False)

    outcome = evaluate_check(
        TextContainsCheck(
            type="text_contains",
            values=["860", "900"],
            scope=ScopeSpec(kind="pptx_slide", selector="经营结果"),
        ),
        view,
        validation,
    )

    assert "860" not in view.text
    assert not outcome.passed


@pytest.mark.parametrize(
    ("body", "expected"),
    [
        ("ARR 860，计划 900", True),
        ("ARR 900，计划 860", False),
        ("ARR 8600，计划 900", False),
    ],
)
def test_labeled_values_bind_numbers_to_their_business_meaning(
    tmp_path: Path,
    body: str,
    expected: bool,
) -> None:
    path = tmp_path / f"relation-{expected}.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "经营结果"
    text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    text_box.text_frame.text = body
    presentation.save(str(path))
    view = extract_artifact(path, "pptx")
    validation = validate_artifact(path, render_visual=False)
    check = TextLabeledValuesCheck.model_validate(
        {
            "type": "text_labeled_values",
            "requirements": [
                {
                    "labels": ["ARR"],
                    "values": ["860"],
                    "direction": "after",
                    "max_distance": 30,
                },
                {
                    "labels": ["计划"],
                    "values": ["900"],
                    "direction": "after",
                    "max_distance": 30,
                },
            ],
            "scope": {"kind": "pptx_slide", "selector": "经营结果"},
        }
    )

    assert evaluate_check(check, view, validation).passed is expected


def test_pptx_chart_check_reads_native_data_and_percentage_fractions(
    tmp_path: Path,
) -> None:
    path = tmp_path / "native-chart.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "试点证据"
    data = ChartData()
    data.categories = ["下限", "上限"]
    data.add_series("试点工时下降", (0.18, 0.24))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(1.8),
        Inches(8),
        Inches(4.5),
        data,
    )
    presentation.save(str(path))
    view = extract_artifact(path, "pptx")
    validation = validate_artifact(path, render_visual=False)

    expected = evaluate_check(
        PptxChartDataCheck(
            type="pptx_chart_data",
            required_values=[18, 24],
            required_labels=["工时"],
        ),
        view,
        validation,
    )
    wrong = evaluate_check(
        PptxChartDataCheck(
            type="pptx_chart_data",
            required_values=[15, 21],
            required_labels=["工时"],
        ),
        view,
        validation,
    )

    assert expected.passed
    assert not wrong.passed


def test_claim_checks_do_not_penalize_explicit_negation(tmp_path: Path) -> None:
    path = tmp_path / "negated.docx"
    document = Document()
    document.add_heading("证据边界", level=1)
    document.add_paragraph("该试点尚未证明生产率提升，例外也不是自动批准。")
    document.save(str(path))
    view = extract_artifact(path, "docx")
    validation = validate_artifact(path, render_visual=False)

    penalty_trigger = evaluate_check(
        TextClaimContainsCheck(
            type="text_claim_contains",
            values=["生产率提升", "自动批准"],
            scope=ScopeSpec(kind="artifact"),
        ),
        view,
        validation,
    )
    fidelity = evaluate_check(
        TextClaimNotContainsCheck(
            type="text_claim_not_contains",
            values=["生产率提升", "自动批准"],
            scope=ScopeSpec(kind="artifact"),
        ),
        view,
        validation,
    )

    assert not penalty_trigger.passed
    assert fidelity.passed


def test_pdf_extractor_localizes_evidence_by_page(tmp_path: Path) -> None:
    path = tmp_path / "pages.pdf"
    document = pymupdf.open()
    document.new_page().insert_text((72, 72), "Approved impact: 14%")
    document.new_page().insert_text((72, 72), "Next update: 2026-09-02")
    document.save(path)
    document.close()
    view = extract_artifact(path, "pdf")
    validation = validate_artifact(path, render_visual=False)

    outcome = evaluate_check(
        TextContainsCheck(
            type="text_contains",
            values=["2026-09-02"],
            scope=ScopeSpec(kind="pdf_page", index=2),
        ),
        view,
        validation,
    )

    assert view.counts["pdf_page_count"] == 2
    assert outcome.passed
    assert outcome.evidence[0].startswith("PDF page 2")


def _seed_submission_path(root: Path, item_id: str, output_file: str) -> Path:
    path = root / item_id / "submission" / output_file
    path.parent.mkdir(parents=True)
    return path


def _save_seed_docx(path: Path) -> None:
    document = Document()
    document.add_heading("执行摘要", level=1)
    document.add_paragraph(
        "数据截至 2026 年 8 月 31 日。收入 1,240 万元，低于 1,300 万元目标，达成率 95.4%。"
    )
    document.add_heading("KPI", level=1)
    table = document.add_table(rows=5, cols=4)
    for cell, value in zip(table.rows[0].cells, ("指标", "实际", "目标/对比", "说明"), strict=True):
        cell.text = value
    values = (
        ("收入", "1,240 万元", "1,300 万元", "达成率 95.4%"),
        ("续约率", "91%", "93%", "低于目标"),
        ("工单积压", "38", "上季度 52", "下降 26.9%"),
        ("新签企业客户", "7", "3 家待评审", "需推进安全评审"),
    )
    for row, row_values in zip(table.rows[1:], values, strict=True):
        for cell, value in zip(row.cells, row_values, strict=True):
            cell.text = value
    document.add_heading("风险", level=1)
    document.add_paragraph("安全评审排期可能拖延 3 家客户上线；华南续约率偏低。")
    document.add_heading("行动", level=1)
    document.add_paragraph("李雯于 2026-09-10 前给出上线计划。")
    document.add_paragraph("周成于 2026-09-15 前完成流失客户复盘。")
    document.save(str(path))


def _add_seed_slide(presentation: PresentationDocument, title: str, body: str) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = RGBColor(250, 247, 241)

    title_box = slide.shapes.add_textbox(Inches(0.85), Inches(0.55), Inches(8.7), Inches(0.75))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    title_run = title_frame.paragraphs[0].add_run()
    title_run.text = title
    title_run.font.name = "Aptos Display"
    title_run.font.size = Pt(25)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(47, 43, 37)

    eyebrow = slide.shapes.add_textbox(Inches(0.88), Inches(1.35), Inches(2.0), Inches(0.45))
    eyebrow_frame = eyebrow.text_frame
    eyebrow_frame.clear()
    eyebrow_frame.margin_top = 0
    eyebrow_frame.margin_bottom = 0
    eyebrow_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    eyebrow_run = eyebrow_frame.paragraphs[0].add_run()
    eyebrow_run.text = "关键事实"
    eyebrow_run.font.name = "Aptos"
    eyebrow_run.font.size = Pt(16)
    eyebrow_run.font.bold = True
    eyebrow_run.font.color.rgb = RGBColor(184, 116, 35)

    for index, line in enumerate(value for value in body.splitlines() if value.strip()):
        top = 1.9 + index * 1.25
        badge = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.9),
            Inches(top + 0.18),
            Inches(0.5),
            Inches(0.5),
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(217, 147, 55)
        badge.line.fill.background()
        badge_frame = badge.text_frame
        badge_frame.clear()
        badge_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        badge_paragraph = badge_frame.paragraphs[0]
        badge_paragraph.alignment = PP_ALIGN.CENTER
        badge_run = badge_paragraph.add_run()
        badge_run.text = f"{index + 1:02d}"
        badge_run.font.name = "Aptos"
        badge_run.font.size = Pt(16)
        badge_run.font.bold = True
        badge_run.font.color.rgb = RGBColor(255, 255, 255)

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.6),
            Inches(top),
            Inches(7.8),
            Inches(0.88),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = RGBColor(231, 222, 208)
        card_frame = card.text_frame
        card_frame.clear()
        card_frame.margin_left = Inches(0.25)
        card_frame.margin_right = Inches(0.25)
        card_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        card_run = card_frame.paragraphs[0].add_run()
        card_run.text = line.strip()
        card_run.font.name = "Aptos"
        card_run.font.size = Pt(18)
        card_run.font.color.rgb = RGBColor(60, 55, 48)


def _save_seed_pptx(path: Path) -> None:
    presentation = Presentation()
    _add_seed_slide(presentation, "Northstar 月度经营复盘", "2026 年 8 月")
    _add_seed_slide(
        presentation,
        "经营结果｜ARR 与流失率均未达目标",
        "ARR 860 万元 / 计划 900 万元\n流失率 2.8% / 目标 2.2%\nAPI P95 延迟下降 18%",
    )
    _add_seed_slide(
        presentation,
        "客户与交付｜转化仍待推进",
        "企业试点 7 家，预计转付费 3 家\n12 个里程碑完成 10 个",
    )
    _add_seed_slide(
        presentation,
        "主要风险｜API 准备度黄灯",
        "API 准备度黄灯\n两个延期里程碑共同依赖数据迁移",
    )
    _add_seed_slide(
        presentation,
        "未来 30 天｜锁定复核与恢复计划",
        "王蕾：9 月 9 日复核 API 准备度\n刘川：9 月 12 日给出数据迁移恢复计划",
    )
    presentation.save(str(path))


def _save_seed_xlsx(path: Path) -> None:
    workbook = Workbook()
    detail = workbook.active
    detail.title = "明细"
    detail.append(["订单号", "区域", "产品", "数量", "单价", "销售额"])
    rows = (
        ("O-1001", "华东", "A", 10, 120),
        ("O-1002", "华南", "B", 8, 250),
        ("O-1003", "华东", "B", 5, 250),
        ("O-1004", "华北", "A", 20, 120),
        ("O-1005", "华南", "C", 3, 800),
    )
    for row_index, row in enumerate(rows, start=2):
        detail.append([*row, f"=D{row_index}*E{row_index}"])
        detail[f"F{row_index}"].number_format = "¥#,##0"
    for column in "ABCDEF":
        detail.column_dimensions[column].width = 16
    summary = workbook.create_sheet("汇总")
    summary.append(["指标", "数值"])
    summary.append(["总销售额", "=SUM(明细!F2:F6)"])
    summary.append(["华东", '=SUMIF(明细!B2:B6,"华东",明细!F2:F6)'])
    summary.append(["华南", '=SUMIF(明细!B2:B6,"华南",明细!F2:F6)'])
    summary.append(["华北", '=SUMIF(明细!B2:B6,"华北",明细!F2:F6)'])
    for row_index in range(2, 6):
        summary[f"B{row_index}"].number_format = "¥#,##0"
    chart = BarChart()
    chart.add_data(Reference(summary, min_col=2, min_row=3, max_row=5))
    chart.set_categories(Reference(summary, min_col=1, min_row=3, max_row=5))
    summary.add_chart(chart, "D2")
    summary.column_dimensions["A"].width = 16
    summary.column_dimensions["B"].width = 16
    workbook.save(path)


def _save_seed_pdf(path: Path) -> None:
    pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
    writer = canvas.Canvas(str(path))
    writer.setFont("STSong-Light", 15)
    writer.drawString(72, 780, "影响")
    writer.setFont("STSong-Light", 11)
    writer.drawString(72, 750, "2026-08-29，约 14% 的 API 请求经历延迟升高。")
    writer.drawString(72, 725, "未发现数据丢失。")
    writer.setFont("STSong-Light", 15)
    writer.drawString(72, 680, "时间线")
    writer.setFont("STSong-Light", 11)
    writer.drawString(72, 650, "09:12—10:06 UTC 出现服务降级，共 54 分钟。")
    writer.showPage()
    writer.setFont("STSong-Light", 15)
    writer.drawString(72, 780, "处置")
    writer.setFont("STSong-Light", 11)
    writer.drawString(72, 750, "团队通过回滚最近一次配置变更恢复服务。")
    writer.setFont("STSong-Light", 15)
    writer.drawString(72, 700, "后续更新")
    writer.setFont("STSong-Light", 11)
    writer.drawString(72, 670, "下一次进展更新：2026-09-02。")
    writer.save()


def test_seed_representative_oracles_pass_each_format(tmp_path: Path) -> None:
    suite = load_suite(DEFAULT_SUITE)
    representative_ids = {
        "office-docx-001-quarterly-brief",
        "office-pptx-001-monthly-review",
        "office-xlsx-001-sales-summary",
        "office-pdf-001-external-incident",
    }
    representative_suite = suite.model_copy(
        update={"items": [item for item in suite.items if item.id in representative_ids]}
    )
    submission_root = tmp_path / "representatives"
    _save_seed_docx(
        _seed_submission_path(
            submission_root,
            "office-docx-001-quarterly-brief",
            "q3-business-brief.docx",
        )
    )
    _save_seed_pptx(
        _seed_submission_path(
            submission_root,
            "office-pptx-001-monthly-review",
            "northstar-monthly-review.pptx",
        )
    )
    _save_seed_xlsx(
        _seed_submission_path(
            submission_root,
            "office-xlsx-001-sales-summary",
            "sales-summary.xlsx",
        )
    )
    _save_seed_pdf(
        _seed_submission_path(
            submission_root,
            "office-pdf-001-external-incident",
            "external-incident-summary.pdf",
        )
    )

    report = evaluate_suite(
        representative_suite,
        submission_root,
        tmp_path / "representative-report",
    )

    gate_failures = {
        result["id"]: result["gate"]["reasons"]
        for result in report["results"]
        if not result["gate"]["passed"]
    }
    assert report["summary"]["gate_passed"] == 4, gate_failures
    assert report["summary"]["automatic_passed"] == 4
    assert all(item["automatic_score"] == 100 for item in report["results"])
    assert len(report["suite_sha256"]) == 64
    assert len(report["scorer_fingerprint"]) == 64
    assert report["generated_at"].endswith("+00:00")


def _calibration_suite() -> OfficeContentSuite:
    return OfficeContentSuite.model_validate(
        {
            "schema_version": "workpilot-office-content-suite.v1",
            "name": "office-content-calibration",
            "version": "1.0.0",
            "origin": "synthetic",
            "data_classification": "synthetic",
            "review_status": "pending_human_review",
            "reviewer": None,
            "reviewed_at": None,
            "methodology": ["unit calibration"],
            "dimension_weights": {
                "fundamentals": 0.2,
                "completeness": 0.2,
                "correctness": 0.2,
                "fidelity": 0.2,
                "usability": 0.2,
            },
            "automatic_weight": 0.6,
            "review_weight": 0.4,
            "items": [
                {
                    "id": "calibration-docx",
                    "split": "dev",
                    "artifact_type": "docx",
                    "task_type": "analyze",
                    "category": "calibration",
                    "difficulty": "easy",
                    "prompt": "制作包含摘要、KPI、风险和行动的经营文档，并保留正确的数据关系。",
                    "output_file": "brief.docx",
                    "fixtures": [{"path": "source.txt", "content": "ARR 860 / 目标 900；API 黄灯"}],
                    "gate": {
                        "render_visual": False,
                        "fail_on_dimensions": ["structural", "security"],
                        "require_measured_dimensions": ["structural", "security"],
                        "min_validator_quality": 0,
                    },
                    "rubric": [
                        {
                            "id": "fund-structure",
                            "dimension": "fundamentals",
                            "description": "至少四个标题",
                            "source_refs": ["prompt"],
                            "critical": True,
                            "check": {
                                "type": "structure_count",
                                "metric": "docx_heading_count",
                                "operator": "gte",
                                "value": 4,
                            },
                        },
                        {
                            "id": "complete-risk",
                            "dimension": "completeness",
                            "description": "风险段包含 API 黄灯",
                            "source_refs": ["source.txt"],
                            "critical": True,
                            "check": {
                                "type": "text_contains",
                                "values": ["API 黄灯"],
                                "scope": {"kind": "docx_section", "selector": "风险"},
                            },
                        },
                        {
                            "id": "correct-kpi",
                            "dimension": "correctness",
                            "description": "KPI 表保留实际与目标",
                            "source_refs": ["source.txt"],
                            "critical": True,
                            "check": {
                                "type": "text_contains",
                                "values": ["860", "900"],
                                "scope": {"kind": "docx_table", "index": 1},
                            },
                        },
                        {
                            "id": "faithful-no-overclaim",
                            "dimension": "fidelity",
                            "description": "不声称超额完成",
                            "source_refs": ["source.txt"],
                            "critical": True,
                            "check": {
                                "type": "text_not_contains",
                                "values": ["超额完成"],
                                "scope": {"kind": "artifact"},
                            },
                        },
                        {
                            "id": "usable-order",
                            "dimension": "usability",
                            "description": "决策阅读顺序正确",
                            "source_refs": ["prompt"],
                            "check": {
                                "type": "text_ordered",
                                "values": ["执行摘要", "KPI", "风险", "行动"],
                                "scope": {"kind": "artifact"},
                            },
                        },
                    ],
                    "penalties": [
                        {
                            "id": "overclaim",
                            "description": "错误声称超额完成",
                            "source_refs": ["source.txt"],
                            "points": 20,
                            "blocking": True,
                            "trigger": {
                                "type": "text_contains",
                                "values": ["超额完成"],
                                "scope": {"kind": "artifact"},
                            },
                        }
                    ],
                    "review_criteria": [
                        {
                            "id": "coherence",
                            "dimension": "coherence",
                            "description": "结论与行动是否连贯",
                            "anchors": ["0：不连贯", "1：基本连贯", "2：清楚连贯"],
                            "minimum_score": 1,
                        },
                        {
                            "id": "decision",
                            "dimension": "decision_utility",
                            "description": "是否支持经营决策",
                            "anchors": ["0：不可决策", "1：可初步判断", "2：可直接行动"],
                            "minimum_score": 1,
                        },
                        {
                            "id": "visual",
                            "dimension": "visual_quality",
                            "description": "版面是否可读",
                            "anchors": ["0：不可读", "1：基本可读", "2：层次清晰"],
                            "minimum_score": 1,
                        },
                    ],
                    "pass_threshold": 75,
                }
            ],
        }
    )


def test_schema_rejects_format_incompatible_checks() -> None:
    payload = _calibration_suite().model_dump(mode="json")
    payload["items"][0]["rubric"][0]["check"] = {
        "type": "structure_count",
        "metric": "xlsx_sheet_count",
        "operator": "gte",
        "value": 1,
    }

    with pytest.raises(ValueError, match="artifact_type 不兼容"):
        OfficeContentSuite.model_validate(payload)


def _save_calibration_docx(path: Path, *, good: bool) -> None:
    document = Document()
    if good:
        document.add_heading("执行摘要", level=1)
        document.add_paragraph("ARR 未达目标，需要跟进。")
        document.add_heading("KPI", level=1)
        table = document.add_table(rows=2, cols=3)
        table.rows[0].cells[0].text = "指标"
        table.rows[0].cells[1].text = "实际"
        table.rows[0].cells[2].text = "目标"
        table.rows[1].cells[0].text = "ARR"
        table.rows[1].cells[1].text = "860"
        table.rows[1].cells[2].text = "900"
        document.add_heading("风险", level=1)
        document.add_paragraph("API 黄灯")
        document.add_heading("行动", level=1)
        document.add_paragraph("安排复核")
    else:
        document.add_paragraph("执行摘要 KPI 860 900 风险 API 黄灯 行动，ARR 超额完成。")
    document.save(str(path))


def _submission(root: Path, *, good: bool) -> Path:
    path = root / "calibration-docx" / "submission" / "brief.docx"
    path.parent.mkdir(parents=True)
    _save_calibration_docx(path, good=good)
    return path


def _calibration_reviews(sha256: str, *, score: int | dict[str, int]) -> ReviewFile:
    return ReviewFile.model_validate(
        {
            "schema_version": "workpilot-office-content-reviews.v1",
            "reviews": [
                {
                    "item_id": "calibration-docx",
                    "criterion_id": criterion_id,
                    "artifact_sha256": sha256,
                    "score": score[criterion_id] if isinstance(score, dict) else score,
                    "evidence": evidence,
                    "reviewer": "calibration-reviewer",
                }
                for criterion_id, evidence in (
                    ("coherence", "层级和叙事完整"),
                    ("decision", "行动信息足以支持判断"),
                    ("visual", "文档层次清楚可读"),
                )
            ],
        }
    )


def _calibration_model_reviews(sha256: str, *, score: int) -> ReviewFile:
    return ReviewFile.model_validate(
        {
            "schema_version": "workpilot-office-content-reviews.v2",
            "reviews": [
                {
                    "item_id": "calibration-docx",
                    "criterion_id": criterion_id,
                    "artifact_sha256": sha256,
                    "score": score,
                    "evidence": evidence,
                    "reviewer": "model:openai_compatible/vision-test",
                    "source": "model",
                    "provider": "openai_compatible",
                    "model": "vision-test",
                    "prompt_fingerprint": "a" * 64,
                    "authorization_note_fingerprint": "b" * 64,
                    "calibration_status": "uncalibrated",
                    "render_mode": "office_preview",
                }
                for criterion_id, evidence in (
                    ("coherence", "结构顺序能够回扣结论"),
                    ("decision", "行动项包含明确判断入口"),
                    ("visual", "第 1 页层次清晰可读"),
                )
            ],
        }
    )


def test_calibration_ranks_structured_grounded_output_above_keyword_stuffing(
    tmp_path: Path,
) -> None:
    suite = _calibration_suite()
    good_root = tmp_path / "good"
    bad_root = tmp_path / "bad"
    _submission(good_root, good=True)
    _submission(bad_root, good=False)

    good = evaluate_suite(suite, good_root, tmp_path / "good-report")
    bad = evaluate_suite(suite, bad_root, tmp_path / "bad-report")

    good_result = good["results"][0]
    bad_result = bad["results"][0]
    assert good_result["automatic_score"] == 100
    assert good_result["automatic_pass"] is True
    assert good_result["final_score"] is None
    assert bad_result["automatic_score"] == 0
    assert bad_result["automatic_pass"] is False
    assert set(bad_result["critical_failures"]) == {
        "fund-structure",
        "complete-risk",
        "correct-kpi",
        "faithful-no-overclaim",
    }
    assert bad_result["blocking_penalties"] == ["overclaim"]


def test_complete_hash_bound_review_produces_final_score(tmp_path: Path) -> None:
    suite = _calibration_suite()
    submission_root = tmp_path / "submission"
    artifact = _submission(submission_root, good=True)
    first = evaluate_suite(suite, submission_root, tmp_path / "unreviewed")
    sha256 = first["results"][0]["artifact_sha256"]
    assert sha256 is not None
    reviews = _calibration_reviews(sha256, score=2)

    reviewed = evaluate_suite(
        suite,
        submission_root,
        tmp_path / "reviewed",
        reviews=reviews,
    )

    result = reviewed["results"][0]
    assert artifact.is_file()
    assert result["review_complete"] is True
    assert result["review_score"] == 100
    assert result["review_dimension_scores"] == {
        "coherence": 100,
        "decision_utility": 100,
        "visual_quality": 100,
    }
    assert result["review_pass"] is True
    assert result["final_score"] == 100
    assert result["benchmark_pass"] is True

    _save_calibration_docx(artifact, good=False)
    stale = evaluate_suite(
        suite,
        submission_root,
        tmp_path / "stale-review",
        reviews=reviews,
    )["results"][0]
    assert {criterion["status"] for criterion in stale["review"]} == {"stale"}
    assert stale["review_complete"] is False
    assert stale["final_score"] is None


def test_review_minimums_block_a_high_automatic_score(tmp_path: Path) -> None:
    suite = _calibration_suite()
    submission_root = tmp_path / "submission"
    _submission(submission_root, good=True)
    first = evaluate_suite(suite, submission_root, tmp_path / "unreviewed")
    sha256 = first["results"][0]["artifact_sha256"]
    assert sha256 is not None

    reviewed = evaluate_suite(
        suite,
        submission_root,
        tmp_path / "reviewed",
        reviews=_calibration_reviews(
            sha256,
            score={"coherence": 0, "decision": 2, "visual": 2},
        ),
    )

    result = reviewed["results"][0]
    assert result["automatic_score"] == 100
    assert result["review_complete"] is True
    assert result["review_score"] == 66.67
    assert result["final_score"] == 86.67
    assert result["review_failures"] == ["coherence"]
    assert result["review_pass"] is False
    assert result["benchmark_pass"] is False
    assert reviewed["summary"]["review_passed"] == 0
    assert reviewed["summary"]["benchmark_passed"] == 0


def test_uncalibrated_model_review_scores_engineering_but_not_benchmark(
    tmp_path: Path,
) -> None:
    suite = _calibration_suite()
    submission_root = tmp_path / "submission"
    _submission(submission_root, good=True)
    first = evaluate_suite(suite, submission_root, tmp_path / "unreviewed")
    sha256 = first["results"][0]["artifact_sha256"]
    assert sha256 is not None

    reviewed = evaluate_suite(
        suite,
        submission_root,
        tmp_path / "model-reviewed",
        reviews=_calibration_model_reviews(sha256, score=2),
    )

    result = reviewed["results"][0]
    assert result["final_score"] == 100
    assert result["engineering_pass"] is True
    assert result["review_eligible_for_benchmark"] is False
    assert result["benchmark_pass"] is False
    assert result["review_sources"] == ["model"]
    assert set(result["review_eligibility_failures"]) == {
        "coherence:uncalibrated_model_review",
        "decision:uncalibrated_model_review",
        "visual:uncalibrated_model_review",
    }
    assert reviewed["summary"]["engineering_passed"] == 1
    assert reviewed["summary"]["benchmark_eligible_items"] == 0
    assert reviewed["summary"]["benchmark_passed"] == 0


def test_release_cli_requires_complete_passing_reviews(
    tmp_path: Path,
    capsys: pytest.CaptureFixture[str],
) -> None:
    suite = _calibration_suite()
    suite_path = tmp_path / "suite.json"
    suite_path.write_text(
        json.dumps(suite.model_dump(mode="json"), ensure_ascii=False),
        encoding="utf-8",
    )
    submission_root = tmp_path / "submission"
    _submission(submission_root, good=True)
    unreviewed_output = tmp_path / "unreviewed-output"

    assert (
        _main(
            [
                "score",
                "--suite",
                str(suite_path),
                "--submission-root",
                str(submission_root),
                "--output-dir",
                str(unreviewed_output),
                "--require-complete-reviews",
            ]
        )
        == 1
    )
    unreviewed = json.loads((unreviewed_output / "report.json").read_text(encoding="utf-8"))
    sha256 = unreviewed["results"][0]["artifact_sha256"]
    assert sha256 is not None
    reviews_path = tmp_path / "reviews.json"
    reviews_path.write_text(
        _calibration_reviews(sha256, score=2).model_dump_json(indent=2),
        encoding="utf-8",
    )

    reviewed_output = tmp_path / "reviewed-output"
    assert (
        _main(
            [
                "score",
                "--suite",
                str(suite_path),
                "--submission-root",
                str(submission_root),
                "--output-dir",
                str(reviewed_output),
                "--reviews",
                str(reviews_path),
                "--require-complete-reviews",
            ]
        )
        == 0
    )
    reviewed = json.loads((reviewed_output / "report.json").read_text(encoding="utf-8"))
    assert reviewed["summary"]["review_complete_items"] == 1
    assert reviewed["summary"]["benchmark_passed"] == 1
    capsys.readouterr()


def _judge_response(item: OfficeContentItem, *, score: int = 2) -> str:
    return json.dumps(
        {
            "reviews": [
                {
                    "criterion_id": criterion.id,
                    "score": score,
                    "evidence": f"第 1 页为 {criterion.id} 提供了具体可见证据",
                }
                for criterion in item.review_criteria
            ]
        },
        ensure_ascii=False,
        separators=(",", ":"),
    )


class _FakeOfficeJudge:
    chat_provider = "openai_compatible"
    chat_model = "vision-test"

    def __init__(self, responses: list[str]) -> None:
        self.responses = responses
        self.messages: list[list[Message]] = []

    async def complete(
        self,
        messages: list[Message],
        *,
        task_type: str,
        max_tokens: int,
        temperature: float,
        cache_retention: CacheRetention = "none",
    ) -> CompletionResult:
        assert task_type == "judge"
        assert max_tokens > 0
        assert temperature == 0
        assert cache_retention == "none"
        self.messages.append(messages)
        response = self.responses.pop(0)
        return CompletionResult(
            text=response,
            model=self.chat_model,
            provider=self.chat_provider,
            usage=Usage(input_tokens=100, output_tokens=30),
        )


def _fake_office_renderer_sync(
    item: OfficeContentItem,
    source: Path,
    output_root: Path,
    max_pages: int,
) -> RenderedArtifact:
    assert source.is_file()
    assert max_pages >= 1
    output_root.mkdir(parents=True)
    page = output_root / "page-001.png"
    Image.new("RGB", (320, 180), "white").save(page)
    return RenderedArtifact(
        pages=(page,),
        mode="office_preview",
        warnings=(f"fake-render:{item.artifact_type}",),
    )


async def _fake_office_renderer(
    item: OfficeContentItem,
    source: Path,
    output_root: Path,
    max_pages: int,
) -> RenderedArtifact:
    return await asyncio.to_thread(_fake_office_renderer_sync, item, source, output_root, max_pages)


def test_model_judge_requires_explicit_artifact_send_authorization(
    tmp_path: Path,
) -> None:
    suite = _calibration_suite()
    submission_root = tmp_path / "submission"
    _submission(submission_root, good=True)

    with pytest.raises(PermissionError, match="显式授权"):
        asyncio.run(
            run_model_reviews(
                suite,
                submission_root,
                tmp_path / "rendered",
                gateway=_FakeOfficeJudge([_judge_response(suite.items[0])]),
                allow_model_send=False,
                authorization_note="",
                expected_provider="openai_compatible",
                expected_model="vision-test",
                renderer=_fake_office_renderer,
            )
        )


def test_model_judge_repairs_invalid_json_and_binds_visual_review_to_hash(
    tmp_path: Path,
) -> None:
    suite = _calibration_suite()
    item = suite.items[0]
    submission_root = tmp_path / "submission"
    artifact = _submission(submission_root, good=True)
    gateway = _FakeOfficeJudge(["not json", _judge_response(item)])

    reviews, run = asyncio.run(
        run_model_reviews(
            suite,
            submission_root,
            tmp_path / "rendered",
            gateway=gateway,
            allow_model_send=True,
            authorization_note="仅允许发送本次 synthetic Office calibration 文件",
            expected_provider="openai_compatible",
            expected_model="vision-test",
            max_model_calls=2,
            renderer=_fake_office_renderer,
        )
    )

    assert reviews.schema_version == "workpilot-office-content-reviews.v2"
    assert len(reviews.reviews) == 3
    assert {review.source for review in reviews.reviews} == {"model"}
    assert {review.artifact_sha256 for review in reviews.reviews} == {
        hashlib.sha256(artifact.read_bytes()).hexdigest()
    }
    assert {review.calibration_status for review in reviews.reviews} == {"uncalibrated"}
    assert run["model_calls"] == 2
    assert run["repair_retries"] == 1
    assert len(str(run["implementation_fingerprint"])) == 64
    assert run["benchmark_eligible"] is False
    assert len(run["items"][0]["page_images"][0]["sha256"]) == 64
    assert len(gateway.messages) == 2
    assert all(len(call[1].attachments) == 1 for call in gateway.messages)


def test_structural_fallback_cannot_claim_top_visual_score() -> None:
    item = _calibration_suite().items[0]

    with pytest.raises(ValueError, match="visual_quality"):
        parse_model_response(
            _judge_response(item, score=2),
            item,
            render_mode="structural_fallback",
        )


def test_model_review_renderer_rasterizes_every_pptx_slide(tmp_path: Path) -> None:
    suite = load_suite(DEFAULT_SUITE)
    item = next(value for value in suite.items if value.id == "office-pptx-001-monthly-review")
    source = tmp_path / item.output_file
    _save_seed_pptx(source)

    rendered = asyncio.run(render_artifact_for_review(item, source, tmp_path / "rendered", 12))

    assert rendered.mode == "native_pptx"
    assert len(rendered.pages) == 5
    assert all(page.is_file() and page.suffix == ".png" for page in rendered.pages)


def test_one_click_office_eval_writes_engineering_result_but_not_benchmark(
    tmp_path: Path,
) -> None:
    suite = _calibration_suite()
    submission_root = tmp_path / "submission"
    _submission(submission_root, good=True)
    output_dir = tmp_path / "one-click-output"
    gateway = _FakeOfficeJudge([_judge_response(suite.items[0])])

    report = asyncio.run(
        run_one_click_evaluation(
            suite,
            submission_root,
            output_dir,
            gateway=gateway,
            allow_model_send=True,
            authorization_note="synthetic Office 文件可发送到 vision-test",
            expected_provider="openai_compatible",
            expected_model="vision-test",
            renderer=_fake_office_renderer,
        )
    )

    assert report["summary"]["automatic_passed"] == 1
    assert report["summary"]["engineering_passed"] == 1
    assert report["summary"]["benchmark_eligible_items"] == 0
    assert report["summary"]["benchmark_passed"] == 0
    assert report["model_review"]["benchmark_eligible"] is False
    assert (output_dir / "report.json").is_file()
    assert (output_dir / "model-reviews.json").is_file()
    assert (output_dir / "model-review-run.json").is_file()
    persisted = json.loads((output_dir / "report.json").read_text(encoding="utf-8"))
    assert persisted["results"][0]["engineering_pass"] is True
    assert persisted["results"][0]["benchmark_pass"] is False


def test_office_eval_configures_native_gemini_judge(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    monkeypatch.setenv("TEST_GEMINI_API_KEY", "test-only-key")
    args = _parse_args(
        [
            "--submission-root",
            str(tmp_path / "submissions"),
            "--output-dir",
            str(tmp_path / "output"),
            "--judge-provider",
            "gemini",
            "--judge-model",
            "models/gemini-3.6-flash",
            "--api-key-env",
            "TEST_GEMINI_API_KEY",
        ]
    )

    provider, gateway, provider_name, model = _configured_judge(args)
    try:
        assert isinstance(provider, GeminiProvider)
        assert provider_name == "gemini"
        assert model == "gemini-3.6-flash"
    finally:
        asyncio.run(gateway.aclose())
