import asyncio
import html
from collections.abc import Callable
from pathlib import Path
from typing import Annotated
from urllib.parse import quote
from uuid import UUID

from docx import Document
from fastapi import APIRouter, Depends, File, HTTPException, Response, UploadFile, status
from fastapi.responses import FileResponse, HTMLResponse
from openpyxl import load_workbook  # type: ignore[import-untyped]
from pptx import Presentation

from app.api.dependencies import require_owner_identity
from app.core.config import Settings, get_settings
from app.core.db import DbSession as AsyncSession
from app.core.db import get_db_session
from app.cowork.approvals import list_approval_rules, revoke_approval_rule
from app.cowork.artifact_formats import TEXT_ARTIFACT_SUFFIXES
from app.cowork.artifacts import ArtifactRegistrationError, list_artifacts, resolve_artifact_file
from app.cowork.attachments import CoworkAttachmentError, store_attachment
from app.cowork.memory import (
    MemoryNotFoundError,
    MemoryScopeError,
    default_workspace_path,
    forget_memory,
    list_memories,
    remember,
    update_memory,
)
from app.cowork.office_preview import OfficePreviewError, render_office_preview
from app.cowork.permissions import (
    CapabilityDeniedError,
    ConversationNotFoundError,
    CoworkPermissionError,
    SessionRootNotFoundError,
    authorize_path,
    create_session_root,
    grant_capability,
    list_capability_grants,
    list_session_roots,
    revoke_capability_grant,
    revoke_session_root,
)
from app.cowork.reading import (
    Material,
    ReadingError,
    default_material_cache,
    render_units,
)
from app.cowork.workspace_trust import (
    WorkspaceTrustError,
    is_workspace_trusted,
    read_workspace_allowlist,
    set_workspace_trust,
)
from app.ingest.pdf_render import PdfRenderError, render_pdf_page
from app.knowledge_contracts import KnowledgeUnavailableError
from app.rag.kb import KbManifest, local_kb_service
from app.rag.kb.jobs import IndexingJob, default_indexing_jobs
from app.rag.kb.service import SkippedSource, expand_sources
from app.runstore.conversations import get_conversation_kb, set_conversation_kb
from app.schemas.cowork import (
    ApprovalRuleListResponse,
    ApprovalRuleResponse,
    ArtifactListResponse,
    ArtifactResponse,
    AttachmentResponse,
    CapabilityGrantCreate,
    CapabilityGrantListResponse,
    CapabilityGrantResponse,
    ConversationKnowledgeBase,
    KnowledgeBaseAddDocuments,
    KnowledgeBaseCreate,
    KnowledgeBaseDocumentResponse,
    KnowledgeBaseIndexingJob,
    KnowledgeBaseListResponse,
    KnowledgeBaseResponse,
    KnowledgeBaseSkipped,
    MemoryCreate,
    MemoryListResponse,
    MemoryPatch,
    MemoryResponse,
    ReadingMaterialResponse,
    ReadingOutlineEntry,
    ReadingUnitResponse,
    SessionRootCreate,
    SessionRootListResponse,
    SessionRootResponse,
    WorkspaceTrustEntry,
    WorkspaceTrustListResponse,
    WorkspaceTrustUpdate,
)


async def require_cowork_enabled(
    settings: Annotated[Settings, Depends(get_settings)],
) -> None:
    if not settings.cowork_enabled:
        raise HTTPException(status_code=404, detail="Cowork 功能尚未启用")


router = APIRouter(
    prefix="/api/v1/cowork",
    tags=["cowork"],
    dependencies=[Depends(require_owner_identity), Depends(require_cowork_enabled)],
)
DbSession = Annotated[AsyncSession, Depends(get_db_session)]
_PREVIEW_SECURITY_HEADERS = {
    "Content-Security-Policy": (
        "default-src 'none'; img-src data:; font-src data:; style-src 'unsafe-inline'; sandbox"
    ),
    "Referrer-Policy": "no-referrer",
    "X-Content-Type-Options": "nosniff",
}


def _root_response(value: object) -> SessionRootResponse:
    return SessionRootResponse.model_validate(value, from_attributes=True)


def _grant_response(value: object) -> CapabilityGrantResponse:
    return CapabilityGrantResponse.model_validate(value, from_attributes=True)


def _not_found(error: ConversationNotFoundError) -> HTTPException:
    return HTTPException(status_code=404, detail="Cowork 会话不存在")


@router.post(
    "/sessions/{conversation_id}/attachments",
    response_model=AttachmentResponse,
    status_code=status.HTTP_201_CREATED,
)
async def post_attachment(
    conversation_id: UUID,
    session: DbSession,
    settings: Annotated[Settings, Depends(get_settings)],
    upload: Annotated[UploadFile, File()],
) -> AttachmentResponse:
    """接收一份待绑定输入附件；run 创建成功时才会把它绑定到消息。"""

    try:
        raw = await upload.read(settings.cowork_attachment_max_bytes + 1)
        attachment = await store_attachment(
            session,
            conversation_id=conversation_id,
            filename=upload.filename or "attachment",
            declared_media_type=upload.content_type,
            raw=raw,
            settings=settings,
        )
    except CoworkAttachmentError as error:
        await session.rollback()
        status_code = 404 if str(error) == "Cowork 会话不存在" else 422
        raise HTTPException(status_code=status_code, detail=str(error)) from error
    finally:
        await upload.close()
    await session.commit()
    return AttachmentResponse.model_validate(attachment, from_attributes=True)


@router.get("/sessions/{conversation_id}/roots", response_model=SessionRootListResponse)
async def get_session_roots(conversation_id: UUID, session: DbSession) -> SessionRootListResponse:
    try:
        items = await list_session_roots(session, conversation_id=conversation_id)
    except ConversationNotFoundError as error:
        raise _not_found(error) from error
    return SessionRootListResponse(items=[_root_response(item) for item in items])


@router.post(
    "/sessions/{conversation_id}/roots",
    response_model=SessionRootResponse,
    status_code=status.HTTP_201_CREATED,
)
async def post_session_root(
    conversation_id: UUID, request: SessionRootCreate, session: DbSession
) -> SessionRootResponse:
    try:
        root = await create_session_root(
            session,
            conversation_id=conversation_id,
            requested_path=request.path,
            access_mode=request.access_mode,
            label=request.label,
        )
    except ConversationNotFoundError as error:
        raise _not_found(error) from error
    except (CoworkPermissionError, ValueError) as error:
        raise HTTPException(status_code=422, detail=str(error)) from error
    await session.commit()
    return _root_response(root)


@router.delete(
    "/sessions/{conversation_id}/roots/{root_id}",
    status_code=status.HTTP_204_NO_CONTENT,
)
async def delete_session_root(conversation_id: UUID, root_id: UUID, session: DbSession) -> Response:
    try:
        deleted = await revoke_session_root(
            session, conversation_id=conversation_id, root_id=root_id
        )
    except ConversationNotFoundError as error:
        raise _not_found(error) from error
    if not deleted:
        raise HTTPException(status_code=404, detail="会话目录不存在")
    await session.commit()
    return Response(status_code=status.HTTP_204_NO_CONTENT)


@router.get("/sessions/{conversation_id}/grants", response_model=CapabilityGrantListResponse)
async def get_capability_grants(
    conversation_id: UUID, session: DbSession
) -> CapabilityGrantListResponse:
    try:
        items = await list_capability_grants(session, conversation_id=conversation_id)
    except ConversationNotFoundError as error:
        raise _not_found(error) from error
    return CapabilityGrantListResponse(items=[_grant_response(item) for item in items])


@router.post(
    "/sessions/{conversation_id}/grants",
    response_model=CapabilityGrantResponse,
    status_code=status.HTTP_201_CREATED,
)
async def post_capability_grant(
    conversation_id: UUID, request: CapabilityGrantCreate, session: DbSession
) -> CapabilityGrantResponse:
    try:
        grant = await grant_capability(
            session,
            conversation_id=conversation_id,
            capability=request.capability,
            session_root_id=request.session_root_id,
            expires_in_s=request.expires_in_s,
        )
    except ConversationNotFoundError as error:
        raise _not_found(error) from error
    except SessionRootNotFoundError as error:
        raise HTTPException(status_code=404, detail="会话目录不存在") from error
    except (CapabilityDeniedError, ValueError) as error:
        raise HTTPException(status_code=422, detail=str(error)) from error
    await session.commit()
    return _grant_response(grant)


@router.delete(
    "/sessions/{conversation_id}/grants/{grant_id}",
    status_code=status.HTTP_204_NO_CONTENT,
)
async def delete_capability_grant(
    conversation_id: UUID, grant_id: UUID, session: DbSession
) -> Response:
    try:
        deleted = await revoke_capability_grant(
            session, conversation_id=conversation_id, grant_id=grant_id
        )
    except ConversationNotFoundError as error:
        raise _not_found(error) from error
    if not deleted:
        raise HTTPException(status_code=404, detail="能力授权不存在")
    await session.commit()
    return Response(status_code=status.HTTP_204_NO_CONTENT)


def _approval_rule_response(record: object) -> ApprovalRuleResponse:
    return ApprovalRuleResponse.model_validate(record, from_attributes=True)


@router.get(
    "/sessions/{conversation_id}/approval-rules",
    response_model=ApprovalRuleListResponse,
)
async def get_approval_rules(conversation_id: UUID, session: DbSession) -> ApprovalRuleListResponse:
    """会话攒下的常驻审批规则。

    没有创建端点：规则只能在审批卡片上产生。让 API 能凭空造一条，等于给"用户当场看过
    这次调用"这个前提开了一个后门——而那正是常驻授权唯一的正当性来源。
    """

    items = await list_approval_rules(session, conversation_id=conversation_id)
    return ApprovalRuleListResponse(items=[_approval_rule_response(item) for item in items])


@router.delete(
    "/sessions/{conversation_id}/approval-rules/{rule_id}",
    status_code=status.HTTP_204_NO_CONTENT,
)
async def delete_approval_rule(
    conversation_id: UUID, rule_id: UUID, session: DbSession
) -> Response:
    revoked = await revoke_approval_rule(session, conversation_id=conversation_id, rule_id=rule_id)
    if not revoked:
        raise HTTPException(status_code=404, detail="常驻审批规则不存在")
    await session.commit()
    return Response(status_code=status.HTTP_204_NO_CONTENT)


@router.get(
    "/sessions/{conversation_id}/workspace-trust",
    response_model=WorkspaceTrustListResponse,
)
async def get_workspace_trust(
    conversation_id: UUID, session: DbSession
) -> WorkspaceTrustListResponse:
    """本会话每个已授权目录的信任状态，以及它自己声明了哪些命令前缀。"""

    roots = await list_session_roots(session, conversation_id=conversation_id)
    items: list[WorkspaceTrustEntry] = []
    for root in roots:
        trusted = await is_workspace_trusted(session, canonical_path=root.canonical_path)
        declared: list[str] = []
        rejected: list[str] = []
        config_error: str | None = None
        try:
            allowlist = await asyncio.to_thread(read_workspace_allowlist, Path(root.canonical_path))
        except WorkspaceTrustError as error:
            config_error = str(error)
        else:
            declared = list(allowlist.entries)
            rejected = [f"{entry}：{reason}" for entry, reason in allowlist.rejected]
        items.append(
            WorkspaceTrustEntry(
                canonical_path=root.canonical_path,
                trusted=trusted,
                declared=declared,
                rejected=rejected,
                config_error=config_error,
            )
        )
    return WorkspaceTrustListResponse(items=items)


@router.put(
    "/sessions/{conversation_id}/workspace-trust",
    response_model=WorkspaceTrustListResponse,
)
async def put_workspace_trust(
    conversation_id: UUID, request: WorkspaceTrustUpdate, session: DbSession
) -> WorkspaceTrustListResponse:
    """信任或撤销信任一个目录。

    只接受**本会话已经授权过的 root**：信任一个用户还没选进来的目录，等于让这个接口
    凭空扩大可执行范围。
    """

    roots = await list_session_roots(session, conversation_id=conversation_id)
    if all(root.canonical_path != request.canonical_path for root in roots):
        raise HTTPException(status_code=404, detail="该目录不在本会话的已授权目录里")
    await set_workspace_trust(
        session, canonical_path=request.canonical_path, trusted=request.trusted
    )
    await session.commit()
    return await get_workspace_trust(conversation_id, session)


def _memory_response(record: object) -> MemoryResponse:
    return MemoryResponse.model_validate(record, from_attributes=True)


@router.get("/sessions/{conversation_id}/memories", response_model=MemoryListResponse)
async def get_memories(
    conversation_id: UUID,
    session: DbSession,
    settings: Annotated[Settings, Depends(get_settings)],
    include_forgotten: bool = False,
) -> MemoryListResponse:
    """列出当前会话可见的记忆。

    可见范围和注入给模型的完全一致：global + 本会话授权目录的 workspace + 本会话。
    面板里看到的就是模型看到的，不然用户没法判断"它为什么会这么以为"。
    """

    try:
        roots = await list_session_roots(session, conversation_id=conversation_id)
        items = await list_memories(
            session,
            conversation_id=conversation_id,
            workspace_paths=[root.canonical_path for root in roots],
            include_forgotten=include_forgotten,
            limit=settings.cowork_memory_max_items,
        )
    except ConversationNotFoundError as error:
        raise _not_found(error) from error
    return MemoryListResponse(items=[_memory_response(item) for item in items])


@router.post(
    "/sessions/{conversation_id}/memories",
    response_model=MemoryResponse,
    status_code=status.HTTP_201_CREATED,
)
async def post_memory(
    conversation_id: UUID, request: MemoryCreate, session: DbSession
) -> MemoryResponse:
    try:
        workspace_path = (
            await default_workspace_path(session, conversation_id=conversation_id)
            if request.scope == "workspace"
            else None
        )
        record, _ = await remember(
            session,
            conversation_id=conversation_id,
            scope=request.scope,
            content=request.content,
            key=request.key,
            workspace_path=workspace_path,
            source="user",
        )
    except ConversationNotFoundError as error:
        raise _not_found(error) from error
    except MemoryScopeError as error:
        raise HTTPException(status_code=422, detail=str(error)) from error
    await session.commit()
    return _memory_response(record)


@router.patch("/memories/{memory_id}", response_model=MemoryResponse)
async def patch_memory(memory_id: UUID, request: MemoryPatch, session: DbSession) -> MemoryResponse:
    """改写或恢复一条记忆——客户端的「撤销」也走这里。"""

    try:
        record, _ = await update_memory(
            session,
            memory_id=memory_id,
            content=request.content,
            restore=request.restore,
        )
    except MemoryNotFoundError as error:
        raise HTTPException(status_code=404, detail="记忆不存在") from error
    except MemoryScopeError as error:
        raise HTTPException(status_code=422, detail=str(error)) from error
    await session.commit()
    return _memory_response(record)


@router.delete("/memories/{memory_id}", status_code=status.HTTP_204_NO_CONTENT)
async def delete_memory(memory_id: UUID, session: DbSession) -> Response:
    """软删除。记录留着，撤销和事后追查都还能拿到原文。"""

    record = await forget_memory(session, memory_id=memory_id)
    if record is None:
        raise HTTPException(status_code=404, detail="记忆不存在或已经 retire")
    await session.commit()
    return Response(status_code=status.HTTP_204_NO_CONTENT)


@router.get("/sessions/{conversation_id}/artifacts", response_model=ArtifactListResponse)
async def get_artifacts(conversation_id: UUID, session: DbSession) -> ArtifactListResponse:
    try:
        items = await list_artifacts(session, conversation_id=conversation_id)
    except ConversationNotFoundError as error:
        raise _not_found(error) from error
    return ArtifactListResponse(
        items=[ArtifactResponse.model_validate(item, from_attributes=True) for item in items]
    )


@router.get("/artifacts/{artifact_id}/preview")
async def get_artifact_preview(
    artifact_id: UUID,
    session: DbSession,
    settings: Annotated[Settings, Depends(get_settings)],
) -> Response:
    try:
        resolved = await resolve_artifact_file(session, artifact_id=artifact_id)
    except ArtifactRegistrationError as error:
        raise HTTPException(status_code=410, detail=str(error)) from error
    if resolved is None:
        raise HTTPException(status_code=404, detail="交付物不存在")
    artifact, path = resolved
    suffix = path.suffix.casefold()
    # 预览类型只能由经过白名单校验的扩展名决定，不能信任模型登记的 mime_type。
    if suffix == ".pdf":
        return FileResponse(
            path,
            media_type="application/pdf",
            headers={
                **_PREVIEW_SECURITY_HEADERS,
                "Content-Disposition": _preview_content_disposition(path),
                "X-WorkPilot-Preview-Mode": "native-pdf",
            },
        )
    if suffix in {".docx", ".xlsx", ".pptx"}:
        try:
            rendered = await asyncio.to_thread(
                render_office_preview,
                path,
                cache_root=settings.office_preview_cache_path,
                timeout_s=settings.office_preview_timeout_s,
                max_source_bytes=settings.office_preview_max_source_bytes,
                max_cache_entries=settings.office_preview_max_cache_entries,
            )
        except OfficePreviewError as error:
            raise HTTPException(status_code=422, detail=str(error)) from error
        if rendered is not None:
            return FileResponse(
                rendered.path,
                media_type=rendered.media_type,
                headers={
                    **_PREVIEW_SECURITY_HEADERS,
                    "Content-Disposition": _preview_content_disposition(path),
                    "X-WorkPilot-Preview-Mode": rendered.mode,
                },
            )
        body = (
            _docx_preview(path)
            if suffix == ".docx"
            else _xlsx_preview(path)
            if suffix == ".xlsx"
            else _pptx_preview(path)
        )
        preview_mode = "structure"
    elif suffix in TEXT_ARTIFACT_SUFFIXES:
        if path.stat().st_size > 5 * 1024 * 1024:
            raise HTTPException(status_code=413, detail="交付物过大，无法在线预览")
        body = f"<pre>{html.escape(path.read_text(encoding='utf-8', errors='replace'))}</pre>"
        preview_mode = "text"
    else:
        raise HTTPException(status_code=415, detail="该交付物格式不支持在线预览")
    document = (
        "<!doctype html><meta charset='utf-8'><style>"
        "body{font:15px/1.65 system-ui;color:#26332f;padding:32px;max-width:960px;margin:auto}"
        "table{border-collapse:collapse;width:100%}td,th{border:1px solid #d9dedb;padding:6px 8px}"
        "pre{white-space:pre-wrap;word-break:break-word}</style>"
        f"<title>{html.escape(artifact.title)}</title><h1>{html.escape(artifact.title)}</h1>{body}"
    )
    return HTMLResponse(
        document,
        headers={**_PREVIEW_SECURITY_HEADERS, "X-WorkPilot-Preview-Mode": preview_mode},
    )


def _docx_preview(path: Path) -> str:
    document = Document(str(path))
    blocks: list[str] = []
    for paragraph in document.paragraphs[:1000]:
        text = html.escape(paragraph.text)
        if not text:
            continue
        style = paragraph.style.name.casefold() if paragraph.style is not None else ""
        tag = "h2" if "heading" in style else "p"
        blocks.append(f"<{tag}>{text}</{tag}>")
    return "".join(blocks)


def _xlsx_preview(path: Path) -> str:
    workbook = load_workbook(path, read_only=True, data_only=False)
    try:
        blocks: list[str] = []
        for sheet in workbook.worksheets[:5]:
            blocks.append(f"<h2>{html.escape(sheet.title)}</h2><table>")
            for row in sheet.iter_rows(min_row=1, max_row=100, max_col=20, values_only=True):
                blocks.append(
                    "<tr>"
                    + "".join(
                        f"<td>{html.escape(str(value if value is not None else ''))}</td>"
                        for value in row
                    )
                    + "</tr>"
                )
            blocks.append("</table>")
        return "".join(blocks)
    finally:
        workbook.close()


def _pptx_preview(path: Path) -> str:
    presentation = Presentation(str(path))
    blocks: list[str] = []
    for index, slide in enumerate(list(presentation.slides)[:100], start=1):
        blocks.append(f"<h2>第 {index} 页</h2>")
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                blocks.append(f"<p>{html.escape(str(text))}</p>")
    return "".join(blocks)


def _safe_preview_name(path: Path) -> str:
    value = "".join(
        character if character.isascii() and (character.isalnum() or character in "._-") else "_"
        for character in path.name
    ).strip("._")
    if value:
        return value
    return f"artifact{path.suffix if path.suffix.isascii() else ''}"


def _preview_content_disposition(path: Path) -> str:
    """生成只含 ASCII 的响应头，同时保留 RFC 5987 中文文件名。"""

    fallback = _safe_preview_name(path)
    encoded = quote(path.name, safe="")
    return f"inline; filename=\"{fallback}\"; filename*=UTF-8''{encoded}"


# --- 阅读器面板 -----------------------------------------------------------------
#
# 面板要渲染文档，就必须能按路径取内容——于是这里是除了工具执行边界之外，**第二个**
# 接受用户可控路径的地方。三个端点因此都先过 `authorize_path`：不这么做，任何人只要
# 知道会话 id 就能把本机任意文件的内容读出来，而工具那一侧每次调用都在校验。
#
# 授权失败一律 403 且不区分"没授权"与"文件不存在"：区分开就等于给了一个探测本机文件
# 是否存在的接口。


async def _authorized_material(
    session: AsyncSession,
    *,
    conversation_id: UUID,
    path: str,
    settings: Settings,
) -> Material:
    try:
        authorization = await authorize_path(
            session,
            conversation_id=conversation_id,
            target_path=Path(path),
            capability="filesystem.read",
        )
    except CapabilityDeniedError as error:
        raise HTTPException(status_code=403, detail=str(error)) from error
    except ConversationNotFoundError as error:
        raise _not_found(error) from error
    try:
        return await default_material_cache().load(authorization.target_path, settings=settings)
    except ReadingError as error:
        raise HTTPException(status_code=422, detail=str(error)) from error


@router.get(
    "/sessions/{conversation_id}/reading/material",
    response_model=ReadingMaterialResponse,
)
async def get_reading_material(
    conversation_id: UUID,
    path: str,
    session: DbSession,
    settings: Annotated[Settings, Depends(get_settings)],
) -> ReadingMaterialResponse:
    material = await _authorized_material(
        session, conversation_id=conversation_id, path=path, settings=settings
    )
    return ReadingMaterialResponse(
        path=str(material.path),
        material_id=material.material_id,
        filename=material.filename,
        title=material.title,
        unit=material.unit,
        unit_count=material.unit_count,
        parser=material.parser,
        has_page_image=material.unit == "page",
        outline=[
            ReadingOutlineEntry(
                locator=entry.locator,
                title=entry.title,
                level=entry.level,
                synthesised=entry.synthesised,
            )
            for entry in material.outline
        ],
    )


@router.get(
    "/sessions/{conversation_id}/reading/units/{locator}",
    response_model=ReadingUnitResponse,
)
async def get_reading_unit(
    conversation_id: UUID,
    locator: int,
    path: str,
    session: DbSession,
    settings: Annotated[Settings, Depends(get_settings)],
) -> ReadingUnitResponse:
    material = await _authorized_material(
        session, conversation_id=conversation_id, path=path, settings=settings
    )
    try:
        # 复用工具那条渲染路径，locator 语义和截断上限都不会和模型看到的分叉。
        rendered = render_units(material, locator, max_chars=settings.cowork_pdf_text_max_chars)
    except ReadingError as error:
        raise HTTPException(status_code=404, detail=str(error)) from error
    return ReadingUnitResponse(locator=locator, unit=material.unit, text=rendered.text)


@router.get("/sessions/{conversation_id}/reading/pages/{locator}.png")
async def get_reading_page(
    conversation_id: UUID,
    locator: int,
    path: str,
    session: DbSession,
    settings: Annotated[Settings, Depends(get_settings)],
) -> Response:
    material = await _authorized_material(
        session, conversation_id=conversation_id, path=path, settings=settings
    )
    if material.unit != "page":
        raise HTTPException(status_code=409, detail="这份材料没有可渲染的原页，请读文本视图")
    try:
        content = await asyncio.to_thread(render_pdf_page, material.path, locator)
    except PdfRenderError as error:
        raise HTTPException(status_code=404, detail=str(error)) from error
    return Response(
        content,
        media_type="image/png",
        # material_id 是内容哈希，文件一改就换 URL，所以可以放心长缓存。
        headers={"Cache-Control": "private, max-age=3600"},
    )


def _kb_response(manifest: KbManifest) -> KnowledgeBaseResponse:
    return KnowledgeBaseResponse(
        slug=manifest.slug,
        name=manifest.name,
        description=manifest.description,
        document_count=len(manifest.documents),
        is_indexed=manifest.is_indexed,
        embedding=None if manifest.embedding is None else manifest.embedding.describe(),
        documents=[
            KnowledgeBaseDocumentResponse(
                doc_id=document.doc_id,
                filename=document.filename,
                title=document.title,
                parser=document.parser,
                char_count=document.char_count,
            )
            for document in manifest.documents
        ],
    )


@router.get("/knowledge-bases", response_model=KnowledgeBaseListResponse)
async def list_knowledge_bases(
    settings: Annotated[Settings, Depends(get_settings)],
) -> KnowledgeBaseListResponse:
    service = local_kb_service(settings)
    manifests = await asyncio.to_thread(service.list_kbs)
    return KnowledgeBaseListResponse(items=[_kb_response(item) for item in manifests])


@router.get(
    "/sessions/{conversation_id}/knowledge-base",
    response_model=ConversationKnowledgeBase,
)
async def get_session_knowledge_base(
    conversation_id: UUID,
    session: DbSession,
) -> ConversationKnowledgeBase:
    slug = await get_conversation_kb(
        session,
        conversation_id=conversation_id,
    )
    return ConversationKnowledgeBase(slug=slug)


@router.put(
    "/sessions/{conversation_id}/knowledge-base",
    response_model=ConversationKnowledgeBase,
)
async def put_session_knowledge_base(
    conversation_id: UUID,
    payload: ConversationKnowledgeBase,
    session: DbSession,
    settings: Annotated[Settings, Depends(get_settings)],
) -> ConversationKnowledgeBase:
    """把一个知识库挂到会话上；`slug=null` 卸载。

    挂之前先确认这个库真的存在。数据库那一列没有外键（KB 的事实来源是磁盘上的
    manifest，不是表），所以校验必须发生在这里——否则用户会在挂载时看到成功、
    在下一次提问时才看到检索失败。
    """
    slug = (payload.slug or "").strip() or None
    if slug is not None:
        try:
            await asyncio.to_thread(local_kb_service(settings).get, slug)
        except KnowledgeUnavailableError as error:
            raise HTTPException(status_code=404, detail=str(error)) from error
        except ValueError as error:
            raise HTTPException(status_code=422, detail=str(error)) from error
    changed = await set_conversation_kb(
        session,
        conversation_id=conversation_id,
        kb_slug=slug,
    )
    if not changed:
        raise HTTPException(status_code=404, detail="会话不存在")
    await session.commit()
    return ConversationKnowledgeBase(slug=slug)


def _expand_requested_sources(paths: list[str]) -> list[Path]:
    """`~` 展开 + 目录递归。整个跑在线程里：rglob 扫一个大目录是真的会阻塞事件循环。"""
    return expand_sources([Path(item).expanduser() for item in paths])


def _job_response(job: IndexingJob) -> KnowledgeBaseIndexingJob:
    return KnowledgeBaseIndexingJob(
        slug=job.slug,
        status=job.status,
        stage=job.stage,
        done=job.done,
        total=job.total,
        added=job.added,
        error=job.error,
        skipped=[
            KnowledgeBaseSkipped(filename=item.filename, reason=item.reason)
            for item in job.skipped
        ],
    )


@router.post(
    "/knowledge-bases",
    response_model=KnowledgeBaseResponse,
    status_code=status.HTTP_201_CREATED,
)
async def create_knowledge_base(
    payload: KnowledgeBaseCreate,
    settings: Annotated[Settings, Depends(get_settings)],
) -> KnowledgeBaseResponse:
    service = local_kb_service(settings)
    try:
        manifest = await asyncio.to_thread(
            service.create,
            payload.name.strip() or payload.slug,
            slug=payload.slug.strip(),
            description=payload.description,
        )
    except ValueError as error:
        # KbNameError 就是 ValueError，消息按约束 4 已经写成了可执行指令。
        raise HTTPException(status_code=422, detail=str(error)) from error
    return _kb_response(manifest)


@router.delete("/knowledge-bases/{slug}", status_code=status.HTTP_204_NO_CONTENT)
async def delete_knowledge_base(
    slug: str,
    settings: Annotated[Settings, Depends(get_settings)],
) -> Response:
    """删库连索引一起删。

    **挂着这个库的会话不会被自动解绑。** 要解绑就得扫两套 store 里的所有会话，而收益是
    零：检索侧本来就会给出"知识库 X 不存在"这条可执行错误，管理界面也会把失效的挂载标出来。
    与其为此加一次全表扫描，不如让失效状态可见。
    """
    if default_indexing_jobs().is_running(slug):
        raise HTTPException(status_code=409, detail=f"知识库 {slug} 正在建索引，等它完成再删。")
    try:
        removed = await asyncio.to_thread(local_kb_service(settings).delete, slug)
    except ValueError as error:
        raise HTTPException(status_code=422, detail=str(error)) from error
    if not removed:
        raise HTTPException(status_code=404, detail=f"知识库 {slug} 不存在")
    return Response(status_code=status.HTTP_204_NO_CONTENT)


@router.post(
    "/knowledge-bases/{slug}/documents",
    response_model=KnowledgeBaseIndexingJob,
    status_code=status.HTTP_202_ACCEPTED,
)
async def add_knowledge_base_documents(
    slug: str,
    payload: KnowledgeBaseAddDocuments,
    settings: Annotated[Settings, Depends(get_settings)],
) -> KnowledgeBaseIndexingJob:
    """导入文档并重建索引；立刻返回作业状态，用 GET .../indexing 轮询。

    一个文件夹的论文解析加 embedding 是分钟级的活，挂在 HTTP 请求上必然超时
    （CLAUDE.md：worker 不依附 HTTP 连接）。

    路径不设根目录限制：这是本机 owner 亲自点的导入，和他直接跑 CLI 是同一件事，
    而 KB 的产品前提就是"资料留在你自己的目录里"。Agent 够不到这个接口——它是 HTTP
    端点不是工具，模型没有任何路径能调到。
    """
    service = local_kb_service(settings)
    try:
        await asyncio.to_thread(service.get, slug)
        sources = await asyncio.to_thread(_expand_requested_sources, payload.paths)
    except KnowledgeUnavailableError as error:
        raise HTTPException(status_code=404, detail=str(error)) from error
    except (ValueError, OSError) as error:
        raise HTTPException(status_code=422, detail=str(error)) from error
    if not sources:
        raise HTTPException(
            status_code=422,
            detail="这些路径里没有可导入的文件。支持 .pdf/.md/.markdown/.txt，目录会递归展开。",
        )

    async def work(
        progress: Callable[[str, int, int], None],
    ) -> tuple[int, tuple[SkippedSource, ...]]:
        # skip_failures=True：一整个文件夹里混进一个扫描件，把另外二十九篇一起退回去
        # 是最没用的行为。跳过并逐条报告。
        result = await service.add_documents(
            slug, sources, skip_failures=True, progress=progress
        )
        return len(result.added), result.skipped

    try:
        job = default_indexing_jobs().start(slug, work, stage=f"准备导入 {len(sources)} 个文件")
    except KnowledgeUnavailableError as error:
        raise HTTPException(status_code=409, detail=str(error)) from error
    return _job_response(job)


@router.post(
    "/knowledge-bases/{slug}/rebuild",
    response_model=KnowledgeBaseIndexingJob,
    status_code=status.HTTP_202_ACCEPTED,
)
async def rebuild_knowledge_base(
    slug: str,
    settings: Annotated[Settings, Depends(get_settings)],
) -> KnowledgeBaseIndexingJob:
    """按清单里的源路径重新解析并建索引。换了 embedding 模型之后必须做这一步。"""
    service = local_kb_service(settings)
    try:
        manifest = await asyncio.to_thread(service.get, slug)
    except KnowledgeUnavailableError as error:
        raise HTTPException(status_code=404, detail=str(error)) from error
    if not manifest.documents:
        raise HTTPException(status_code=422, detail="这个知识库还没有文档，没有可重建的内容。")

    async def work(
        progress: Callable[[str, int, int], None],
    ) -> tuple[int, tuple[SkippedSource, ...]]:
        rebuilt = await service.rebuild(slug, progress=progress)
        return len(rebuilt.documents), ()

    try:
        job = default_indexing_jobs().start(slug, work, stage="准备重建")
    except KnowledgeUnavailableError as error:
        raise HTTPException(status_code=409, detail=str(error)) from error
    return _job_response(job)


@router.get(
    "/knowledge-bases/{slug}/indexing",
    response_model=KnowledgeBaseIndexingJob | None,
)
async def get_knowledge_base_indexing(slug: str) -> KnowledgeBaseIndexingJob | None:
    """当前或最近一次建索引作业；没有则返回 null。

    作业表在进程内，重启会丢。丢的只是记录：索引要么写成了、要么没写成，清单里的
    embedding 签名就是答案。
    """
    job = default_indexing_jobs().get(slug)
    return None if job is None else _job_response(job)
