"""原生 PPTX、DOCX、XLSX、PDF 交付物生成与原子替换。"""

from __future__ import annotations

import hashlib
import html
import os
import re
import stat
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Literal

import fitz  # type: ignore[import-untyped]
from docx import Document
from openpyxl import Workbook  # type: ignore[import-untyped]
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from app.services.cowork_files import create_file_backup


@dataclass(frozen=True)
class NativeArtifactResult:
    path: Path
    sha256: str
    size_bytes: int
    mime_type: str
    backup_path: Path | None


def create_native_artifact(
    path: Path,
    *,
    format: Literal["docx", "xlsx", "pptx", "pdf"],
    title: str,
    content: str,
    sheets: list[dict[str, Any]],
    baseline_sha256: str | None,
    slides: list[dict[str, Any]] | None = None,
    backup_versions: int = 5,
) -> NativeArtifactResult:
    expected_suffix = f".{format}"
    if path.suffix.casefold() != expected_suffix:
        raise ValueError(f"{format} 交付物路径必须以 {expected_suffix} 结尾")
    path.parent.mkdir(parents=True, exist_ok=True)
    _check_baseline(path, baseline_sha256)
    descriptor, temporary_name = tempfile.mkstemp(
        prefix=f".{path.name}.", suffix=".tmp", dir=path.parent
    )
    os.close(descriptor)
    temporary = Path(temporary_name)
    try:
        if format == "docx":
            _write_docx(temporary, title=title, content=content)
            mime_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        elif format == "xlsx":
            _write_xlsx(temporary, title=title, sheets=sheets)
            mime_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        elif format == "pptx":
            _write_pptx(temporary, title=title, content=content, slides=slides or [])
            mime_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        else:
            _write_pdf(temporary, title=title, content=content)
            mime_type = "application/pdf"
        payload = temporary.read_bytes()
        # 生成可能耗时；替换前再次核对 baseline，防止覆盖期间的并发修改。
        _check_baseline(path, baseline_sha256)
        previous_mode = stat.S_IMODE(path.stat().st_mode) if path.exists() else None
        backup_path = create_file_backup(path, backup_versions) if path.exists() else None
        if previous_mode is not None:
            os.chmod(temporary, previous_mode)
        os.replace(temporary, path)
        directory_fd = os.open(path.parent, os.O_RDONLY)
        try:
            os.fsync(directory_fd)
        finally:
            os.close(directory_fd)
        return NativeArtifactResult(
            path=path,
            sha256=hashlib.sha256(payload).hexdigest(),
            size_bytes=len(payload),
            mime_type=mime_type,
            backup_path=backup_path,
        )
    finally:
        temporary.unlink(missing_ok=True)


def _check_baseline(path: Path, baseline: str | None) -> None:
    if not path.exists():
        if baseline is not None:
            raise ValueError("原生交付物尚不存在，baseline_sha256 必须省略")
        return
    if path.is_symlink() or not path.is_file():
        raise ValueError("原生交付物目标必须是普通文件")
    if baseline is None:
        raise ValueError("覆盖已有原生交付物必须提供 baseline_sha256")
    current = hashlib.sha256(path.read_bytes()).hexdigest()
    if current != baseline:
        raise ValueError("原生交付物已被修改；baseline_sha256 不匹配")


def _write_docx(path: Path, *, title: str, content: str) -> None:
    document = Document()
    document.core_properties.title = title
    document.add_heading(title, level=0)
    for raw_line in content.splitlines():
        line = raw_line.strip()
        if not line:
            continue
        if line.startswith("### "):
            document.add_heading(line[4:], level=3)
        elif line.startswith("## "):
            document.add_heading(line[3:], level=2)
        elif line.startswith("# "):
            document.add_heading(line[2:], level=1)
        elif line.startswith(("- ", "* ")):
            document.add_paragraph(line[2:], style="List Bullet")
        else:
            document.add_paragraph(line)
    document.save(str(path))


def _write_xlsx(path: Path, *, title: str, sheets: list[dict[str, Any]]) -> None:
    if not sheets:
        sheets = [{"name": title[:31] or "Sheet1", "rows": []}]
    workbook = Workbook()
    default = workbook.active
    workbook.remove(default)
    total_cells = 0
    for index, raw_sheet in enumerate(sheets):
        name = str(raw_sheet.get("name") or f"Sheet{index + 1}")[:31]
        sheet = workbook.create_sheet(name)
        rows = raw_sheet.get("rows") or []
        if not isinstance(rows, list):
            raise ValueError("XLSX sheets.rows 必须是二维数组")
        for row in rows:
            if not isinstance(row, list):
                raise ValueError("XLSX sheets.rows 必须是二维数组")
            total_cells += len(row)
            if total_cells > 100_000:
                raise ValueError("XLSX 单次生成不能超过 100000 个单元格")
            sheet.append(row)
    workbook.save(path)


def _write_pptx(path: Path, *, title: str, content: str, slides: list[dict[str, Any]]) -> None:
    """生成 16:9、可继续编辑的原生演示文稿。"""

    presentation = Presentation()
    presentation.slide_width = Inches(40 / 3)
    presentation.slide_height = Inches(7.5)
    presentation.core_properties.title = title
    blank = presentation.slide_layouts[6]
    normalized = _normalize_slides(content, slides)
    cover_subtitle = str(slides[0].get("subtitle") or "") if slides else ""
    _add_ppt_cover(presentation.slides.add_slide(blank), title, cover_subtitle)
    for index, item in enumerate(normalized, start=1):
        _add_ppt_content_slide(presentation.slides.add_slide(blank), item, index)
    presentation.save(str(path))


def _normalize_slides(content: str, slides: list[dict[str, Any]]) -> list[dict[str, Any]]:
    if slides:
        return [dict(item) for item in slides]
    sections: list[dict[str, Any]] = []
    current: dict[str, Any] = {"title": "主要内容", "bullets": []}
    for raw in content.splitlines():
        line = raw.strip()
        if line.startswith(("# ", "## ")):
            if current.get("body") or current.get("bullets"):
                sections.append(current)
            current = {"title": line.lstrip("# ").strip(), "bullets": []}
        elif line.startswith(("- ", "* ")):
            current.setdefault("bullets", []).append(line[2:].strip())
        elif line:
            current["body"] = f"{current.get('body', '')}\n{line}".strip()
    if current.get("body") or current.get("bullets"):
        sections.append(current)
    if not sections:
        sections.append({"title": "主要内容", "body": content or "内容待补充"})
    return sections[:100]


def _set_slide_background(slide: Any, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rgb(red: int, green: int, blue: int) -> RGBColor:
    return RGBColor(red, green, blue)  # type: ignore[no-untyped-call]


def _add_textbox(
    slide: Any,
    text: str,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    size: int,
    color: RGBColor,
    bold: bool = False,
    align: Any = PP_ALIGN.LEFT,
) -> Any:
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = shape.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    run = paragraph.runs[0]
    run.font.name = "PingFang SC"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def _add_ppt_cover(slide: Any, title: str, subtitle: Any) -> None:
    _set_slide_background(slide, _rgb(255, 248, 235))
    for left, top, size, color in (
        (0.6, 0.7, 1.2, _rgb(255, 111, 97)),
        (11.6, 0.5, 0.8, _rgb(63, 188, 171)),
        (10.9, 5.8, 1.6, _rgb(255, 196, 61)),
    ):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
    _add_textbox(
        slide,
        title,
        left=1.45,
        top=2.15,
        width=10.4,
        height=1.5,
        size=34,
        color=_rgb(39, 54, 65),
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide,
        str(subtitle or "快乐成长 · 放飞梦想"),
        left=2.5,
        top=3.75,
        width=8.3,
        height=0.7,
        size=17,
        color=_rgb(63, 142, 132),
        align=PP_ALIGN.CENTER,
    )


def _add_ppt_content_slide(slide: Any, item: dict[str, Any], index: int) -> None:
    _set_slide_background(slide, _rgb(252, 252, 249))
    accent = (_rgb(255, 111, 97), _rgb(63, 188, 171), _rgb(255, 196, 61))[index % 3]
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.24), Inches(7.5)
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()
    title = str(item.get("title") or f"第 {index} 部分")
    _add_textbox(
        slide,
        title,
        left=0.85,
        top=0.55,
        width=11.4,
        height=0.8,
        size=26,
        color=_rgb(39, 54, 65),
        bold=True,
    )
    body = str(item.get("body") or "").strip()
    bullets = item.get("bullets") or []
    lines = ([body] if body else []) + [f"•  {value}" for value in bullets if str(value).strip()]
    _add_textbox(
        slide,
        "\n\n".join(lines) or "内容待补充",
        left=1.0,
        top=1.55,
        width=11.1,
        height=4.9,
        size=19,
        color=_rgb(61, 72, 78),
    )
    _add_textbox(
        slide,
        f"{index:02d}",
        left=11.7,
        top=6.65,
        width=0.7,
        height=0.4,
        size=11,
        color=accent,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )


def _write_pdf(path: Path, *, title: str, content: str) -> None:
    story = fitz.Story(
        html=_markdown_document(title=title, content=content),
        user_css=_PDF_CSS,
        em=11,
    )
    writer = fitz.DocumentWriter(str(path))
    page_rect = fitz.paper_rect("a4")
    content_rect = fitz.Rect(
        page_rect.x0 + 48,
        page_rect.y0 + 48,
        page_rect.x1 - 48,
        page_rect.y1 - 50,
    )
    more = True
    pages = 0
    try:
        while more:
            if pages >= 500:
                raise ValueError("PDF 内容过长，超过 500 页限制")
            device = writer.begin_page(page_rect)
            more, _ = story.place(content_rect)
            story.draw(device)
            writer.end_page()
            pages += 1
    finally:
        writer.close()

    document = fitz.open(path)
    try:
        metadata = document.metadata
        metadata["title"] = title
        document.set_metadata(metadata)
        for index, page in enumerate(document, start=1):
            page.insert_text(
                (page.rect.width - 72, page.rect.height - 25),
                f"{index} / {document.page_count}",
                fontname="helv",
                fontsize=8,
                color=(0.42, 0.47, 0.45),
            )
        document.saveIncr()
    finally:
        document.close()


_PDF_CSS = """
body { color: #24312e; font-family: sans-serif; font-size: 10.5pt; line-height: 1.55; }
h1 { color: #123e35; font-size: 20pt; line-height: 1.2; margin: 0 0 16pt 0;
     padding: 0 0 8pt 0; border-bottom: 1pt solid #b9d2cb; }
h2 { color: #185b4d; font-size: 16pt; line-height: 1.25; margin: 17pt 0 7pt 0; }
h3 { color: #276d5f; font-size: 13pt; line-height: 1.3; margin: 13pt 0 5pt 0; }
h4, h5, h6 { color: #315f56; font-size: 11pt; margin: 10pt 0 4pt 0; }
p { margin: 0 0 7pt 0; }
ul, ol { margin: 3pt 0 9pt 17pt; padding: 0; }
li { margin: 0 0 3pt 0; }
blockquote { color: #41645c; background-color: #eef6f3; border-left: 3pt solid #76aa9d;
             margin: 8pt 0; padding: 7pt 10pt; }
code { color: #7b3f16; background-color: #f3f1ed; font-family: monospace; font-size: 9pt; }
pre { color: #f3f7f5; background-color: #21312d; font-family: monospace; font-size: 8.5pt;
      line-height: 1.35; margin: 8pt 0 11pt 0; padding: 9pt; }
table { border-collapse: collapse; width: 100%; margin: 8pt 0 12pt 0; font-size: 9.3pt; }
th { color: #ffffff; background-color: #23705f; font-weight: bold; }
th, td { border: 0.6pt solid #b8c9c4; padding: 5pt 6pt; vertical-align: top; }
tr:nth-child(even) td { background-color: #f4f8f6; }
hr { color: #cbd9d5; margin: 11pt 0; }
"""


_TABLE_SEPARATOR = re.compile(r"^\s*\|?\s*:?-{3,}:?\s*(?:\|\s*:?-{3,}:?\s*)+\|?\s*$")
_ORDERED_ITEM = re.compile(r"^\s*\d+[.)]\s+(.+)$")
_BULLET_ITEM = re.compile(r"^\s*[-+*]\s+(.+)$")
_HEADING = re.compile(r"^(#{1,6})\s+(.+)$")


def _markdown_document(*, title: str, content: str) -> str:
    """把常用 Markdown 安全转换成适合打印的 HTML。"""

    lines = content.replace("\r\n", "\n").replace("\r", "\n").split("\n")
    for position, line in enumerate(lines):
        if not line.strip():
            continue
        first_heading = _HEADING.match(line.strip())
        if first_heading is not None and first_heading.group(1) == "#":
            # title 参数是交付物的唯一文档标题；模型常会在 content 再写一个 H1，
            # 即使措辞略有不同也不应形成双标题。
            lines[position] = ""
        break
    blocks: list[str] = [f"<h1>{_inline_markdown(title)}</h1>"]
    index = 0
    title_key = _plain_markdown(title).casefold()
    while index < len(lines):
        raw = lines[index]
        stripped = raw.strip()
        if not stripped:
            index += 1
            continue

        if stripped.startswith("```"):
            code: list[str] = []
            index += 1
            while index < len(lines) and not lines[index].strip().startswith("```"):
                code.append(lines[index])
                index += 1
            if index < len(lines):
                index += 1
            blocks.append(f"<pre>{html.escape(chr(10).join(code))}</pre>")
            continue

        heading = _HEADING.match(stripped)
        if heading is not None:
            text = heading.group(2).strip()
            if not (heading.group(1) == "#" and _plain_markdown(text).casefold() == title_key):
                level = len(heading.group(1))
                blocks.append(f"<h{level}>{_inline_markdown(text)}</h{level}>")
            index += 1
            continue

        if stripped in {"---", "***", "___"}:
            blocks.append("<hr>")
            index += 1
            continue

        if index + 1 < len(lines) and "|" in raw and _TABLE_SEPARATOR.match(lines[index + 1]):
            table_lines = [raw]
            index += 2
            while index < len(lines) and "|" in lines[index] and lines[index].strip():
                table_lines.append(lines[index])
                index += 1
            blocks.append(_markdown_table(table_lines))
            continue

        if _BULLET_ITEM.match(raw) is not None:
            items: list[str] = []
            while index < len(lines):
                match = _BULLET_ITEM.match(lines[index])
                if match is None:
                    break
                items.append(f"<li>{_inline_markdown(match.group(1))}</li>")
                index += 1
            blocks.append(f"<ul>{''.join(items)}</ul>")
            continue

        if _ORDERED_ITEM.match(raw) is not None:
            items = []
            while index < len(lines):
                match = _ORDERED_ITEM.match(lines[index])
                if match is None:
                    break
                items.append(f"<li>{_inline_markdown(match.group(1))}</li>")
                index += 1
            blocks.append(f"<ol>{''.join(items)}</ol>")
            continue

        if stripped.startswith(">"):
            quoted: list[str] = []
            while index < len(lines) and lines[index].lstrip().startswith(">"):
                quoted.append(lines[index].lstrip()[1:].strip())
                index += 1
            blocks.append(f"<blockquote>{_inline_markdown(' '.join(quoted))}</blockquote>")
            continue

        paragraph = [stripped]
        index += 1
        while (
            index < len(lines) and lines[index].strip() and not _starts_markdown_block(lines, index)
        ):
            paragraph.append(lines[index].strip())
            index += 1
        blocks.append(f"<p>{_inline_markdown(' '.join(paragraph))}</p>")
    return "<!doctype html><html><body>" + "".join(blocks) + "</body></html>"


def _starts_markdown_block(lines: list[str], index: int) -> bool:
    stripped = lines[index].strip()
    return bool(
        stripped.startswith(("```", ">"))
        or _HEADING.match(stripped)
        or _BULLET_ITEM.match(lines[index])
        or _ORDERED_ITEM.match(lines[index])
        or stripped in {"---", "***", "___"}
        or (
            index + 1 < len(lines)
            and "|" in lines[index]
            and _TABLE_SEPARATOR.match(lines[index + 1])
        )
    )


def _markdown_table(lines: list[str]) -> str:
    rows = [[cell.strip() for cell in line.strip().strip("|").split("|")] for line in lines]
    width = max(len(row) for row in rows)
    normalized = [row + [""] * (width - len(row)) for row in rows]
    head = "".join(f"<th>{_inline_markdown(cell)}</th>" for cell in normalized[0])
    body = "".join(
        "<tr>" + "".join(f"<td>{_inline_markdown(cell)}</td>" for cell in row) + "</tr>"
        for row in normalized[1:]
    )
    return f"<table><thead><tr>{head}</tr></thead><tbody>{body}</tbody></table>"


def _inline_markdown(value: str) -> str:
    escaped = html.escape(value, quote=True)
    escaped = re.sub(r"`([^`]+)`", r"<code>\1</code>", escaped)
    escaped = re.sub(r"\*\*([^*]+)\*\*", r"<strong>\1</strong>", escaped)
    escaped = re.sub(r"__([^_]+)__", r"<strong>\1</strong>", escaped)
    escaped = re.sub(r"(?<!\*)\*([^*]+)\*(?!\*)", r"<em>\1</em>", escaped)
    escaped = re.sub(r"(?<!_)_([^_]+)_(?!_)", r"<em>\1</em>", escaped)
    escaped = re.sub(r"\[([^]]+)]\((https?://[^\s)]+)\)", r'<a href="\2">\1</a>', escaped)
    return escaped


def _plain_markdown(value: str) -> str:
    return re.sub(r"[*_`]", "", value).strip()
