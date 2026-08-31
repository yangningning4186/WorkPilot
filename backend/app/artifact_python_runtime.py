"""随 WorkPilot 发布的最小 artifact Python 命令行运行时。

该入口由 PyInstaller 冻结成独立可执行文件。它只提供生成 Office/PDF 产物所需的 Python
执行语义，不承载 API、模型凭据或用户数据库，因此可以作为原生沙箱中的低权限子进程运行。
"""

from __future__ import annotations

import json
import runpy
import sys
import tempfile
from importlib import metadata
from pathlib import Path
from typing import Any, NoReturn

RUNTIME_PROFILE = "artifact-python:1.0.0"
EXPECTED_DISTRIBUTIONS = {
    "python-pptx": "1.0.2",
    "openpyxl": "3.1.5",
    "python-docx": "1.2.0",
    "reportlab": "5.0.1",
    "PyMuPDF": "1.28.2",
    "Pillow": "12.3.0",
    "XlsxWriter": "3.2.9",
}


def _runtime_info() -> dict[str, object]:
    installed = {
        name: metadata.version(name)
        for name in EXPECTED_DISTRIBUTIONS
    }
    return {
        "profile": RUNTIME_PROFILE,
        "python": sys.version.split()[0],
        "dependencies": installed,
    }


def _assert_runtime_versions() -> None:
    actual = _runtime_info()["dependencies"]
    assert isinstance(actual, dict)
    mismatches = {
        name: {"expected": expected, "actual": actual.get(name)}
        for name, expected in EXPECTED_DISTRIBUTIONS.items()
        if actual.get(name) != expected
    }
    if mismatches:
        raise RuntimeError(f"artifact runtime 依赖版本不一致: {mismatches}")


def _selftest() -> None:
    _assert_runtime_versions()
    import pymupdf
    from docx import Document
    from openpyxl import Workbook, load_workbook  # type: ignore[import-untyped]
    from pptx import Presentation
    from reportlab.pdfgen.canvas import Canvas  # type: ignore[import-untyped]

    pymupdf_api: Any = pymupdf

    with tempfile.TemporaryDirectory(prefix="workpilot-artifact-runtime-") as raw:
        root = Path(raw)

        presentation_path = root / "probe.pptx"
        presentation = Presentation()
        presentation.slides.add_slide(presentation.slide_layouts[6])
        presentation.save(str(presentation_path))
        assert len(Presentation(str(presentation_path)).slides) == 1

        workbook_path = root / "probe.xlsx"
        workbook = Workbook()
        workbook.active["A1"] = "WorkPilot"
        workbook.save(workbook_path)
        assert load_workbook(workbook_path, read_only=True).active["A1"].value == "WorkPilot"

        document_path = root / "probe.docx"
        document = Document()
        document.add_paragraph("WorkPilot")
        document.save(str(document_path))
        assert Document(str(document_path)).paragraphs[0].text == "WorkPilot"

        pdf_path = root / "probe.pdf"
        canvas = Canvas(str(pdf_path))
        canvas.drawString(72, 720, "WorkPilot")
        canvas.save()
        with pymupdf_api.open(pdf_path) as pdf:
            assert pdf.page_count == 1


def _usage_error(message: str) -> NoReturn:
    raise SystemExit(f"{message}\n用法: workpilot-artifact-python [-c code | -m module | script.py] [args...]")


def _run_code(arguments: list[str]) -> int:
    if len(arguments) < 2:
        _usage_error("-c 后缺少 Python 代码")
    sys.argv = ["-c", *arguments[2:]]
    sys.path.insert(0, str(Path.cwd()))
    namespace = {"__name__": "__main__", "__package__": None, "__builtins__": __builtins__}
    exec(compile(arguments[1], "<string>", "exec"), namespace, namespace)
    return 0


def _run_module(arguments: list[str]) -> int:
    if len(arguments) < 2:
        _usage_error("-m 后缺少模块名")
    sys.argv = [arguments[1], *arguments[2:]]
    sys.path.insert(0, str(Path.cwd()))
    runpy.run_module(arguments[1], run_name="__main__", alter_sys=True)
    return 0


def _run_script(arguments: list[str]) -> int:
    script = Path(arguments[0]).expanduser().resolve(strict=True)
    if not script.is_file():
        _usage_error(f"脚本不是普通文件: {script}")
    sys.argv = [str(script), *arguments[1:]]
    sys.path.insert(0, str(script.parent))
    runpy.run_path(str(script), run_name="__main__")
    return 0


def main(argv: list[str] | None = None) -> int:
    arguments = list(sys.argv[1:] if argv is None else argv)
    if not arguments:
        _usage_error("缺少要执行的代码或脚本")
    if arguments == ["--workpilot-runtime-info"]:
        print(json.dumps(_runtime_info(), ensure_ascii=False, sort_keys=True))
        return 0
    if arguments == ["--workpilot-selftest"]:
        _selftest()
        print(json.dumps({"ok": True, **_runtime_info()}, ensure_ascii=False, sort_keys=True))
        return 0
    if arguments[0] in {"-V", "--version"}:
        print(f"Python {sys.version.split()[0]} ({RUNTIME_PROFILE})")
        return 0
    if arguments[0] == "-c":
        return _run_code(arguments)
    if arguments[0] == "-m":
        return _run_module(arguments)
    if arguments[0].startswith("-"):
        _usage_error(f"不支持的解释器参数: {arguments[0]}")
    return _run_script(arguments)


if __name__ == "__main__":  # pragma: no cover - 独立运行时入口
    raise SystemExit(main())
