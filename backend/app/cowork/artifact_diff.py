"""把工作区交付物转换成有界、可在右栏审阅的语义 diff。"""

from __future__ import annotations

import difflib
import hashlib
from io import BytesIO
from pathlib import Path
from typing import Any

import fitz  # type: ignore[import-untyped]
from docx import Document
from openpyxl import load_workbook  # type: ignore[import-untyped]
from pptx import Presentation

MAX_DIFF_SOURCE_BYTES = 2 * 1024 * 1024
MAX_DIFF_TEXT_CHARS = 120_000
MAX_DIFF_LINES = 500
MAX_DIFF_CHARS = 48_000

_TEXT_SUFFIXES = frozenset(
    {".txt", ".md", ".markdown", ".csv", ".tsv", ".json", ".xml", ".yaml", ".yml", ".html"}
)
_DIFFABLE_SUFFIXES = _TEXT_SUFFIXES | {".docx", ".xlsx", ".pptx", ".pdf"}


def capture_artifact_baseline(path: Path, *, remaining_bytes: int) -> bytes | None:
    """为即将执行的命令保留一份内存基线；大文件不拖慢每次 Shell。"""

    try:
        size = path.stat().st_size
    except OSError:
        return None
    if (
        path.suffix.casefold() not in _DIFFABLE_SUFFIXES
        or size > MAX_DIFF_SOURCE_BYTES
        or size > remaining_bytes
    ):
        return None
    try:
        return path.read_bytes()
    except OSError:
        return None


def build_artifact_diff(
    *,
    after_path: Path,
    before_bytes: bytes | None = None,
    before_path: Path | None = None,
    created: bool = False,
) -> dict[str, Any]:
    """生成统一 diff；Office/PDF 先抽成稳定的语义行，再比较。"""

    suffix = after_path.suffix.casefold()
    if suffix not in _DIFFABLE_SUFFIXES:
        return _unavailable("该格式没有安全的文本差异视图")
    try:
        after_bytes = after_path.read_bytes()
    except OSError:
        return _unavailable("无法读取交付物")
    if len(after_bytes) > MAX_DIFF_SOURCE_BYTES:
        return _unavailable("文件超过 2 MB 差异快照上限")
    if created:
        baseline = b""
    elif before_bytes is not None:
        baseline = before_bytes
    elif before_path is not None and before_path.is_file():
        try:
            baseline = before_path.read_bytes()
        except OSError:
            return _unavailable("执行前副本不可读")
    else:
        return _unavailable("执行前没有可比较快照")
    if len(baseline) > MAX_DIFF_SOURCE_BYTES:
        return _unavailable("执行前文件超过 2 MB 差异快照上限")

    try:
        before_text = "" if created else _semantic_text(baseline, suffix)
        after_text = _semantic_text(after_bytes, suffix)
    # 这里跨的是四套第三方文档解析器的信任边界；它们对损坏 OOXML/PDF 抛出的异常类型
    # 并不统一。diff 是辅助审阅信息，解析器拒绝旧快照不能反过来让已成功的 Shell 动作失败。
    except Exception as error:
        return _unavailable(f"无法提取差异：{error}")

    before_lines = before_text.splitlines()
    after_lines = after_text.splitlines()
    rendered = list(
        difflib.unified_diff(
            before_lines,
            after_lines,
            fromfile="修改前",
            tofile="修改后",
            lineterm="",
        )
    )
    added = sum(line.startswith("+") and not line.startswith("+++") for line in rendered)
    removed = sum(line.startswith("-") and not line.startswith("---") for line in rendered)
    text, truncated = _bounded_diff(rendered)
    return {
        "schema_version": 1,
        "available": True,
        "format": "unified",
        "view": "semantic" if suffix in {".docx", ".xlsx", ".pptx", ".pdf"} else "text",
        "created": created,
        "before_sha256": None if created else hashlib.sha256(baseline).hexdigest(),
        "after_sha256": hashlib.sha256(after_bytes).hexdigest(),
        "added_lines": added,
        "removed_lines": removed,
        "truncated": truncated,
        "text": text,
        "reason": (
            "文件字节发生变化，但抽取出的文字与公式没有变化"
            if not created and not rendered and baseline != after_bytes
            else None
        ),
    }


def _unavailable(reason: str) -> dict[str, Any]:
    return {
        "schema_version": 1,
        "available": False,
        "format": "unified",
        "view": "unavailable",
        "created": False,
        "before_sha256": None,
        "after_sha256": None,
        "added_lines": 0,
        "removed_lines": 0,
        "truncated": False,
        "text": "",
        "reason": reason,
    }


def _bounded_diff(lines: list[str]) -> tuple[str, bool]:
    selected: list[str] = []
    size = 0
    truncated = False
    for line in lines:
        addition = len(line) + 1
        if len(selected) >= MAX_DIFF_LINES or size + addition > MAX_DIFF_CHARS:
            truncated = True
            break
        selected.append(line)
        size += addition
    return ("\n".join(selected), truncated)


def _semantic_text(raw: bytes, suffix: str) -> str:
    if suffix in _TEXT_SUFFIXES:
        text = raw.decode("utf-8-sig")
    elif suffix == ".docx":
        document = Document(BytesIO(raw))
        lines = [paragraph.text for paragraph in document.paragraphs if paragraph.text]
        for table_index, table in enumerate(document.tables, start=1):
            lines.append(f"[表格 {table_index}]")
            lines.extend(" | ".join(cell.text for cell in row.cells) for row in table.rows)
        text = "\n".join(lines)
    elif suffix == ".xlsx":
        workbook = load_workbook(BytesIO(raw), read_only=True, data_only=False)
        try:
            lines = []
            for sheet in workbook.worksheets:
                lines.append(f"[工作表] {sheet.title}")
                for row in sheet.iter_rows():
                    values = [
                        f"{cell.coordinate}={cell.value}" for cell in row if cell.value is not None
                    ]
                    if values:
                        lines.append(" | ".join(values))
            text = "\n".join(lines)
        finally:
            workbook.close()
    elif suffix == ".pptx":
        presentation = Presentation(BytesIO(raw))
        lines = []
        for index, slide in enumerate(presentation.slides, start=1):
            lines.append(f"[幻灯片 {index}]")
            lines.extend(
                str(value) for shape in slide.shapes if (value := getattr(shape, "text", ""))
            )
        text = "\n".join(lines)
    elif suffix == ".pdf":
        document = fitz.open(stream=raw, filetype="pdf")
        try:
            text = "\n".join(
                f"[第 {index + 1} 页]\n{page.get_text('text')}"
                for index, page in enumerate(document)
            )
        finally:
            document.close()
    else:  # pragma: no cover - 调用前已封闭 suffix
        raise ValueError("不支持的差异格式")
    if len(text) > MAX_DIFF_TEXT_CHARS:
        return text[:MAX_DIFF_TEXT_CHARS] + "\n[语义视图已截断]"
    return text


__all__ = [
    "MAX_DIFF_SOURCE_BYTES",
    "build_artifact_diff",
    "capture_artifact_baseline",
]
