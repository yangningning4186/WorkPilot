"""独立于 Skill/Renderer 的 Artifact Validator 与质量评分。"""

from __future__ import annotations

import json
import os
import re
import signal
import subprocess
import sys
import tempfile
import time
import xml.etree.ElementTree as ET
import zipfile
from html.parser import HTMLParser
from itertools import pairwise
from pathlib import Path
from typing import Any, Literal
from unicodedata import east_asian_width
from urllib.parse import urlsplit

import pymupdf
from docx import Document
from docx.oxml.ns import qn
from openpyxl import load_workbook  # type: ignore[import-untyped]
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR
from pydantic import BaseModel, ConfigDict, Field

from app.cowork.artifact_renderers.contracts import (
    ArtifactSpec,
    DocumentSpec,
    PresentationSpec,
    WorkbookSpec,
)
from app.cowork.office_preview import OfficePreviewError, render_office_preview
from app.cowork.process_limits import read_process_tree_usage
from app.cowork.skills.builtin.pptx.scripts.pptx2image import (
    PptxRasterError,
    render_presentation_pages,
)

_pymupdf: Any = pymupdf

ValidationStatus = Literal["passed", "warning", "failed", "not_run"]
_FORMULA_ERRORS = ("#REF!", "#DIV/0!", "#VALUE!", "#NAME?")
_NETWORK_FORMULA = re.compile(
    r"(?<![A-Z0-9_.])(?:_XLFN\.)?"
    r"(?:WEBSERVICE|RTD|STOCKHISTORY|IMAGE|CUBEMEMBER|CUBEVALUE|CUBESET)\s*\(",
    re.IGNORECASE,
)
_CSS_REMOTE_URL = re.compile(
    r"url\s*\(\s*(['\"]?)(?P<url>[^)'\"]+)\1\s*\)|"
    r"@import\s+(?:url\s*\(\s*)?(['\"]?)(?P<import>[^;'\")\s]+)",
    re.IGNORECASE,
)
_REMOTE_SCHEMES = frozenset({"http", "https", "ftp", "ftps", "ws", "wss", "file"})
_URL_ATTRIBUTES = frozenset(
    {
        "action",
        "background",
        "cite",
        "data",
        "formaction",
        "href",
        "manifest",
        "ping",
        "poster",
        "profile",
        "src",
        "longdesc",
        "codebase",
        "archive",
        "usemap",
        "xlink:href",
    }
)
_FORBIDDEN_HTML_TAGS = frozenset({"base", "embed", "iframe", "object", "script"})
_OOXML_MAX_ENTRIES = 10_000
_OOXML_MAX_EXPANDED_BYTES = 256 * 1024 * 1024
_OOXML_MAX_EXPANSION_RATIO = 250


def _is_remote_reference(value: str) -> bool:
    normalized = value.strip().strip("'\"")
    if not normalized:
        return False
    if normalized.startswith("//"):
        return True
    try:
        return urlsplit(normalized).scheme.casefold() in _REMOTE_SCHEMES
    except ValueError:
        return True


def _css_remote_references(value: str) -> list[str]:
    findings: list[str] = []
    for match in _CSS_REMOTE_URL.finditer(value):
        candidate = match.group("url") or match.group("import") or ""
        if _is_remote_reference(candidate):
            findings.append(candidate[:240])
    return findings


class _OfflineHtmlSecurityParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.findings: list[str] = []
        self._style_depth = 0

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        lowered = tag.casefold()
        if lowered in _FORBIDDEN_HTML_TAGS:
            self.findings.append(f"tag:{lowered}")
        if lowered == "style":
            self._style_depth += 1
        attributes = {name.casefold(): value or "" for name, value in attrs}
        if any(name.startswith("on") for name in attributes):
            self.findings.append(f"event:{lowered}")
        if lowered == "meta" and attributes.get("http-equiv", "").casefold() == "refresh":
            self.findings.append("meta:refresh")
        for name, value in attributes.items():
            if name == "style":
                self.findings.extend(f"css:{item}" for item in _css_remote_references(value))
            elif name == "srcset":
                for candidate in value.split(","):
                    reference = candidate.strip().split(maxsplit=1)[0] if candidate.strip() else ""
                    if _is_remote_reference(reference):
                        self.findings.append(f"srcset:{reference[:240]}")
            elif name in _URL_ATTRIBUTES and _is_remote_reference(value):
                self.findings.append(f"{name}:{value[:240]}")
            elif "javascript:" in value.casefold():
                self.findings.append(f"javascript:{name}")

    def handle_startendtag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        self.handle_starttag(tag, attrs)
        self.handle_endtag(tag)

    def handle_endtag(self, tag: str) -> None:
        if tag.casefold() == "style":
            self._style_depth = max(0, self._style_depth - 1)

    def handle_data(self, data: str) -> None:
        if self._style_depth:
            self.findings.extend(f"css:{item}" for item in _css_remote_references(data))


class _StrictReport(BaseModel):
    model_config = ConfigDict(extra="forbid")


class ValidationCheck(_StrictReport):
    name: str
    status: ValidationStatus
    message: str
    value: int | float | str | bool | None = None


class ValidationDimension(_StrictReport):
    status: ValidationStatus
    checks: list[ValidationCheck] = Field(default_factory=list)


class ArtifactQuality(_StrictReport):
    score: int = Field(ge=0, le=100)
    warnings: list[str] = Field(default_factory=list)
    #: 各维度的独立得分，让"哪一维在掉分"可比较。综合分只是它们的加权和，
    #: 单看一个 0-100 分不知道该往哪修。
    dimension_scores: dict[str, int] = Field(default_factory=dict)
    #: 触发一票否决的原因。非空即表示综合分被强制归零，与加权结果无关。
    blocking: list[str] = Field(default_factory=list)


class ArtifactValidationReport(_StrictReport):
    schema_version: Literal[1] = 1
    artifact_type: Literal["docx", "xlsx", "pptx", "pdf", "html"]
    structural: ValidationDimension
    semantic: ValidationDimension
    visual: ValidationDimension
    evidence: ValidationDimension
    security: ValidationDimension
    quality: ArtifactQuality

    @property
    def deliverable(self) -> bool:
        if self.artifact_type == "pptx" and self.visual.status != "passed":
            return False
        return all(
            dimension.status != "failed"
            for dimension in (
                self.structural,
                self.semantic,
                self.visual,
                self.evidence,
                self.security,
            )
        )


def validate_artifact_in_subprocess(
    path: Path,
    *,
    render_visual: bool = False,
    timeout_s: float = 30.0,
    memory_mb: int = 512,
    cpu_seconds: int = 20,
    pids_limit: int = 32,
    max_file_bytes: int = 200 * 1024 * 1024,
) -> ArtifactValidationReport:
    """在有墙钟、CPU、内存和进程数上限的子进程中解析不受信产物。"""

    resolved = path.resolve(strict=True)
    command = (
        sys.executable,
        "-m",
        "app.cowork.artifact_validation_worker",
        str(resolved),
        "1" if render_visual else "0",
        str(memory_mb),
        str(cpu_seconds),
        str(pids_limit),
        str(max_file_bytes),
    )
    creationflags = 0
    if os.name == "nt":  # pragma: no cover - Windows CI 不在当前矩阵
        subprocess_api: Any = subprocess
        creationflags = int(subprocess_api.CREATE_NEW_PROCESS_GROUP)
    process = subprocess.Popen(
        command,
        cwd=Path(__file__).resolve().parents[2],
        env={
            "PATH": os.defpath,
            "PYTHONDONTWRITEBYTECODE": "1",
            "PYTHONNOUSERSITE": "1",
        },
        stdin=subprocess.DEVNULL,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        start_new_session=os.name == "posix",
        creationflags=creationflags,
    )
    def terminate() -> None:
        if os.name == "posix":
            try:
                os.killpg(process.pid, signal.SIGKILL)
            except ProcessLookupError:
                pass
        else:  # pragma: no cover - Windows CI 不在当前矩阵
            process.kill()
        process.wait()

    deadline = time.monotonic() + timeout_s
    stdout = b""
    stderr = b""
    while True:
        remaining = deadline - time.monotonic()
        if remaining <= 0:
            terminate()
            raise ValueError(f"Artifact 校验超过 {timeout_s:g} 秒，已终止受限子进程")
        try:
            stdout, stderr = process.communicate(timeout=min(0.1, remaining))
            break
        except subprocess.TimeoutExpired:
            if os.name != "posix":
                continue
            try:
                usage = read_process_tree_usage(process.pid)
            except (OSError, subprocess.SubprocessError, ValueError) as error:
                terminate()
                raise ValueError(f"无法验证 Artifact 校验进程资源用量：{error}") from error
            violation = (
                f"内存 {usage.rss_bytes} bytes 超过 {memory_mb * 1024 * 1024} bytes"
                if usage.rss_bytes > memory_mb * 1024 * 1024
                else f"进程数 {usage.pids} 超过 {pids_limit}"
                if usage.pids > pids_limit
                else f"CPU 时间 {usage.cpu_seconds:g} 秒超过 {cpu_seconds} 秒"
                if usage.cpu_seconds > cpu_seconds
                else None
            )
            if violation is not None:
                terminate()
                raise ValueError(
                    f"Artifact 校验子进程资源超限（{violation}），已终止"
                ) from None
    if len(stdout) > 4 * 1024 * 1024 or len(stderr) > 64 * 1024:
        raise ValueError("Artifact 校验子进程输出超过上限")
    try:
        payload = json.loads(stdout)
    except (UnicodeDecodeError, json.JSONDecodeError) as error:
        detail = stderr.decode("utf-8", errors="replace")[-2_000:]
        raise ValueError(f"Artifact 校验子进程返回无效结果：{detail or process.returncode}") from error
    if process.returncode != 0 or not payload.get("ok"):
        raise ValueError(str(payload.get("error") or "Artifact 校验子进程失败"))
    return ArtifactValidationReport.model_validate(payload.get("report"))


def _dimension(checks: list[ValidationCheck], *, empty: ValidationStatus = "not_run") -> ValidationDimension:
    if not checks:
        return ValidationDimension(status=empty)
    statuses = {check.status for check in checks}
    status: ValidationStatus = (
        "failed"
        if "failed" in statuses
        else "warning"
        if "warning" in statuses
        else "not_run"
        if "not_run" in statuses
        else "passed"
    )
    return ValidationDimension(status=status, checks=checks)


def _validate_ooxml_container_budget(path: Path) -> None:
    """在 Office 解析器分配解压内存前先拒绝 zip bomb 与异常容器。"""

    with zipfile.ZipFile(path) as archive:
        entries = archive.infolist()
        if len(entries) > _OOXML_MAX_ENTRIES:
            raise ValueError(f"OOXML 条目数超过上限 {_OOXML_MAX_ENTRIES}")
        expanded = 0
        compressed = 0
        for info in entries:
            normalized = info.filename.replace("\\", "/")
            parts = tuple(part for part in normalized.split("/") if part)
            if normalized.startswith("/") or ".." in parts:
                raise ValueError("OOXML 包含逃逸路径")
            if info.flag_bits & 0x1:
                raise ValueError("OOXML 包含加密条目")
            expanded += info.file_size
            compressed += max(1, info.compress_size)
            if expanded > _OOXML_MAX_EXPANDED_BYTES:
                raise ValueError(
                    f"OOXML 解压后大小超过上限 {_OOXML_MAX_EXPANDED_BYTES} bytes"
                )
        if (
            expanded > 16 * 1024 * 1024
            and expanded > compressed * _OOXML_MAX_EXPANSION_RATIO
        ):
            raise ValueError("OOXML 压缩比异常，疑似 zip bomb")


def _external_ooxml_relationships(path: Path) -> list[str]:
    found: list[str] = []
    with zipfile.ZipFile(path) as archive:
        for name in archive.namelist():
            if not name.endswith(".rels"):
                continue
            info = archive.getinfo(name)
            if info.file_size > 2 * 1024 * 1024:
                raise ValueError(f"OOXML relationship 文件过大：{name}")
            try:
                root = ET.fromstring(archive.read(name))
            except ET.ParseError as error:
                raise ValueError(f"OOXML relationship XML 无效：{name}") from error
            if any(
                str(relationship.attrib.get("TargetMode", "")).casefold() == "external"
                for relationship in root.iter()
            ):
                found.append(name)
    return found


def _render_check(path: Path) -> ValidationCheck:
    if path.suffix.casefold() == ".pptx":
        try:
            with tempfile.TemporaryDirectory(prefix="workpilot-pptx-qa-") as raw_cache:
                result = render_presentation_pages(path, Path(raw_cache))
        except (OSError, ValueError, PptxRasterError) as error:
            return ValidationCheck(
                name="render",
                status="failed",
                message=f"内置 PPTX 逐页渲染失败：{error}",
            )
        failures: list[str] = []
        if result.overflow_shapes:
            failures.append(f"文本溢出对象 {len(result.overflow_shapes)}")
        if result.unsupported_shapes:
            failures.append(f"不支持对象 {len(result.unsupported_shapes)}")
        if failures:
            details = tuple((*result.overflow_shapes, *result.unsupported_shapes))[:5]
            suffix = "；" + "、".join(details) if details else ""
            return ValidationCheck(
                name="render",
                status="failed",
                message="PPTX 逐页视觉检查失败：" + "，".join(failures) + suffix,
                value=len(result.pages),
            )
        return ValidationCheck(
            name="render",
            status="passed",
            message=f"内置 Renderer 已逐页渲染并检查 {len(result.pages)} 页",
            value=len(result.pages),
        )
    try:
        with tempfile.TemporaryDirectory(prefix="workpilot-visual-qa-") as raw_cache:
            preview = render_office_preview(
                path,
                cache_root=Path(raw_cache),
                timeout_s=30,
                max_source_bytes=64 * 1024 * 1024,
                max_cache_entries=2,
            )
            if preview is None:
                return ValidationCheck(
                    name="render",
                    status="warning",
                    message="本机没有可用的 Quick Look/LibreOffice 版面渲染器",
                )
            return ValidationCheck(
                name="render",
                status="passed",
                message=f"已生成真实版面预览（{preview.mode}）",
            )
    except (OSError, OfficePreviewError, ValueError) as error:
        return ValidationCheck(
            name="render",
            status="warning",
            message=f"真实版面渲染器未能完成检查：{error}",
        )


def _pptx_checks(
    path: Path,
    spec: PresentationSpec | None,
) -> tuple[list[ValidationCheck], list[ValidationCheck], list[ValidationCheck], list[ValidationCheck]]:
    presentation = Presentation(str(path))
    structural = [ValidationCheck(name="reopen", status="passed", message="PPTX 可重新打开")]
    semantic: list[ValidationCheck] = []
    visual: list[ValidationCheck] = []
    security: list[ValidationCheck] = []
    blank = 0
    missing_title = 0
    out_of_bounds = 0
    overlaps = 0
    overloaded = 0
    sparse = 0
    hollow_container_pages = 0
    focus_like_pages = 0
    rendered_fingerprints: list[tuple[tuple[int, int, int, int, int], ...]] = []
    min_font: float | None = None
    slide_width = int(presentation.slide_width or 0)
    slide_height = int(presentation.slide_height or 0)
    for slide_index, slide in enumerate(presentation.slides):
        text_shapes: list[tuple[int, int, int, int]] = []
        text_shape_values: list[tuple[tuple[int, int, int, int], str]] = []
        large_containers: list[tuple[int, int, int, int]] = []
        fingerprint: list[tuple[int, int, int, int, int]] = []
        total_text = 0
        has_title = False
        meaningful_shapes = 0
        visual_objects = 0
        large_centered_text = 0
        native_content_nodes = 0
        for shape in slide.shapes:
            left = int(shape.left or 0)
            top = int(shape.top or 0)
            width = int(shape.width or 0)
            height = int(shape.height or 0)
            if width > 0 and height > 0:
                fingerprint.append(
                    (
                        int(shape.shape_type),
                        round(left / slide_width * 20) if slide_width else 0,
                        round(top / slide_height * 20) if slide_height else 0,
                        round(width / slide_width * 20) if slide_width else 0,
                        round(height / slide_height * 20) if slide_height else 0,
                    )
                )
            if (
                shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                and width * height >= slide_width * slide_height * 0.08
                and not str(getattr(shape, "text", "")).strip()
            ):
                large_containers.append((left, top, width, height))
            if left < 0 or top < 0 or left + width > slide_width or top + height > slide_height:
                out_of_bounds += 1
            text = str(getattr(shape, "text", "")).strip()
            if text:
                meaningful_shapes += 1
                total_text += len(text)
                if (
                    shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                    and width * height >= slide_width * slide_height * 0.015
                ):
                    native_content_nodes += 1
                if top < slide_height * 0.28:
                    has_title = True
                if getattr(shape, "has_text_frame", False):
                    text_shapes.append((left, top, width, height))
                    text_shape_values.append(((left, top, width, height), text))
                    shape_max_font = 0.0
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.size is not None:
                                points = run.font.size.pt
                                shape_max_font = max(shape_max_font, points)
                                min_font = points if min_font is None else min(min_font, points)
                    if (
                        shape_max_font >= 26
                        and shape.text_frame.vertical_anchor == MSO_ANCHOR.MIDDLE
                    ):
                        large_centered_text += 1
            elif shape.shape_type in {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.CHART}:
                meaningful_shapes += 1
                visual_objects += 1
        if meaningful_shapes == 0:
            blank += 1
        if not has_title:
            missing_title += 1
        if total_text > 1_200 or len(text_shapes) > 14:
            overloaded += 1
        focus_signal = (
            large_centered_text == 1
            and len(text_shapes) <= 4
            and visual_objects == 0
            and total_text <= 72
        )
        focus_like = slide_index == 0 or focus_signal
        focus_like_pages += int(focus_like)
        if (
            slide_index > 0
            and visual_objects == 0
            and not focus_like
            and large_centered_text < 2
            and native_content_nodes < 2
            and (
                (len(text_shapes) <= 2 and total_text < 160)
                # 标题加四个短标签的“目录式九宫格”仍然是空洞页。旧门禁只看
                # <=4 个文本框，会把这类一张大底板上摆四个词的页面误判为
                # 有结构内容；把标题计入后应允许 5 个文本框，但仍要求有说明。
                or (len(text_shapes) <= 5 and total_text < 45)
            )
        ):
            sparse += 1
        if len(large_containers) >= 3:
            hollow = 0
            for container in large_containers:
                contained_text = "".join(
                    text
                    for box, text in text_shape_values
                    if _overlap_ratio(container, box) >= 0.65
                )
                hollow += int(len(contained_text.strip()) < 6)
            if hollow / len(large_containers) >= 0.75:
                hollow_container_pages += 1
        for index, first in enumerate(text_shapes):
            for second in text_shapes[index + 1 :]:
                if _overlap_ratio(first, second) >= 0.35:
                    overlaps += 1
        rendered_fingerprints.append(tuple(sorted(fingerprint)))
    rendered_repetition = sum(
        first == second
        for first, second in pairwise(rendered_fingerprints)
        if first and second
    )
    visual.extend(
        [
            ValidationCheck(name="blank_slides", status="failed" if blank else "passed", message=f"空白页 {blank}", value=blank),
            ValidationCheck(name="out_of_bounds", status="failed" if out_of_bounds else "passed", message=f"越界元素 {out_of_bounds}", value=out_of_bounds),
            ValidationCheck(name="large_overlap", status="failed" if overlaps else "passed", message=f"疑似大面积文本重叠 {overlaps}", value=overlaps),
            ValidationCheck(name="min_font_pt", status="failed" if min_font is not None and min_font < 16 else "passed", message="未发现小于 16 pt 的正文" if min_font is None or min_font >= 16 else f"最小字号 {min_font:g} pt", value=min_font),
            ValidationCheck(name="text_overload", status="failed" if overloaded else "passed", message=f"文字过载页 {overloaded}", value=overloaded),
            ValidationCheck(name="content_density", status="failed" if sparse else "passed", message=f"非 focus 的空洞内容页 {sparse}", value=sparse),
            ValidationCheck(
                name="hollow_containers",
                status="failed" if hollow_container_pages else "passed",
                message=f"大容器仅有标签、缺少说明的页面 {hollow_container_pages}",
                value=hollow_container_pages,
            ),
            ValidationCheck(
                name="rendered_layout_repetition",
                status="warning" if rendered_repetition else "passed",
                message=f"从最终版面检测到相邻同构页面 {rendered_repetition}",
                value=rendered_repetition,
            ),
        ]
    )
    inferred_focus_limit = max(2, (len(presentation.slides) + 2) // 3)
    visual.append(
        ValidationCheck(
            name="inferred_focus_ratio",
            status=(
                "warning"
                if len(presentation.slides) >= 5 and focus_like_pages > inferred_focus_limit
                else "passed"
            ),
            message=(
                f"从最终版面推断 focus 页 {focus_like_pages}/{len(presentation.slides)}"
                f"（建议上限 {inferred_focus_limit}）"
            ),
            value=focus_like_pages,
        )
    )
    semantic.append(ValidationCheck(name="titles", status="warning" if missing_title else "passed", message=f"缺少可识别标题的页面 {missing_title}", value=missing_title))
    repeated = 0
    if spec is not None:
        repeated = sum(first.layout == second.layout for first, second in zip(spec.slides, spec.slides[1:], strict=False))
    semantic.append(ValidationCheck(name="layout_repetition", status="warning" if repeated else "passed", message=f"相邻重复 layout {repeated}", value=repeated))
    if spec is not None:
        focus_layouts = {"title", "statement", "section", "quote"}
        focus_pages = sum(
            slide.layout in focus_layouts or slide.role in {"hero", "transition"}
            for slide in spec.slides
        )
        focus_limit = max(2, (len(spec.slides) + 2) // 3)
        focus_overuse = len(spec.slides) >= 5 and focus_pages > focus_limit
        semantic.append(
            ValidationCheck(
                name="focus_page_ratio",
                status="warning" if focus_overuse else "passed",
                message=f"focus/hero/transition 页 {focus_pages}/{len(spec.slides)}（建议上限 {focus_limit}）",
                value=focus_pages,
            )
        )

        evidence_layouts = {
            "big_number",
            "chart",
            "image_text",
            "timeline",
            "matrix",
            "cards",
            "activity",
            "diagram",
            "canvas",
        }
        supporting_pages = [
            slide
            for slide in spec.slides
            if slide.layout not in focus_layouts and slide.role == "supporting"
        ]
        evidence_pages = sum(slide.layout in evidence_layouts for slide in supporting_pages)
        evidence_floor = (len(supporting_pages) + 2) // 3
        weak_visual_story = len(supporting_pages) >= 3 and evidence_pages < evidence_floor
        semantic.append(
            ValidationCheck(
                name="visual_story",
                status="warning" if weak_visual_story else "passed",
                message=(
                    f"supporting 内容页 {len(supporting_pages)}，其中证据型视觉页 {evidence_pages}"
                    f"（建议至少 {evidence_floor}）"
                ),
                value=evidence_pages,
            )
        )

        sparse_supporting: list[str] = []
        for slide in spec.slides:
            if slide.role != "supporting" or slide.layout in focus_layouts:
                continue
            if slide.layout in {"two_column", "comparison"} and (
                len(slide.left_items) + len(slide.right_items) < 3
            ):
                sparse_supporting.append(slide.id)
            elif slide.layout == "big_number" and len(slide.metrics) == 1 and not slide.body:
                sparse_supporting.append(slide.id)
            elif slide.layout == "image_text":
                explanation = "".join(slide.bullets) + (slide.body or "") + (slide.image_caption or "")
                if len(slide.bullets) < 2 and len(explanation.strip()) < 24:
                    sparse_supporting.append(slide.id)
            elif slide.layout == "timeline" and len(slide.timeline) == 2 and all(
                not entry.detail for entry in slide.timeline
            ):
                sparse_supporting.append(slide.id)
            elif slide.layout == "matrix" and len(slide.matrix) == 2:
                sparse_supporting.append(slide.id)
            elif slide.layout == "cards" and any(
                len(card.detail.strip()) < 8 for card in slide.cards
            ):
                sparse_supporting.append(slide.id)
            elif slide.layout == "activity" and (
                sum(len(step.strip()) for step in slide.activity_steps) < 20
                or len((slide.activity_debrief or "").strip()) < 6
            ):
                sparse_supporting.append(slide.id)
        semantic.append(
            ValidationCheck(
                name="supporting_density",
                status="failed" if sparse_supporting else "passed",
                message=(
                    "未发现信息单元不足的普通内容页"
                    if not sparse_supporting
                    else "普通内容页信息单元不足：" + "、".join(sparse_supporting[:8])
                ),
                value=len(sparse_supporting),
            )
        )

        generic_titles = {
            "背景介绍",
            "项目背景",
            "概述",
            "相关内容",
            "数据分析",
            "方案介绍",
            "总结",
        }
        generic = [
            slide.id for slide in spec.slides if slide.title.strip() in generic_titles
        ]
        semantic.append(
            ValidationCheck(
                name="generic_titles",
                status="warning" if generic else "passed",
                message=(
                    "未发现栏目式空标题"
                    if not generic
                    else "栏目式空标题页面：" + "、".join(generic[:8])
                ),
                value=len(generic),
            )
        )
    external = _external_ooxml_relationships(path)
    security.append(ValidationCheck(name="external_relationships", status="failed" if external else "passed", message=f"外部关系 {len(external)}", value=len(external)))
    return structural, semantic, visual, security


def _overlap_ratio(first: tuple[int, int, int, int], second: tuple[int, int, int, int]) -> float:
    left = max(first[0], second[0])
    top = max(first[1], second[1])
    right = min(first[0] + first[2], second[0] + second[2])
    bottom = min(first[1] + first[3], second[1] + second[3])
    if right <= left or bottom <= top:
        return 0.0
    intersection = (right - left) * (bottom - top)
    smaller = min(first[2] * first[3], second[2] * second[3])
    return intersection / smaller if smaller else 0.0


def _docx_checks(
    path: Path,
    *,
    strict: bool,
) -> tuple[list[ValidationCheck], list[ValidationCheck], list[ValidationCheck]]:
    document = Document(str(path))
    structural = [ValidationCheck(name="reopen", status="passed", message="DOCX 可重新打开")]
    headings = [paragraph for paragraph in document.paragraphs if paragraph.style and paragraph.style.name.startswith("Heading")]
    heading_levels: list[int] = []
    for paragraph in headings:
        style = paragraph.style
        match = re.match(r"Heading\s+([1-9])", style.name if style is not None else "")
        if match:
            heading_levels.append(int(match.group(1)))
    heading_jumps = sum(
        current > previous + 1
        for previous, current in pairwise(heading_levels)
    )
    letter_sections = 0
    safe_margins = 0
    for section in document.sections:
        width = float(section.page_width.inches) if section.page_width is not None else 0
        height = float(section.page_height.inches) if section.page_height is not None else 0
        letter_sections += int(abs(width - 8.5) <= 0.05 and abs(height - 11) <= 0.05)
        margins = (section.top_margin, section.bottom_margin, section.left_margin, section.right_margin)
        safe_margins += int(all(value is not None and float(value.inches) >= 0.7 for value in margins))
    fixed_tables = 0
    fixed_row_heights = 0
    excessive_columns = 0
    for table in document.tables:
        layout = table._tbl.tblPr.find(qn("w:tblLayout"))
        fixed_tables += int(layout is not None and layout.get(qn("w:type")) == "fixed")
        fixed_row_heights += len(table._tbl.xpath(".//w:trHeight"))
        excessive_columns += int(bool(table.rows) and len(table.rows[0].cells) > 10)
    structural.extend(
        [
            ValidationCheck(
                name="page_geometry",
                status=("failed" if strict else "warning") if letter_sections != len(document.sections) or safe_margins != len(document.sections) else "passed",
                message=f"Letter 页型 {letter_sections}/{len(document.sections)}，安全页边距 {safe_margins}/{len(document.sections)}",
            ),
            ValidationCheck(
                name="table_layout",
                status=("failed" if strict else "warning") if fixed_tables != len(document.tables) or fixed_row_heights or excessive_columns else "passed",
                message=f"固定网格表格 {fixed_tables}/{len(document.tables)}，固定行高 {fixed_row_heights}，超宽表格 {excessive_columns}",
            ),
        ]
    )
    semantic = [
        ValidationCheck(name="heading_hierarchy", status=("failed" if strict else "warning") if not headings or heading_jumps else "passed", message=f"标题段落 {len(headings)}，跳级 {heading_jumps}", value=len(headings)),
        ValidationCheck(name="tables", status="passed", message=f"表格 {len(document.tables)}", value=len(document.tables)),
    ]
    external = _external_ooxml_relationships(path)
    with zipfile.ZipFile(path) as archive:
        embedded = [name for name in archive.namelist() if name.startswith("word/embeddings/")]
    security = [
        ValidationCheck(name="external_relationships", status="failed" if external else "passed", message=f"外部关系 {len(external)}", value=len(external)),
        ValidationCheck(name="embedded_objects", status="failed" if embedded else "passed", message=f"嵌入对象 {len(embedded)}", value=len(embedded)),
    ]
    return structural, semantic, security


def _xlsx_checks(
    path: Path,
    *,
    strict: bool,
) -> tuple[list[ValidationCheck], list[ValidationCheck], list[ValidationCheck]]:
    workbook = load_workbook(path, read_only=False, data_only=False, keep_links=True)
    try:
        formulas = 0
        errors = 0
        dde = 0
        network_formulas = 0
        merged = 0
        charts = 0
        blank_sheets = 0
        clipped_cells = 0
        for sheet in workbook.worksheets:
            merged += len(sheet.merged_cells.ranges)
            charts += len(sheet._charts)
            populated = 0
            for row in sheet.iter_rows():
                for cell in row:
                    value = cell.value
                    populated += int(value is not None)
                    if isinstance(value, str) and value.startswith("="):
                        formulas += 1
                        upper = value.upper()
                        errors += sum(marker in upper for marker in _FORMULA_ERRORS)
                        dde += int("|" in value or "[" in value or "]" in value)
                        network_formulas += int(_NETWORK_FORMULA.search(upper) is not None)
                    if value is not None and isinstance(value, str):
                        width = float(sheet.column_dimensions[cell.column_letter].width or 13)
                        if _display_width(value) > max(1, width - 1) and not cell.alignment.wrap_text:
                            clipped_cells += 1
            blank_sheets += int(populated == 0)
        structural = [ValidationCheck(name="reopen", status="passed", message="XLSX 可重新打开")]
        semantic = [
            ValidationCheck(name="sheets", status=("failed" if strict else "warning") if blank_sheets else "passed", message=f"工作表 {len(workbook.sheetnames)}，空白 {blank_sheets}", value=len(workbook.sheetnames)),
            ValidationCheck(name="formulas", status="failed" if errors else "passed", message=f"公式 {formulas}，显式错误 {errors}", value=formulas),
            ValidationCheck(name="merged_cells", status="passed", message=f"合并区域 {merged}", value=merged),
            ValidationCheck(name="charts", status="passed", message=f"图表 {charts}", value=charts),
            ValidationCheck(name="potential_clipping", status=("failed" if strict else "warning") if clipped_cells else "passed", message=f"疑似被列宽截断的文本单元格 {clipped_cells}", value=clipped_cells),
        ]
        external = _external_ooxml_relationships(path)
        security = [
            ValidationCheck(name="external_relationships", status="failed" if external else "passed", message=f"外部关系 {len(external)}", value=len(external)),
            ValidationCheck(name="dde_or_external_formula", status="failed" if dde else "passed", message=f"疑似 DDE/外部公式 {dde}", value=dde),
            ValidationCheck(
                name="network_formula",
                status="failed" if network_formulas else "passed",
                message=f"可联网或读取外部数据的公式 {network_formulas}",
                value=network_formulas,
            ),
        ]
        return structural, semantic, security
    finally:
        workbook.close()


_PDF_ACTIVE_OBJECT = re.compile(
    r"/(?:OpenAction|AA|JavaScript|JS|Launch|URI|EmbeddedFiles|Filespec|RichMedia|GoToR)\b"
)


def _pdf_active_content(document: object) -> list[str]:
    findings: list[str] = []
    xref_length = int(document.xref_length())  # type: ignore[attr-defined]
    for xref in range(1, xref_length):
        try:
            raw = str(document.xref_object(xref, compressed=False))  # type: ignore[attr-defined]
        except (RuntimeError, ValueError):
            findings.append(f"xref:{xref}:unreadable")
            continue
        findings.extend(
            f"xref:{xref}:{match.group(0)}" for match in _PDF_ACTIVE_OBJECT.finditer(raw)
        )
    try:
        findings.extend(
            f"embedded:{name}" for name in document.embfile_names()  # type: ignore[attr-defined]
        )
    except (RuntimeError, ValueError):
        findings.append("embedded-files:unreadable")
    return findings


def _pdf_checks(
    path: Path,
) -> tuple[
    list[ValidationCheck],
    list[ValidationCheck],
    list[ValidationCheck],
    list[ValidationCheck],
]:
    document = _pymupdf.open(path)
    try:
        encrypted = bool(document.needs_pass)
        blank = 0
        raster_failures = 0
        out_of_bounds = 0
        for page in document:
            blank += int(
                not page.get_text("text").strip()
                and not page.get_images()
                and not page.get_drawings()
            )
            try:
                pixmap = page.get_pixmap(matrix=_pymupdf.Matrix(1.25, 1.25), alpha=False)
                raster_failures += int(pixmap.width <= 0 or pixmap.height <= 0)
            except (RuntimeError, ValueError):
                raster_failures += 1
            rectangle = page.rect
            for block in page.get_text("blocks"):
                x0, y0, x1, y1 = block[:4]
                out_of_bounds += int(
                    x0 < rectangle.x0 - 1
                    or y0 < rectangle.y0 - 1
                    or x1 > rectangle.x1 + 1
                    or y1 > rectangle.y1 + 1
                )
        structural = [ValidationCheck(name="reopen", status="failed" if encrypted else "passed", message="PDF 已加密" if encrypted else "PDF 可重新打开")]
        semantic = [ValidationCheck(name="pages", status="failed" if document.page_count == 0 else "passed", message=f"页面 {document.page_count}", value=document.page_count)]
        visual = [
            ValidationCheck(name="blank_pages", status="failed" if blank else "passed", message=f"空白页 {blank}", value=blank),
            ValidationCheck(name="page_raster", status="failed" if raster_failures else "passed", message=f"逐页栅格化失败 {raster_failures}", value=raster_failures),
            ValidationCheck(name="text_bounds", status="failed" if out_of_bounds else "passed", message=f"越界文本块 {out_of_bounds}", value=out_of_bounds),
        ]
        active_content = _pdf_active_content(document)
        security = [
            ValidationCheck(
                name="active_content",
                status="failed" if active_content else "passed",
                message=f"PDF 主动内容或外部动作 {len(active_content)}",
                value=len(active_content),
            )
        ]
        return structural, semantic, visual, security
    finally:
        document.close()


def _html_checks(path: Path) -> tuple[list[ValidationCheck], list[ValidationCheck], list[ValidationCheck]]:
    raw = path.read_text(encoding="utf-8")
    parser = _OfflineHtmlSecurityParser()
    parser.feed(raw)
    parser.close()
    active = parser.findings
    structural = [ValidationCheck(name="utf8", status="passed", message="HTML 是 UTF-8 文本")]
    semantic = [ValidationCheck(name="title", status="warning" if "<title" not in raw.casefold() else "passed", message="已检查 title")]
    security = [ValidationCheck(name="active_or_external_content", status="failed" if active else "passed", message=f"主动或外部内容 {len(active)}", value=len(active))]
    return structural, semantic, security


def _display_width(value: str) -> int:
    return max(
        (
            sum(2 if east_asian_width(character) in {"W", "F", "A"} else 1 for character in line)
            for line in value.splitlines() or [""]
        ),
        default=0,
    )


def _evidence_checks(
    spec: ArtifactSpec | None,
    *,
    bound_claim_ids: frozenset[str],
) -> list[ValidationCheck]:
    if spec is None or spec.evidence_policy == "none":
        return []
    total = len(spec.claims)
    bound = sum(claim.claim_id in bound_claim_ids and bool(claim.evidence_ids) for claim in spec.claims)
    coverage = 100.0 if total == 0 else bound * 100.0 / total
    status: ValidationStatus = "passed"
    if coverage < 100 and spec.evidence_policy == "required":
        status = "failed"
    elif coverage < 100:
        status = "warning"
    return [ValidationCheck(name="claim_coverage", status=status, message=f"Claim 证据覆盖 {coverage:.1f}%（{bound}/{total}）", value=round(coverage, 1))]


#: 各维度在综合分里的权重。**security 刻意不在这张表里**——它是一票否决维度：
#: 任一安全检查 failed，综合分直接归零。安全问题不允许被其它维度的高分稀释成
#: "92 分，小瑕疵"，那正是它最容易被当成 warning 放过去的方式。
_DIMENSION_WEIGHTS: dict[str, int] = {
    "structural": 30,
    "semantic": 25,
    "visual": 25,
    "evidence": 20,
}
_FAILED_PENALTY = 50
_WARNING_PENALTY = 10


def _dimension_score(dimension: ValidationDimension) -> int:
    """维度内部先各自算分，再按权重合成。

    没有检查项或整维 ``not_run`` 记满分：**"没测量"不该和"测过且有问题"同罚**。
    缺 LibreOffice 的机器上 visual 恒为 not_run，若按 0 分计，扣的是环境不是产物，
    而且会让同一份文件在两台机器上给出不同质量分，评测直接失去可比性。
    "到底验没验过"由 ``ValidationDimension.status`` 与 ``deliverable`` 回答，
    不由分数兼职承担。
    """

    score = 100
    for check in dimension.checks:
        if check.status == "failed":
            score -= _FAILED_PENALTY
        elif check.status == "warning":
            score -= _WARNING_PENALTY
    return max(0, score)


def _quality(
    *,
    structural: ValidationDimension,
    semantic: ValidationDimension,
    visual: ValidationDimension,
    evidence: ValidationDimension,
    security: ValidationDimension,
) -> ArtifactQuality:
    scored = {
        "structural": structural,
        "semantic": semantic,
        "visual": visual,
        "evidence": evidence,
    }
    dimension_scores = {name: _dimension_score(item) for name, item in scored.items()}
    dimension_scores["security"] = _dimension_score(security)
    weighted = sum(
        dimension_scores[name] * weight for name, weight in _DIMENSION_WEIGHTS.items()
    ) / sum(_DIMENSION_WEIGHTS.values())

    warnings: list[str] = []
    for dimension in (structural, semantic, visual, evidence, security):
        warnings.extend(
            check.message for check in dimension.checks if check.status in {"failed", "warning"}
        )
    blocking = [check.message for check in security.checks if check.status == "failed"]
    score = 0 if blocking else round(weighted)
    return ArtifactQuality(
        score=score,
        warnings=warnings,
        dimension_scores=dimension_scores,
        blocking=blocking,
    )


def validate_artifact(
    path: Path,
    *,
    spec: ArtifactSpec | None = None,
    bound_claim_ids: frozenset[str] = frozenset(),
    render_visual: bool = False,
) -> ArtifactValidationReport:
    suffix = path.suffix.casefold()
    structural: list[ValidationCheck]
    semantic: list[ValidationCheck]
    visual: list[ValidationCheck] = []
    security: list[ValidationCheck] = []
    artifact_type: Literal["docx", "xlsx", "pptx", "pdf", "html"]
    if suffix in {".pptx", ".docx", ".xlsx"}:
        _validate_ooxml_container_budget(path)
    if suffix == ".pptx":
        artifact_type = "pptx"
        ppt_spec = spec if isinstance(spec, PresentationSpec) else None
        structural, semantic, visual, security = _pptx_checks(path, ppt_spec)
    elif suffix == ".docx":
        artifact_type = "docx"
        structural, semantic, security = _docx_checks(
            path,
            strict=isinstance(spec, DocumentSpec),
        )
    elif suffix == ".xlsx":
        artifact_type = "xlsx"
        structural, semantic, security = _xlsx_checks(
            path,
            strict=isinstance(spec, WorkbookSpec),
        )
    elif suffix == ".pdf":
        artifact_type = "pdf"
        structural, semantic, visual, security = _pdf_checks(path)
    elif suffix in {".html", ".htm"}:
        artifact_type = "html"
        structural, semantic, security = _html_checks(path)
    else:
        raise ValueError("该格式没有 Artifact Validator")
    if render_visual and suffix in {".pptx", ".docx", ".xlsx"}:
        visual.append(_render_check(path))
    elif suffix in {".pptx", ".docx", ".xlsx"}:
        visual.append(ValidationCheck(name="render", status="not_run", message="未请求真实版面渲染"))
    structural_dimension = _dimension(structural)
    semantic_dimension = _dimension(semantic)
    visual_dimension = _dimension(visual)
    evidence_dimension = _dimension(_evidence_checks(spec, bound_claim_ids=bound_claim_ids))
    security_dimension = _dimension(security)
    return ArtifactValidationReport(
        artifact_type=artifact_type,
        structural=structural_dimension,
        semantic=semantic_dimension,
        visual=visual_dimension,
        evidence=evidence_dimension,
        security=security_dimension,
        quality=_quality(
            structural=structural_dimension,
            semantic=semantic_dimension,
            visual=visual_dimension,
            evidence=evidence_dimension,
            security=security_dimension,
        ),
    )


__all__ = [
    "ArtifactValidationReport",
    "validate_artifact",
    "validate_artifact_in_subprocess",
]
