"""WorkPilot 办公内容评测：对最终 DOCX/XLSX/PPTX/PDF 做可审计评分。

这条评测轨与 ``artifact_suite`` 分工明确：后者验证固定 Renderer 能否稳定产出文件，
本模块验证一个办公任务的最终交付物是否完整、正确、忠于资料且可用。自动轨只使用
确定性、实例级检查；连贯性、决策价值和视觉层次留给带锚点的人工/VLM 复核轨。
"""

from __future__ import annotations

import argparse
import ast
import hashlib
import json
import math
import re
import sys
import unicodedata
import zipfile
from collections import Counter
from dataclasses import dataclass, field
from datetime import UTC, datetime
from pathlib import Path, PurePosixPath
from typing import Annotated, Any, Literal, cast

import pymupdf
from docx import Document
from docx.document import Document as DocxDocument
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table
from docx.text.paragraph import Paragraph
from openpyxl import load_workbook  # type: ignore[import-untyped]
from openpyxl.utils.cell import (  # type: ignore[import-untyped]
    get_column_letter,
    range_boundaries,
)
from pptx import Presentation
from pydantic import BaseModel, ConfigDict, Field, ValidationError, model_validator

from app.cowork.artifact_validation import (  # type: ignore[import-untyped]
    ArtifactValidationReport,
    ValidationStatus,
    validate_artifact,
)

SCHEMA_VERSION = "workpilot-office-content-suite.v1"
LEGACY_REVIEW_SCHEMA_VERSION = "workpilot-office-content-reviews.v1"
REVIEW_SCHEMA_VERSION = "workpilot-office-content-reviews.v2"
REPORT_SCHEMA_VERSION = "workpilot-office-content-report.v2"
DEFAULT_SUITE = Path(__file__).parent / "suites" / "office-content-dev-v1.json"

ArtifactType = Literal["docx", "xlsx", "pptx", "pdf"]
Split = Literal["dev", "test"]
Dimension = Literal[
    "fundamentals",
    "completeness",
    "correctness",
    "fidelity",
    "usability",
]
ReviewDimension = Literal["coherence", "decision_utility", "visual_quality"]
ReviewSource = Literal["human", "model"]
ReviewRenderMode = Literal[
    "native_pptx",
    "native_pdf",
    "office_preview",
    "structural_fallback",
]
ValidatorDimensionName = Literal["structural", "semantic", "visual", "security"]
_DIMENSIONS: tuple[Dimension, ...] = (
    "fundamentals",
    "completeness",
    "correctness",
    "fidelity",
    "usability",
)
_REVIEW_DIMENSIONS: tuple[ReviewDimension, ...] = (
    "coherence",
    "decision_utility",
    "visual_quality",
)
_SUFFIXES: dict[ArtifactType, str] = {
    "docx": ".docx",
    "xlsx": ".xlsx",
    "pptx": ".pptx",
    "pdf": ".pdf",
}


class OfficeContentSuiteError(ValueError):
    """套件、提交或复核结果不满足可复现评分契约。"""


class _StrictModel(BaseModel):
    model_config = ConfigDict(extra="forbid")


class SourceFixture(_StrictModel):
    path: str = Field(min_length=1, max_length=180)
    content: str = Field(max_length=1_000_000)

    @model_validator(mode="after")
    def _safe_relative_path(self) -> SourceFixture:
        path = PurePosixPath(self.path)
        if path.is_absolute() or ".." in path.parts or not path.name:
            raise ValueError("fixture path 必须是安全的相对路径")
        return self


class ScopeSpec(_StrictModel):
    kind: Literal[
        "artifact",
        "docx_section",
        "docx_table",
        "pptx_slide",
        "xlsx_sheet",
        "xlsx_range",
        "pdf_page",
    ] = "artifact"
    index: int | None = Field(default=None, ge=1)
    selector: str | None = Field(default=None, min_length=1, max_length=200)
    cell_range: str | None = Field(
        default=None,
        pattern=r"^\$?[A-Z]{1,3}\$?[1-9][0-9]*:\$?[A-Z]{1,3}\$?[1-9][0-9]*$",
    )

    @model_validator(mode="after")
    def _valid_selector(self) -> ScopeSpec:
        if self.kind == "artifact":
            if (
                self.index is not None
                or self.selector is not None
                or self.cell_range is not None
            ):
                raise ValueError("artifact scope 不能带定位字段")
        elif self.kind in {"docx_table", "pdf_page"}:
            if (
                self.index is None
                or self.selector is not None
                or self.cell_range is not None
            ):
                raise ValueError(f"{self.kind} scope 只接受 index")
        elif self.kind == "docx_section":
            if (
                self.selector is None
                or self.index is not None
                or self.cell_range is not None
            ):
                raise ValueError("docx_section scope 只接受 selector")
        elif self.kind == "pptx_slide":
            if (self.index is None) == (
                self.selector is None
            ) or self.cell_range is not None:
                raise ValueError("pptx_slide scope 必须且只能使用 index 或 selector")
        elif self.kind == "xlsx_sheet":
            if (
                self.selector is None
                or self.index is not None
                or self.cell_range is not None
            ):
                raise ValueError("xlsx_sheet scope 只接受 selector")
        elif self.kind == "xlsx_range" and (
            self.selector is None or self.cell_range is None or self.index is not None
        ):
            raise ValueError("xlsx_range scope 必须包含 selector 与 cell_range")
        return self


class TextContainsCheck(_StrictModel):
    type: Literal["text_contains"]
    values: list[str] = Field(min_length=1, max_length=20)
    match: Literal["all", "any"] = "all"
    scope: ScopeSpec = Field(default_factory=ScopeSpec)


class TextNotContainsCheck(_StrictModel):
    type: Literal["text_not_contains"]
    values: list[str] = Field(min_length=1, max_length=20)
    scope: ScopeSpec = Field(default_factory=ScopeSpec)


class TextOrderedCheck(_StrictModel):
    type: Literal["text_ordered"]
    values: list[str] = Field(min_length=2, max_length=20)
    scope: ScopeSpec = Field(default_factory=ScopeSpec)


class LabeledValueRequirement(_StrictModel):
    """Bind one or more values to a nearby semantic label.

    Plain keyword checks cannot distinguish ``actual=860, target=900`` from the
    reversed relation.  A requirement succeeds only when every value appears in
    the bounded context of at least one allowed label.
    """

    labels: list[str] = Field(min_length=1, max_length=10)
    values: list[str] = Field(min_length=1, max_length=10)
    direction: Literal["after", "within"] = "after"
    max_distance: int = Field(default=80, ge=1, le=500)

    @model_validator(mode="after")
    def _non_empty_tokens(self) -> LabeledValueRequirement:
        if any(not value.strip() for value in [*self.labels, *self.values]):
            raise ValueError("labels/values 不能包含空字符串")
        return self


class TextLabeledValuesCheck(_StrictModel):
    type: Literal["text_labeled_values"]
    requirements: list[LabeledValueRequirement] = Field(min_length=1, max_length=20)
    scope: ScopeSpec = Field(default_factory=ScopeSpec)


class PptxChartDataCheck(_StrictModel):
    type: Literal["pptx_chart_data"]
    required_values: list[float] = Field(min_length=1, max_length=100)
    required_labels: list[str] = Field(default_factory=list, max_length=50)
    tolerance: float = Field(default=1e-6, ge=0.0)
    accept_percentage_fractions: bool = True


def _default_negations() -> list[str]:
    return ["不", "未", "没有", "并非", "不是", "尚未", "尚无", "无法", "不能", "不得"]


class TextClaimContainsCheck(_StrictModel):
    type: Literal["text_claim_contains"]
    values: list[str] = Field(min_length=1, max_length=20)
    match: Literal["all", "any"] = "any"
    scope: ScopeSpec = Field(default_factory=ScopeSpec)
    negations: list[str] = Field(
        default_factory=_default_negations, min_length=1, max_length=20
    )
    context_window: int = Field(default=12, ge=1, le=50)


class TextClaimNotContainsCheck(_StrictModel):
    type: Literal["text_claim_not_contains"]
    values: list[str] = Field(min_length=1, max_length=20)
    scope: ScopeSpec = Field(default_factory=ScopeSpec)
    negations: list[str] = Field(
        default_factory=_default_negations, min_length=1, max_length=20
    )
    context_window: int = Field(default=12, ge=1, le=50)


class StructureCountCheck(_StrictModel):
    type: Literal["structure_count"]
    metric: Literal[
        "docx_heading_count",
        "docx_table_count",
        "docx_paragraph_count",
        "pptx_slide_count",
        "pptx_table_count",
        "pptx_chart_count",
        "xlsx_sheet_count",
        "xlsx_formula_count",
        "xlsx_table_count",
        "xlsx_chart_count",
        "pdf_page_count",
    ]
    operator: Literal["eq", "gte", "lte"]
    value: int = Field(ge=0, le=10_000)


class StructureValuesCheck(_StrictModel):
    type: Literal["structure_values"]
    metric: Literal["docx_headings", "pptx_slide_titles", "xlsx_sheet_names"]
    values: list[str] = Field(min_length=1, max_length=50)
    match: Literal["all", "any", "exact_order"] = "all"


class XlsxCellValueCheck(_StrictModel):
    type: Literal["xlsx_cell_value"]
    address: str = Field(pattern=r"^.+!\$?[A-Z]{1,3}\$?[1-9][0-9]*$")
    expected: str | int | float | bool
    tolerance: float = Field(default=0.0, ge=0.0)


class XlsxCellsValueCheck(_StrictModel):
    type: Literal["xlsx_cells_value"]
    cells: dict[str, str | int | float | bool] = Field(min_length=1, max_length=100)
    tolerance: float = Field(default=0.0, ge=0.0)

    @model_validator(mode="after")
    def _valid_addresses(self) -> XlsxCellsValueCheck:
        pattern = re.compile(r"^.+!\$?[A-Z]{1,3}\$?[1-9][0-9]*$")
        invalid = [
            address for address in self.cells if pattern.fullmatch(address) is None
        ]
        if invalid:
            raise ValueError(f"非法 XLSX 单元格地址：{invalid}")
        return self


class XlsxFormulaValueCheck(_StrictModel):
    type: Literal["xlsx_formula_value"]
    address: str = Field(pattern=r"^.+!\$?[A-Z]{1,3}\$?[1-9][0-9]*$")
    expected: float
    tolerance: float = Field(default=1e-6, ge=0.0)
    must_contain: list[str] = Field(default_factory=list, max_length=10)


class XlsxNumberFormatCheck(_StrictModel):
    type: Literal["xlsx_number_format"]
    address: str = Field(pattern=r"^.+!\$?[A-Z]{1,3}\$?[1-9][0-9]*$")
    contains_any: list[str] = Field(min_length=1, max_length=10)


class ValidatorStatusCheck(_StrictModel):
    type: Literal["validator_status"]
    dimension: Literal["structural", "semantic", "visual", "security"]
    check_name: str | None = Field(default=None, min_length=1, max_length=100)
    allowed: list[ValidationStatus] = Field(min_length=1, max_length=4)


CheckSpec = Annotated[
    TextContainsCheck
    | TextNotContainsCheck
    | TextOrderedCheck
    | TextLabeledValuesCheck
    | TextClaimContainsCheck
    | TextClaimNotContainsCheck
    | StructureCountCheck
    | StructureValuesCheck
    | XlsxCellValueCheck
    | XlsxCellsValueCheck
    | XlsxFormulaValueCheck
    | XlsxNumberFormatCheck
    | PptxChartDataCheck
    | ValidatorStatusCheck,
    Field(discriminator="type"),
]


def _check_artifact_type(check: CheckSpec) -> ArtifactType | None:
    if isinstance(
        check,
        (
            TextContainsCheck,
            TextNotContainsCheck,
            TextOrderedCheck,
            TextLabeledValuesCheck,
            TextClaimContainsCheck,
            TextClaimNotContainsCheck,
        ),
    ):
        scope_prefix = check.scope.kind.split("_", maxsplit=1)[0]
        if scope_prefix in _SUFFIXES:
            return scope_prefix
        return None
    if isinstance(check, StructureCountCheck):
        return cast(ArtifactType, check.metric.split("_", maxsplit=1)[0])
    if isinstance(check, StructureValuesCheck):
        return cast(ArtifactType, check.metric.split("_", maxsplit=1)[0])
    if isinstance(
        check,
        (
            XlsxCellValueCheck,
            XlsxCellsValueCheck,
            XlsxFormulaValueCheck,
            XlsxNumberFormatCheck,
        ),
    ):
        return "xlsx"
    if isinstance(check, PptxChartDataCheck):
        return "pptx"
    return None


class RubricItem(_StrictModel):
    id: str = Field(pattern=r"^[a-z0-9][a-z0-9._-]{0,119}$")
    dimension: Dimension
    description: str = Field(min_length=1, max_length=500)
    source_refs: list[str] = Field(min_length=1, max_length=20)
    weight: float = Field(default=1.0, gt=0, le=20)
    critical: bool = False
    check: CheckSpec


class PenaltyItem(_StrictModel):
    id: str = Field(pattern=r"^[a-z0-9][a-z0-9._-]{0,119}$")
    description: str = Field(min_length=1, max_length=500)
    source_refs: list[str] = Field(min_length=1, max_length=20)
    points: float = Field(gt=0, le=100)
    blocking: bool = False
    trigger: CheckSpec


class ReviewCriterion(_StrictModel):
    id: str = Field(pattern=r"^[a-z0-9][a-z0-9._-]{0,119}$")
    dimension: ReviewDimension
    description: str = Field(min_length=1, max_length=800)
    anchors: list[str] = Field(min_length=3, max_length=6)
    minimum_score: int = Field(default=1, ge=0)

    @model_validator(mode="after")
    def _minimum_within_scale(self) -> ReviewCriterion:
        if self.minimum_score > self.max_score:
            raise ValueError("minimum_score 不能超过复核量表最高分")
        return self

    @property
    def max_score(self) -> int:
        return len(self.anchors) - 1


def _default_gate_dimensions() -> list[ValidatorDimensionName]:
    return ["structural", "security"]


class GateSpec(_StrictModel):
    render_visual: bool = True
    fail_on_dimensions: list[ValidatorDimensionName] = Field(
        default_factory=_default_gate_dimensions
    )
    require_measured_dimensions: list[ValidatorDimensionName] = Field(
        default_factory=list
    )
    min_validator_quality: int = Field(default=0, ge=0, le=100)
    max_file_bytes: int = Field(default=64 * 1024 * 1024, ge=1)
    max_uncompressed_bytes: int = Field(default=256 * 1024 * 1024, ge=1)

    @model_validator(mode="after")
    def _no_duplicate_dimensions(self) -> GateSpec:
        if len(self.fail_on_dimensions) != len(set(self.fail_on_dimensions)):
            raise ValueError("fail_on_dimensions 不能重复")
        if len(self.require_measured_dimensions) != len(
            set(self.require_measured_dimensions)
        ):
            raise ValueError("require_measured_dimensions 不能重复")
        return self


class OfficeContentItem(_StrictModel):
    id: str = Field(pattern=r"^[a-z0-9][a-z0-9._-]{0,119}$")
    split: Split
    artifact_type: ArtifactType
    task_type: Literal[
        "create",
        "transform",
        "analyze",
        "conflict_handling",
    ]
    category: str = Field(pattern=r"^[a-z0-9][a-z0-9._-]{0,79}$")
    difficulty: Literal["easy", "medium", "hard"]
    prompt: str = Field(min_length=20, max_length=10_000)
    output_file: str = Field(pattern=r"^[A-Za-z0-9][A-Za-z0-9._-]{0,119}$")
    fixtures: list[SourceFixture] = Field(min_length=1, max_length=30)
    gate: GateSpec = Field(default_factory=GateSpec)
    rubric: list[RubricItem] = Field(min_length=5, max_length=100)
    penalties: list[PenaltyItem] = Field(default_factory=list, max_length=30)
    review_criteria: list[ReviewCriterion] = Field(min_length=1, max_length=10)
    pass_threshold: float = Field(default=75.0, ge=0, le=100)

    @model_validator(mode="after")
    def _consistent_item(self) -> OfficeContentItem:
        if Path(self.output_file).suffix.casefold() != _SUFFIXES[self.artifact_type]:
            raise ValueError("output_file 后缀必须与 artifact_type 一致")
        fixture_paths = [fixture.path for fixture in self.fixtures]
        if len(fixture_paths) != len(set(fixture_paths)):
            raise ValueError("同一任务的 fixture path 必须唯一")
        rubric_ids = [item.id for item in self.rubric]
        penalty_ids = [item.id for item in self.penalties]
        review_ids = [item.id for item in self.review_criteria]
        all_ids = [*rubric_ids, *penalty_ids, *review_ids]
        if len(all_ids) != len(set(all_ids)):
            raise ValueError("rubric、penalty 与 review id 必须在任务内唯一")
        dimensions = {item.dimension for item in self.rubric}
        missing = set(_DIMENSIONS) - dimensions
        if missing:
            raise ValueError(f"每道题必须覆盖五个自动维度，缺少 {sorted(missing)}")
        review_dimensions = {item.dimension for item in self.review_criteria}
        missing_review = set(_REVIEW_DIMENSIONS) - review_dimensions
        if missing_review:
            raise ValueError(
                f"每道题必须覆盖三个人工复核维度，缺少 {sorted(missing_review)}"
            )
        if self.artifact_type == "pptx":
            required_gate_dimensions = {"structural", "visual", "security"}
            if (
                not self.gate.render_visual
                or not required_gate_dimensions.issubset(self.gate.fail_on_dimensions)
                or not required_gate_dimensions.issubset(
                    self.gate.require_measured_dimensions
                )
            ):
                raise ValueError(
                    "PPTX 必须渲染并把 structural/visual/security 设为已测量硬门禁"
                )
        incompatible_checks = [
            rubric.id
            for rubric in self.rubric
            if _check_artifact_type(rubric.check) not in {None, self.artifact_type}
        ] + [
            penalty.id
            for penalty in self.penalties
            if _check_artifact_type(penalty.trigger) not in {None, self.artifact_type}
        ]
        if incompatible_checks:
            raise ValueError(f"检查类型与 artifact_type 不兼容：{incompatible_checks}")
        fixture_names = {fixture.path for fixture in self.fixtures}
        source_refs = [
            source_ref for rubric in self.rubric for source_ref in rubric.source_refs
        ] + [
            source_ref
            for penalty in self.penalties
            for source_ref in penalty.source_refs
        ]
        unknown_refs = sorted(
            {
                source_ref
                for source_ref in source_refs
                if source_ref != "prompt"
                and source_ref.split("#", maxsplit=1)[0] not in fixture_names
            }
        )
        if unknown_refs:
            raise ValueError(f"source_refs 引用了未知 fixture：{unknown_refs}")
        return self


class OfficeContentSuite(_StrictModel):
    schema_version: Literal["workpilot-office-content-suite.v1"]
    name: str = Field(pattern=r"^[a-z0-9][a-z0-9._-]{0,119}$")
    version: str = Field(pattern=r"^[0-9]+\.[0-9]+(?:\.[0-9]+)?$")
    origin: Literal["synthetic", "human", "public"]
    data_classification: Literal["synthetic", "internal", "public"]
    review_status: Literal["pending_human_review", "approved"]
    reviewer: str | None = Field(default=None, min_length=1, max_length=200)
    reviewed_at: str | None = Field(default=None, min_length=1, max_length=100)
    methodology: list[str] = Field(min_length=1, max_length=30)
    dimension_weights: dict[Dimension, float]
    automatic_weight: float = Field(default=0.6, ge=0, le=1)
    review_weight: float = Field(default=0.4, ge=0, le=1)
    items: list[OfficeContentItem] = Field(min_length=1, max_length=500)

    @model_validator(mode="after")
    def _consistent_suite(self) -> OfficeContentSuite:
        ids = [item.id for item in self.items]
        if len(ids) != len(set(ids)):
            raise ValueError("item id 必须唯一")
        if set(self.dimension_weights) != set(_DIMENSIONS):
            raise ValueError("dimension_weights 必须且只能包含五个自动维度")
        if any(weight < 0 for weight in self.dimension_weights.values()):
            raise ValueError("dimension_weights 不能为负数")
        if not math.isclose(sum(self.dimension_weights.values()), 1.0, abs_tol=1e-9):
            raise ValueError("dimension_weights 之和必须为 1")
        if not math.isclose(
            self.automatic_weight + self.review_weight, 1.0, abs_tol=1e-9
        ):
            raise ValueError("automatic_weight + review_weight 必须为 1")
        if self.review_status == "approved":
            if self.reviewer is None or self.reviewed_at is None:
                raise ValueError("approved suite 必须记录 reviewer 与 reviewed_at")
            try:
                reviewed_at = datetime.fromisoformat(
                    self.reviewed_at.replace("Z", "+00:00")
                )
            except ValueError as error:
                raise ValueError("reviewed_at 必须是 ISO-8601 时间") from error
            if reviewed_at.tzinfo is None:
                raise ValueError("reviewed_at 必须包含时区")
        elif self.reviewer is not None or self.reviewed_at is not None:
            raise ValueError("pending suite 不能预填 reviewer 或 reviewed_at")
        return self


class ReviewAnnotation(_StrictModel):
    item_id: str
    criterion_id: str
    artifact_sha256: str = Field(pattern=r"^[0-9a-f]{64}$")
    score: int = Field(ge=0)
    evidence: str = Field(min_length=1, max_length=2_000)
    reviewer: str = Field(min_length=1, max_length=200)
    source: ReviewSource = "human"
    provider: str | None = Field(default=None, min_length=1, max_length=200)
    model: str | None = Field(default=None, min_length=1, max_length=300)
    prompt_fingerprint: str | None = Field(
        default=None, pattern=r"^[0-9a-f]{64}$"
    )
    authorization_note_fingerprint: str | None = Field(
        default=None, pattern=r"^[0-9a-f]{64}$"
    )
    calibration_status: Literal["uncalibrated"] | None = None
    render_mode: ReviewRenderMode | None = None

    @model_validator(mode="after")
    def _valid_provenance(self) -> ReviewAnnotation:
        model_fields = (
            self.provider,
            self.model,
            self.prompt_fingerprint,
            self.authorization_note_fingerprint,
            self.calibration_status,
            self.render_mode,
        )
        if self.source == "human":
            if any(value is not None for value in model_fields):
                raise ValueError("human 复核不能携带模型 provenance")
        elif any(value is None for value in model_fields):
            raise ValueError("model 复核必须记录模型、prompt、授权、校准与渲染 provenance")
        return self


class ReviewFile(_StrictModel):
    schema_version: Literal[
        "workpilot-office-content-reviews.v1",
        "workpilot-office-content-reviews.v2",
    ]
    reviews: list[ReviewAnnotation] = Field(max_length=5_000)

    @model_validator(mode="after")
    def _unique_reviews(self) -> ReviewFile:
        keys = [(item.item_id, item.criterion_id) for item in self.reviews]
        if len(keys) != len(set(keys)):
            raise ValueError("同一 item/criterion 只能有一条复核")
        if self.schema_version == LEGACY_REVIEW_SCHEMA_VERSION and any(
            item.source != "human" for item in self.reviews
        ):
            raise ValueError("v1 review schema 只接受人工复核；模型复核必须使用 v2")
        return self


@dataclass(frozen=True)
class TextScope:
    kind: str
    label: str
    text: str
    index: int | None = None
    selector: str | None = None


@dataclass
class ArtifactView:
    artifact_type: ArtifactType
    text: str
    scopes: list[TextScope]
    counts: dict[str, int]
    collections: dict[str, list[str]]
    cells: dict[str, Any] = field(default_factory=dict)
    number_formats: dict[str, str] = field(default_factory=dict)
    charts: list[dict[str, Any]] = field(default_factory=list)


@dataclass(frozen=True)
class CheckOutcome:
    passed: bool
    detail: str
    evidence: list[str]
    actual: Any = None


def _json_safe(value: Any) -> Any:
    if isinstance(value, float) and not math.isfinite(value):
        return str(value)
    if value is None or isinstance(value, (str, int, float, bool)):
        return value
    if isinstance(value, list):
        return [_json_safe(item) for item in value]
    if isinstance(value, tuple):
        return [_json_safe(item) for item in value]
    if isinstance(value, dict):
        return {str(key): _json_safe(item) for key, item in value.items()}
    return str(value)


def _unique_object(pairs: list[tuple[str, object]]) -> dict[str, object]:
    result: dict[str, object] = {}
    for key, value in pairs:
        if key in result:
            raise OfficeContentSuiteError(f"重复 JSON key: {key}")
        result[key] = value
    return result


def load_suite(path: Path = DEFAULT_SUITE) -> OfficeContentSuite:
    try:
        payload = json.loads(
            path.read_text(encoding="utf-8"), object_pairs_hook=_unique_object
        )
        return OfficeContentSuite.model_validate(payload)
    except (
        OSError,
        UnicodeDecodeError,
        json.JSONDecodeError,
        ValidationError,
    ) as error:
        raise OfficeContentSuiteError(f"Office content suite 无效：{error}") from error


def load_reviews(path: Path | None) -> ReviewFile | None:
    if path is None:
        return None
    try:
        payload = json.loads(
            path.read_text(encoding="utf-8"), object_pairs_hook=_unique_object
        )
        return ReviewFile.model_validate(payload)
    except (
        OSError,
        UnicodeDecodeError,
        json.JSONDecodeError,
        ValidationError,
    ) as error:
        raise OfficeContentSuiteError(
            f"Office content reviews 无效：{error}"
        ) from error


def suite_summary(suite: OfficeContentSuite) -> dict[str, object]:
    return {
        "name": suite.name,
        "version": suite.version,
        "origin": suite.origin,
        "data_classification": suite.data_classification,
        "review_status": suite.review_status,
        "items": len(suite.items),
        "splits": dict(sorted(Counter(item.split for item in suite.items).items())),
        "formats": dict(
            sorted(Counter(item.artifact_type for item in suite.items).items())
        ),
        "task_types": dict(
            sorted(Counter(item.task_type for item in suite.items).items())
        ),
        "difficulties": dict(
            sorted(Counter(item.difficulty for item in suite.items).items())
        ),
        "automatic_checks": sum(len(item.rubric) for item in suite.items),
        "penalties": sum(len(item.penalties) for item in suite.items),
        "review_criteria": sum(len(item.review_criteria) for item in suite.items),
    }


def _suite_sha256(suite: OfficeContentSuite) -> str:
    canonical = json.dumps(
        suite.model_dump(mode="json"),
        ensure_ascii=False,
        separators=(",", ":"),
        sort_keys=True,
    ).encode("utf-8")
    return hashlib.sha256(canonical).hexdigest()


def _scorer_fingerprint() -> str:
    """Fingerprint every implementation file that can change an Office score."""

    repo_root = Path(__file__).resolve().parents[1]
    relative_paths = (
        Path("eval/office_content_suite.py"),
        Path("backend/app/cowork/artifact_validation.py"),
        Path("backend/app/cowork/artifact_renderers/contracts.py"),
        Path("backend/app/cowork/office_preview.py"),
        Path("backend/app/cowork/skills/builtin/pptx/scripts/pptx2image.py"),
        Path("backend/uv.lock"),
    )
    digest = hashlib.sha256()
    for relative in relative_paths:
        digest.update(relative.as_posix().encode("utf-8"))
        digest.update(b"\0")
        digest.update((repo_root / relative).read_bytes())
        digest.update(b"\0")
    return digest.hexdigest()


def _iter_docx_blocks(document: DocxDocument) -> list[Paragraph | Table]:
    blocks: list[Paragraph | Table] = []
    for child in document.element.body.iterchildren():
        if isinstance(child, CT_P):
            blocks.append(Paragraph(child, document))
        elif isinstance(child, CT_Tbl):
            blocks.append(Table(child, document))
    return blocks


def _heading_level(paragraph: Paragraph) -> int | None:
    style_name = paragraph.style.name if paragraph.style is not None else ""
    match = re.match(r"Heading\s+([1-9])", style_name)
    return int(match.group(1)) if match else None


def _table_text(table: Table) -> str:
    return "\n".join("\t".join(cell.text for cell in row.cells) for row in table.rows)


def _extract_docx(path: Path) -> ArtifactView:
    document = Document(str(path))
    blocks = _iter_docx_blocks(document)
    all_values: list[str] = []
    table_scopes: list[TextScope] = []
    section_values: list[dict[str, Any]] = []
    active_sections: list[int] = []
    headings: list[str] = []
    table_index = 0
    paragraph_count = 0
    for block in blocks:
        if isinstance(block, Paragraph):
            paragraph_count += 1
            text = block.text.strip()
            level = _heading_level(block)
            if level is not None and text:
                headings.append(text)
                while (
                    active_sections
                    and section_values[active_sections[-1]]["level"] >= level
                ):
                    active_sections.pop()
                section_values.append(
                    {"heading": text, "level": level, "values": [text]}
                )
                active_sections.append(len(section_values) - 1)
            elif text:
                for index in active_sections:
                    section_values[index]["values"].append(text)
            if text:
                all_values.append(text)
        else:
            table_index += 1
            text = _table_text(block)
            table_scopes.append(
                TextScope(
                    kind="docx_table",
                    label=f"DOCX table {table_index}",
                    text=text,
                    index=table_index,
                )
            )
            if text:
                all_values.append(text)
                for index in active_sections:
                    section_values[index]["values"].append(text)
    section_scopes = [
        TextScope(
            kind="docx_section",
            label=f"DOCX section {item['heading']}",
            text="\n".join(item["values"]),
            selector=str(item["heading"]),
        )
        for item in section_values
    ]
    return ArtifactView(
        artifact_type="docx",
        text="\n".join(all_values),
        scopes=[*section_scopes, *table_scopes],
        counts={
            "docx_heading_count": len(headings),
            "docx_table_count": len(document.tables),
            "docx_paragraph_count": paragraph_count,
        },
        collections={"docx_headings": headings},
    )


def _shape_text(shape: Any) -> list[str]:
    values: list[str] = []
    if bool(getattr(shape, "has_text_frame", False)):
        text = str(getattr(shape, "text", "")).strip()
        if text:
            values.append(text)
    if bool(getattr(shape, "has_table", False)):
        values.extend(
            "\t".join(cell.text for cell in row.cells) for row in shape.table.rows
        )
    return values


def _extract_pptx(path: Path) -> ArtifactView:
    presentation = Presentation(str(path))
    scopes: list[TextScope] = []
    titles: list[str] = []
    all_values: list[str] = []
    table_count = 0
    chart_count = 0
    charts: list[dict[str, Any]] = []
    for index, slide in enumerate(presentation.slides, start=1):
        values: list[str] = []
        title_shape = slide.shapes.title
        title = str(title_shape.text).strip() if title_shape is not None else ""
        for shape in slide.shapes:
            values.extend(_shape_text(shape))
            table_count += int(bool(getattr(shape, "has_table", False)))
            if bool(getattr(shape, "has_chart", False)):
                chart_count += 1
                chart = shape.chart
                chart_values: list[float] = []
                chart_labels: list[str] = []
                series_names: list[str] = []
                for plot in chart.plots:
                    try:
                        chart_labels.extend(
                            str(category.label)
                            for category in plot.categories
                            if category.label is not None
                        )
                    except (AttributeError, TypeError, ValueError):
                        # Some chart types do not expose a category sequence.
                        pass
                    for series in plot.series:
                        name = str(series.name or "").strip()
                        if name:
                            series_names.append(name)
                        chart_values.extend(
                            float(value)
                            for value in series.values
                            if isinstance(value, (int, float)) and math.isfinite(value)
                        )
                charts.append(
                    {
                        "slide_index": index,
                        "title": title,
                        "labels": chart_labels,
                        "series_names": series_names,
                        "values": chart_values,
                    }
                )
        # Speaker notes are deliberately excluded.  Content rubrics evaluate
        # what the audience can see; otherwise a submission can hide every
        # required fact in notes while leaving the actual slide empty.
        if not title:
            title = next(
                (value.splitlines()[0] for value in values if value.strip()),
                f"Slide {index}",
            )
        titles.append(title)
        text = "\n".join(values)
        all_values.append(text)
        scopes.append(
            TextScope(
                kind="pptx_slide",
                label=f"PPTX slide {index}: {title}",
                text=text,
                index=index,
                selector=title,
            )
        )
    return ArtifactView(
        artifact_type="pptx",
        text="\n".join(all_values),
        scopes=scopes,
        counts={
            "pptx_slide_count": len(presentation.slides),
            "pptx_table_count": table_count,
            "pptx_chart_count": chart_count,
        },
        collections={"pptx_slide_titles": titles},
        charts=charts,
    )


def _cell_address(sheet_name: str, coordinate: str) -> str:
    return f"{sheet_name}!{coordinate.replace('$', '').upper()}"


def _extract_xlsx(path: Path) -> ArtifactView:
    workbook = load_workbook(path, data_only=False, read_only=False)
    try:
        scopes: list[TextScope] = []
        all_values: list[str] = []
        cells: dict[str, Any] = {}
        formats: dict[str, str] = {}
        formula_count = 0
        chart_count = 0
        table_count = 0
        for sheet in workbook.worksheets:
            rows: list[str] = []
            for row in sheet.iter_rows():
                row_values: list[str] = []
                for cell in row:
                    if cell.value is None:
                        row_values.append("")
                        continue
                    address = _cell_address(sheet.title, cell.coordinate)
                    cells[address] = cell.value
                    formats[address] = str(cell.number_format)
                    formula_count += int(
                        isinstance(cell.value, str) and cell.value.startswith("=")
                    )
                    row_values.append(str(cell.value))
                if any(row_values):
                    rows.append("\t".join(row_values).rstrip())
            text = "\n".join(rows)
            all_values.append(text)
            scopes.append(
                TextScope(
                    kind="xlsx_sheet",
                    label=f"XLSX sheet {sheet.title}",
                    text=text,
                    selector=sheet.title,
                )
            )
            chart_count += len(sheet._charts)
            table_count += len(sheet.tables)
        return ArtifactView(
            artifact_type="xlsx",
            text="\n".join(all_values),
            scopes=scopes,
            counts={
                "xlsx_sheet_count": len(workbook.sheetnames),
                "xlsx_formula_count": formula_count,
                "xlsx_chart_count": chart_count,
                "xlsx_table_count": table_count,
            },
            collections={"xlsx_sheet_names": list(workbook.sheetnames)},
            cells=cells,
            number_formats=formats,
        )
    finally:
        workbook.close()


def _extract_pdf(path: Path) -> ArtifactView:
    document = pymupdf.open(path)  # type: ignore[no-untyped-call]
    try:
        scopes: list[TextScope] = []
        values: list[str] = []
        for zero_based_index in range(document.page_count):
            index = zero_based_index + 1
            page: pymupdf.Page = document.load_page(  # type: ignore[no-untyped-call]
                zero_based_index
            )
            text = page.get_text("text")  # type: ignore[no-untyped-call]
            values.append(text)
            scopes.append(
                TextScope(
                    kind="pdf_page", label=f"PDF page {index}", text=text, index=index
                )
            )
        return ArtifactView(
            artifact_type="pdf",
            text="\n".join(values),
            scopes=scopes,
            counts={"pdf_page_count": document.page_count},
            collections={},
        )
    finally:
        document.close()  # type: ignore[no-untyped-call]


def extract_artifact(path: Path, artifact_type: ArtifactType) -> ArtifactView:
    if path.suffix.casefold() != _SUFFIXES[artifact_type]:
        raise OfficeContentSuiteError(
            f"提交后缀 {path.suffix or '<none>'} 与期望 {artifact_type} 不一致"
        )
    if artifact_type == "docx":
        return _extract_docx(path)
    if artifact_type == "xlsx":
        return _extract_xlsx(path)
    if artifact_type == "pptx":
        return _extract_pptx(path)
    return _extract_pdf(path)


def _normalized(value: object) -> str:
    text = unicodedata.normalize("NFKC", str(value)).casefold()
    text = re.sub(
        r"(?<!\d)(\d{4})\s*[年./-]\s*(\d{1,2})\s*[月./-]\s*(\d{1,2})\s*日?",
        lambda match: (
            f"{int(match.group(1)):04d}-{int(match.group(2)):02d}-{int(match.group(3)):02d}"
        ),
        text,
    )
    text = re.sub(r"(?<=\d)[,，](?=\d{3}(?:\D|$))", "", text)
    # Preserve structural boundaries before removing presentation whitespace.
    # Otherwise adjacent shapes/cells such as badge ``02`` + value ``12``
    # collapse into ``0212`` and corrupt numeric token boundaries.
    text = re.sub(r"[\r\n\t]+", "|", text)
    return re.sub(r"\s+", "", text)


def _snippet(text: str, needle: str, *, limit: int = 220) -> str:
    folded = text.casefold()
    position = folded.find(needle.casefold())
    if position < 0:
        compact = re.sub(r"\s+", " ", text).strip()
        return compact[:limit]
    start = max(0, position - 70)
    end = min(len(text), position + len(needle) + 110)
    return re.sub(r"\s+", " ", text[start:end]).strip()[:limit]


def _scope_texts(view: ArtifactView, scope: ScopeSpec) -> list[TextScope]:
    if scope.kind == "artifact":
        return [TextScope(kind="artifact", label="whole artifact", text=view.text)]
    candidates = [item for item in view.scopes if item.kind == scope.kind]
    if scope.index is not None:
        return [item for item in candidates if item.index == scope.index]
    if scope.kind == "xlsx_range":
        assert scope.selector is not None and scope.cell_range is not None
        min_col, min_row, max_col, max_row = range_boundaries(scope.cell_range)
        rows: list[str] = []
        for row in range(min_row, max_row + 1):
            values: list[str] = []
            for column in range(min_col, max_col + 1):
                address = _cell_address(
                    scope.selector, f"{get_column_letter(column)}{row}"
                )
                value = view.cells.get(address)
                values.append("" if value is None else str(value))
            rows.append("\t".join(values).rstrip())
        return [
            TextScope(
                kind="xlsx_range",
                label=f"XLSX {scope.selector}!{scope.cell_range}",
                text="\n".join(rows),
                selector=scope.selector,
            )
        ]
    selector = _normalized(scope.selector or "")
    return [
        item
        for item in candidates
        if selector in _normalized(item.selector or item.label)
    ]


def _text_contains(check: TextContainsCheck, view: ArtifactView) -> CheckOutcome:
    scopes = _scope_texts(view, check.scope)
    if not scopes:
        return CheckOutcome(False, "定位范围不存在", [])
    normalized_values = [_normalized(value) for value in check.values]
    matches: list[tuple[str, str, str]] = []
    missing: list[str] = []
    for value, normalized_value in zip(check.values, normalized_values, strict=True):
        scope = next(
            (item for item in scopes if normalized_value in _normalized(item.text)),
            None,
        )
        if scope is None:
            missing.append(value)
        else:
            matches.append((value, scope.label, _snippet(scope.text, value)))
    passed = not missing if check.match == "all" else bool(matches)
    evidence = [f"{label}: {value}｜{snippet}" for value, label, snippet in matches]
    return CheckOutcome(
        passed,
        f"matched={len(matches)}/{len(check.values)} missing={missing!r}",
        evidence,
        {"matched": [item[0] for item in matches], "missing": missing},
    )


def _text_not_contains(check: TextNotContainsCheck, view: ArtifactView) -> CheckOutcome:
    scopes = _scope_texts(view, check.scope)
    if not scopes:
        return CheckOutcome(False, "定位范围不存在", [])
    hits: list[tuple[str, str, str]] = []
    for value in check.values:
        normalized_value = _normalized(value)
        for scope in scopes:
            if normalized_value in _normalized(scope.text):
                hits.append((value, scope.label, _snippet(scope.text, value)))
                break
    evidence = [
        f"{label}: forbidden={value}｜{snippet}" for value, label, snippet in hits
    ]
    return CheckOutcome(
        not hits, f"forbidden_hits={[item[0] for item in hits]!r}", evidence
    )


def _text_ordered(check: TextOrderedCheck, view: ArtifactView) -> CheckOutcome:
    scopes = _scope_texts(view, check.scope)
    if not scopes:
        return CheckOutcome(False, "定位范围不存在", [])
    for scope in scopes:
        text = _normalized(scope.text)
        cursor = 0
        positions: list[int] = []
        for value in check.values:
            position = text.find(_normalized(value), cursor)
            if position < 0:
                break
            positions.append(position)
            cursor = position + len(_normalized(value))
        if len(positions) == len(check.values):
            return CheckOutcome(
                True,
                f"ordered in {scope.label}",
                [f"{scope.label}: {' → '.join(check.values)}"],
                positions,
            )
    return CheckOutcome(False, f"未按顺序出现：{check.values!r}", [])


def _bounded_value_present(text: str, value: str) -> bool:
    """Match a normalized value without accepting a larger adjacent number."""

    if not value:
        return False
    starts_numeric = re.match(r"[+-]?\d", value) is not None
    ends_numeric = value[-1].isdigit() or value.endswith("%")
    left_boundary = r"(?<![\d.])" if starts_numeric else ""
    right_boundary = r"(?![\d.])" if ends_numeric else ""
    return (
        re.search(f"{left_boundary}{re.escape(value)}{right_boundary}", text)
        is not None
    )


def _labeled_value_match(
    scope: TextScope,
    requirement: LabeledValueRequirement,
    boundary_labels: list[str],
) -> tuple[str, list[str]] | None:
    text = _normalized(scope.text)
    normalized_values = [_normalized(value) for value in requirement.values]
    for raw_label in requirement.labels:
        label = _normalized(raw_label)
        cursor = 0
        while True:
            position = text.find(label, cursor)
            if position < 0:
                break
            if requirement.direction == "after":
                start = position + len(label)
                end = min(len(text), start + requirement.max_distance)
                # Stop at the next semantic label in the same check.  Without
                # this boundary, ``ARR 900；计划 860`` can satisfy ARR=860 by
                # borrowing the number from the neighbouring field.
                next_boundaries = [
                    next_position
                    for boundary in boundary_labels
                    if (next_position := text.find(boundary, start)) >= 0
                ]
                if next_boundaries:
                    end = min(end, min(next_boundaries))
            else:
                start = max(0, position - requirement.max_distance)
                end = min(
                    len(text),
                    position + len(label) + requirement.max_distance,
                )
                previous_boundaries = [
                    previous_position + len(boundary)
                    for boundary in boundary_labels
                    if (previous_position := text.rfind(boundary, start, position)) >= 0
                ]
                next_boundaries = [
                    next_position
                    for boundary in boundary_labels
                    if (
                        next_position := text.find(
                            boundary,
                            position + len(label),
                            end,
                        )
                    )
                    >= 0
                ]
                if previous_boundaries:
                    start = max(start, max(previous_boundaries))
                if next_boundaries:
                    end = min(end, min(next_boundaries))
            window = text[start:end]
            if all(
                _bounded_value_present(window, value) for value in normalized_values
            ):
                return raw_label, requirement.values
            cursor = position + max(1, len(label))
    return None


def _text_labeled_values(
    check: TextLabeledValuesCheck,
    view: ArtifactView,
) -> CheckOutcome:
    scopes = _scope_texts(view, check.scope)
    if not scopes:
        return CheckOutcome(False, "定位范围不存在", [])
    matches: list[dict[str, Any]] = []
    missing: list[dict[str, Any]] = []
    evidence: list[str] = []
    boundary_labels = list(
        dict.fromkeys(
            _normalized(label)
            for requirement in check.requirements
            for label in requirement.labels
        )
    )
    for requirement in check.requirements:
        found: tuple[TextScope, str, list[str]] | None = None
        for scope in scopes:
            match = _labeled_value_match(scope, requirement, boundary_labels)
            if match is not None:
                found = (scope, match[0], match[1])
                break
        if found is None:
            missing.append(
                {
                    "labels": requirement.labels,
                    "values": requirement.values,
                    "direction": requirement.direction,
                    "max_distance": requirement.max_distance,
                }
            )
            continue
        scope, label, values = found
        matches.append({"label": label, "values": values, "scope": scope.label})
        evidence.append(
            f"{scope.label}: {label} → {', '.join(values)}｜"
            f"{_snippet(scope.text, label)}"
        )
    return CheckOutcome(
        not missing,
        f"matched={len(matches)}/{len(check.requirements)} missing={missing!r}",
        evidence,
        {"matched": matches, "missing": missing},
    )


def _claim_hits(
    *,
    values: list[str],
    scopes: list[TextScope],
    negations: list[str],
    context_window: int,
) -> list[tuple[str, str, str]]:
    normalized_negations = [_normalized(value) for value in negations]
    hits: list[tuple[str, str, str]] = []
    for value in values:
        needle = _normalized(value)
        for scope in scopes:
            text = _normalized(scope.text)
            cursor = 0
            while True:
                position = text.find(needle, cursor)
                if position < 0:
                    break
                before = text[max(0, position - context_window) : position]
                after = text[
                    position + len(needle) : position + len(needle) + context_window
                ]
                negated = any(
                    negation in before or negation in after
                    for negation in normalized_negations
                )
                if not negated:
                    hits.append((value, scope.label, _snippet(scope.text, value)))
                    break
                cursor = position + len(needle)
    return hits


def _text_claim_contains(
    check: TextClaimContainsCheck, view: ArtifactView
) -> CheckOutcome:
    scopes = _scope_texts(view, check.scope)
    if not scopes:
        return CheckOutcome(False, "定位范围不存在", [])
    hits = _claim_hits(
        values=check.values,
        scopes=scopes,
        negations=check.negations,
        context_window=check.context_window,
    )
    hit_values = {item[0] for item in hits}
    passed = (
        all(value in hit_values for value in check.values)
        if check.match == "all"
        else bool(hits)
    )
    return CheckOutcome(
        passed,
        f"unnegated_claims={sorted(hit_values)!r}",
        [f"{label}: claim={value}｜{snippet}" for value, label, snippet in hits],
        sorted(hit_values),
    )


def _text_claim_not_contains(
    check: TextClaimNotContainsCheck, view: ArtifactView
) -> CheckOutcome:
    scopes = _scope_texts(view, check.scope)
    if not scopes:
        return CheckOutcome(False, "定位范围不存在", [])
    hits = _claim_hits(
        values=check.values,
        scopes=scopes,
        negations=check.negations,
        context_window=check.context_window,
    )
    return CheckOutcome(
        not hits,
        f"unnegated_claims={[item[0] for item in hits]!r}",
        [f"{label}: claim={value}｜{snippet}" for value, label, snippet in hits],
        [item[0] for item in hits],
    )


def _structure_count(check: StructureCountCheck, view: ArtifactView) -> CheckOutcome:
    actual = view.counts.get(check.metric)
    if actual is None:
        return CheckOutcome(False, f"当前格式没有指标 {check.metric}", [], None)
    comparisons = {
        "eq": actual == check.value,
        "gte": actual >= check.value,
        "lte": actual <= check.value,
    }
    return CheckOutcome(
        comparisons[check.operator],
        f"{check.metric}: actual={actual} {check.operator} expected={check.value}",
        [f"{check.metric}={actual}"],
        actual,
    )


def _structure_values(check: StructureValuesCheck, view: ArtifactView) -> CheckOutcome:
    actual = view.collections.get(check.metric)
    if actual is None:
        return CheckOutcome(False, f"当前格式没有集合 {check.metric}", [], None)
    normalized_actual = [_normalized(value) for value in actual]
    expected = [_normalized(value) for value in check.values]
    if check.match == "exact_order":
        passed = normalized_actual == expected
        missing = [] if passed else check.values
    else:
        present = [
            value
            for value in expected
            if any(value in item for item in normalized_actual)
        ]
        passed = (
            len(present) == len(expected) if check.match == "all" else bool(present)
        )
        missing = [
            raw
            for raw, value in zip(check.values, expected, strict=True)
            if value not in present
        ]
    return CheckOutcome(
        passed,
        f"actual={actual!r} missing={missing!r}",
        [f"{check.metric}={actual!r}"],
        actual,
    )


def _pptx_chart_data(check: PptxChartDataCheck, view: ArtifactView) -> CheckOutcome:
    def matches_expected(actual: float, expected: float) -> bool:
        candidates = [expected]
        if check.accept_percentage_fractions and abs(expected) > 1:
            candidates.append(expected / 100.0)
        return any(
            math.isclose(actual, candidate, rel_tol=0.0, abs_tol=check.tolerance)
            for candidate in candidates
        )

    for chart in view.charts:
        actual_values = [
            float(value)
            for value in chart.get("values", [])
            if isinstance(value, (int, float))
        ]
        labels = [
            _normalized(value)
            for value in [
                *chart.get("labels", []),
                *chart.get("series_names", []),
            ]
        ]
        values_pass = all(
            any(matches_expected(actual, expected) for actual in actual_values)
            for expected in check.required_values
        )
        labels_pass = all(
            any(_normalized(expected) in actual for actual in labels)
            for expected in check.required_labels
        )
        if values_pass and labels_pass:
            detail = {
                "slide_index": chart.get("slide_index"),
                "title": chart.get("title"),
                "labels": chart.get("labels"),
                "series_names": chart.get("series_names"),
                "values": actual_values,
            }
            return CheckOutcome(
                True,
                f"matched native chart on slide {chart.get('slide_index')}",
                [f"PPTX chart: {detail!r}"],
                detail,
            )
    return CheckOutcome(
        False,
        "没有同一原生图表同时满足所需标签和数值",
        [],
        {
            "required_labels": check.required_labels,
            "required_values": check.required_values,
            "charts": view.charts,
        },
    )


def _lookup_cell(view: ArtifactView, address: str) -> tuple[str, Any]:
    normalized_address = address.replace("$", "")
    for key, value in view.cells.items():
        if _normalized(key) == _normalized(normalized_address):
            return key, value
    raise KeyError(address)


def _split_address(address: str) -> tuple[str, str]:
    sheet, coordinate = address.rsplit("!", maxsplit=1)
    return sheet.strip("'"), coordinate.replace("$", "").upper()


def _range_values(
    view: ArtifactView,
    sheet: str,
    cell_range: str,
    stack: set[str],
) -> list[float]:
    min_col, min_row, max_col, max_row = range_boundaries(cell_range.replace("$", ""))
    values: list[float] = []
    for row in range(min_row, max_row + 1):
        for column in range(min_col, max_col + 1):
            address = f"{sheet}!{get_column_letter(column)}{row}"
            try:
                value = _numeric_cell_value(view, address, stack)
            except (KeyError, OfficeContentSuiteError):
                continue
            values.append(value)
    return values


_RANGE_REF = re.compile(
    r"(?:(?:'([^']+)'|([A-Za-z0-9_\u4e00-\u9fff .-]+))!)?"
    r"(\$?[A-Z]{1,3}\$?[1-9][0-9]*:\$?[A-Z]{1,3}\$?[1-9][0-9]*)"
)
_CELL_REF = re.compile(
    r"(?<![A-Za-z0-9_])(?:(?:'([^']+)'|([A-Za-z0-9_\u4e00-\u9fff .-]+))!)?"
    r"(\$?[A-Z]{1,3}\$?[1-9][0-9]*)"
)
_FORMULA_RANGE_REFERENCE = (
    r"(?:(?:'[^']+'|[^!,()]+)!)?"
    r"\$?[A-Z]{1,3}\$?[1-9][0-9]*:\$?[A-Z]{1,3}\$?[1-9][0-9]*"
)
_SUMIF_FORMULA = re.compile(
    rf"^\s*=\s*SUMIF\(\s*(?P<criteria>{_FORMULA_RANGE_REFERENCE})\s*,"
    rf'\s*"(?P<criterion>[^"]+)"\s*,\s*(?P<total>{_FORMULA_RANGE_REFERENCE})\s*\)\s*$',
    re.IGNORECASE,
)


FormulaValue = float | list[float]


def _as_scalar(value: FormulaValue) -> float:
    if isinstance(value, list):
        raise OfficeContentSuiteError("公式在需要单值的位置返回了区域")
    return value


def _flatten(values: tuple[FormulaValue, ...]) -> list[float]:
    result: list[float] = []
    for value in values:
        if isinstance(value, list):
            result.extend(float(item) for item in value)
        else:
            result.append(float(value))
    return result


def _evaluate_formula_ast(node: ast.AST, names: dict[str, Any]) -> FormulaValue:
    if isinstance(node, ast.Expression):
        return _evaluate_formula_ast(node.body, names)
    if isinstance(node, ast.Constant) and isinstance(node.value, (int, float)):
        return float(node.value)
    if isinstance(node, ast.Name) and node.id in names:
        value = names[node.id]
        return value if isinstance(value, list) else float(value)
    if isinstance(node, ast.UnaryOp) and isinstance(node.op, (ast.UAdd, ast.USub)):
        value = _as_scalar(_evaluate_formula_ast(node.operand, names))
        return value if isinstance(node.op, ast.UAdd) else -value
    if isinstance(node, ast.BinOp) and isinstance(
        node.op, (ast.Add, ast.Sub, ast.Mult, ast.Div, ast.Pow)
    ):
        left = _as_scalar(_evaluate_formula_ast(node.left, names))
        right = _as_scalar(_evaluate_formula_ast(node.right, names))
        if isinstance(node.op, ast.Add):
            return left + right
        if isinstance(node.op, ast.Sub):
            return left - right
        if isinstance(node.op, ast.Mult):
            return left * right
        if isinstance(node.op, ast.Div):
            return left / right
        result = left**right
        if isinstance(result, complex):
            raise OfficeContentSuiteError("公式幂运算产生复数")
        return float(result)
    if isinstance(node, ast.Call) and isinstance(node.func, ast.Name):
        function = node.func.id.upper()
        arguments = tuple(_evaluate_formula_ast(item, names) for item in node.args)
        values = _flatten(arguments)
        if function == "SUM":
            return sum(values)
        if function == "AVERAGE":
            return sum(values) / len(values)
        if function == "MIN":
            return min(values)
        if function == "MAX":
            return max(values)
        if function == "ROUND" and len(arguments) == 2:
            return float(round(_as_scalar(arguments[0]), int(_as_scalar(arguments[1]))))
    raise OfficeContentSuiteError(
        f"暂不支持的公式表达式：{ast.dump(node, include_attributes=False)}"
    )


def _range_addresses(reference: str, default_sheet: str) -> list[str]:
    if "!" in reference:
        sheet, cell_range = reference.rsplit("!", maxsplit=1)
        sheet = sheet.strip().strip("'")
    else:
        sheet = default_sheet
        cell_range = reference
    min_col, min_row, max_col, max_row = range_boundaries(cell_range.replace("$", ""))
    return [
        f"{sheet}!{get_column_letter(column)}{row}"
        for row in range(min_row, max_row + 1)
        for column in range(min_col, max_col + 1)
    ]


def _evaluate_sumif(
    view: ArtifactView,
    raw: str,
    default_sheet: str,
    stack: set[str],
) -> float | None:
    match = _SUMIF_FORMULA.fullmatch(raw)
    if match is None:
        return None
    criteria_addresses = _range_addresses(match.group("criteria"), default_sheet)
    total_addresses = _range_addresses(match.group("total"), default_sheet)
    if len(criteria_addresses) != len(total_addresses):
        raise OfficeContentSuiteError("SUMIF 条件区域与求和区域大小不一致")
    criterion = _normalized(match.group("criterion"))
    total = 0.0
    for criteria_address, total_address in zip(
        criteria_addresses, total_addresses, strict=True
    ):
        try:
            _, criteria_value = _lookup_cell(view, criteria_address)
        except KeyError:
            continue
        if _normalized(criteria_value) != criterion:
            continue
        try:
            total += _numeric_cell_value(view, total_address, stack)
        except KeyError:
            continue
    return total


def _numeric_cell_value(
    view: ArtifactView, address: str, stack: set[str] | None = None
) -> float:
    key, raw = _lookup_cell(view, address)
    if isinstance(raw, bool):
        return float(raw)
    if isinstance(raw, (int, float)):
        return float(raw)
    if not isinstance(raw, str) or not raw.startswith("="):
        try:
            return float(raw)
        except (TypeError, ValueError) as error:
            raise OfficeContentSuiteError(
                f"{key} 不是数值或可计算公式：{raw!r}"
            ) from error
    active = set() if stack is None else set(stack)
    if key in active:
        raise OfficeContentSuiteError(f"公式循环引用：{key}")
    active.add(key)
    sheet, _ = _split_address(key)
    sumif_value = _evaluate_sumif(view, raw, sheet, active)
    if sumif_value is not None:
        return sumif_value
    expression = raw[1:].replace("^", "**")
    expression = re.sub(r"(?<![A-Za-z0-9_.])(\d+(?:\.\d+)?)%", r"(\1/100)", expression)
    names: dict[str, Any] = {}

    def replace_range(match: re.Match[str]) -> str:
        range_sheet = match.group(1) or match.group(2) or sheet
        name = f"__range_{len(names)}"
        names[name] = _range_values(view, range_sheet, match.group(3), active)
        return name

    expression = _RANGE_REF.sub(replace_range, expression)

    def replace_cell(match: re.Match[str]) -> str:
        cell_sheet = match.group(1) or match.group(2) or sheet
        name = f"__cell_{len(names)}"
        names[name] = _numeric_cell_value(
            view, f"{cell_sheet}!{match.group(3).replace('$', '')}", active
        )
        return name

    expression = _CELL_REF.sub(replace_cell, expression)
    try:
        tree = ast.parse(expression, mode="eval")
        result = _evaluate_formula_ast(tree, names)
    except (SyntaxError, ZeroDivisionError, ValueError) as error:
        raise OfficeContentSuiteError(f"无法计算 {key}={raw!r}：{error}") from error
    if isinstance(result, list):
        raise OfficeContentSuiteError(f"{key} 公式结果是区域而非单值")
    return float(result)


def _xlsx_cell_value(check: XlsxCellValueCheck, view: ArtifactView) -> CheckOutcome:
    try:
        key, raw = _lookup_cell(view, check.address)
    except KeyError:
        return CheckOutcome(False, f"单元格不存在：{check.address}", [])
    if isinstance(check.expected, (int, float)) and not isinstance(
        check.expected, bool
    ):
        try:
            actual: Any = _numeric_cell_value(view, key)
        except OfficeContentSuiteError as error:
            return CheckOutcome(False, str(error), [f"{key}={raw!r}"], raw)
        passed = math.isclose(
            float(actual), float(check.expected), rel_tol=0.0, abs_tol=check.tolerance
        )
    else:
        actual = raw
        passed = _normalized(actual) == _normalized(check.expected)
    return CheckOutcome(
        passed,
        f"{key}: actual={actual!r} expected={check.expected!r} tolerance={check.tolerance}",
        [f"{key}={raw!r}"],
        actual,
    )


def _xlsx_formula_value(
    check: XlsxFormulaValueCheck, view: ArtifactView
) -> CheckOutcome:
    try:
        key, raw = _lookup_cell(view, check.address)
    except KeyError:
        return CheckOutcome(False, f"单元格不存在：{check.address}", [])
    if not isinstance(raw, str) or not raw.startswith("="):
        return CheckOutcome(
            False, f"{key} 必须是公式，实际为 {raw!r}", [f"{key}={raw!r}"], raw
        )
    missing_tokens = [
        token
        for token in check.must_contain
        if _normalized(token) not in _normalized(raw)
    ]
    try:
        actual = _numeric_cell_value(view, key)
    except OfficeContentSuiteError as error:
        return CheckOutcome(False, str(error), [f"{key}={raw}"], raw)
    value_ok = math.isclose(
        actual, check.expected, rel_tol=0.0, abs_tol=check.tolerance
    )
    return CheckOutcome(
        value_ok and not missing_tokens,
        f"{key}: formula={raw!r} result={actual:g} expected={check.expected:g} missing_tokens={missing_tokens!r}",
        [f"{key}={raw} → {actual:g}"],
        actual,
    )


def _xlsx_cells_value(check: XlsxCellsValueCheck, view: ArtifactView) -> CheckOutcome:
    outcomes = [
        _xlsx_cell_value(
            XlsxCellValueCheck(
                type="xlsx_cell_value",
                address=address,
                expected=expected,
                tolerance=check.tolerance,
            ),
            view,
        )
        for address, expected in check.cells.items()
    ]
    failures = [outcome.detail for outcome in outcomes if not outcome.passed]
    return CheckOutcome(
        not failures,
        f"matched={len(outcomes) - len(failures)}/{len(outcomes)} failures={failures!r}",
        [evidence for outcome in outcomes for evidence in outcome.evidence],
        {
            address: outcome.actual
            for address, outcome in zip(check.cells, outcomes, strict=True)
        },
    )


def _xlsx_number_format(
    check: XlsxNumberFormatCheck, view: ArtifactView
) -> CheckOutcome:
    try:
        key, _ = _lookup_cell(view, check.address)
    except KeyError:
        return CheckOutcome(False, f"单元格不存在：{check.address}", [])
    number_format = view.number_formats.get(key, "")
    passed = any(
        token.casefold() in number_format.casefold() for token in check.contains_any
    )
    return CheckOutcome(
        passed,
        f"{key}: number_format={number_format!r} expected_any={check.contains_any!r}",
        [f"{key} number_format={number_format!r}"],
        number_format,
    )


def _validator_status(
    check: ValidatorStatusCheck, report: ArtifactValidationReport
) -> CheckOutcome:
    dimension = getattr(report, check.dimension)
    label: str
    if check.check_name is None:
        actual: ValidationStatus = dimension.status
        label = check.dimension
    else:
        found = next(
            (item for item in dimension.checks if item.name == check.check_name), None
        )
        if found is None:
            return CheckOutcome(
                False,
                f"validator check 不存在：{check.dimension}.{check.check_name}",
                [],
            )
        actual = found.status
        label = f"{check.dimension}.{check.check_name}"
    return CheckOutcome(
        actual in check.allowed,
        f"{label}: actual={actual} allowed={check.allowed!r}",
        [f"{label}={actual}"],
        actual,
    )


def evaluate_check(
    check: CheckSpec,
    view: ArtifactView,
    report: ArtifactValidationReport,
) -> CheckOutcome:
    if isinstance(check, TextContainsCheck):
        return _text_contains(check, view)
    if isinstance(check, TextNotContainsCheck):
        return _text_not_contains(check, view)
    if isinstance(check, TextOrderedCheck):
        return _text_ordered(check, view)
    if isinstance(check, TextLabeledValuesCheck):
        return _text_labeled_values(check, view)
    if isinstance(check, TextClaimContainsCheck):
        return _text_claim_contains(check, view)
    if isinstance(check, TextClaimNotContainsCheck):
        return _text_claim_not_contains(check, view)
    if isinstance(check, StructureCountCheck):
        return _structure_count(check, view)
    if isinstance(check, StructureValuesCheck):
        return _structure_values(check, view)
    if isinstance(check, XlsxCellValueCheck):
        return _xlsx_cell_value(check, view)
    if isinstance(check, XlsxCellsValueCheck):
        return _xlsx_cells_value(check, view)
    if isinstance(check, XlsxFormulaValueCheck):
        return _xlsx_formula_value(check, view)
    if isinstance(check, XlsxNumberFormatCheck):
        return _xlsx_number_format(check, view)
    if isinstance(check, PptxChartDataCheck):
        return _pptx_chart_data(check, view)
    return _validator_status(check, report)


def _artifact_sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _gate(
    item: OfficeContentItem,
    path: Path,
) -> tuple[bool, list[str], ArtifactView | None, ArtifactValidationReport | None]:
    if not path.is_file():
        return False, [f"missing artifact: {path}"], None, None
    file_size = path.stat().st_size
    if file_size > item.gate.max_file_bytes:
        return (
            False,
            [f"file_size={file_size} > max_file_bytes={item.gate.max_file_bytes}"],
            None,
            None,
        )
    if item.artifact_type in {"docx", "xlsx", "pptx"}:
        try:
            with zipfile.ZipFile(path) as archive:
                uncompressed_size = sum(info.file_size for info in archive.infolist())
        except (OSError, zipfile.BadZipFile) as error:
            return False, [f"invalid OOXML archive: {error}"], None, None
        if uncompressed_size > item.gate.max_uncompressed_bytes:
            return (
                False,
                [
                    f"uncompressed_size={uncompressed_size} > max_uncompressed_bytes={item.gate.max_uncompressed_bytes}"
                ],
                None,
                None,
            )
    try:
        view = extract_artifact(path, item.artifact_type)
        report = validate_artifact(path, render_visual=item.gate.render_visual)
    except Exception as error:
        return False, [f"{type(error).__name__}: {error}"], None, None
    reasons: list[str] = []
    if report.artifact_type != item.artifact_type:
        reasons.append(
            f"artifact_type expected={item.artifact_type} actual={report.artifact_type}"
        )
    for name in item.gate.fail_on_dimensions:
        dimension = getattr(report, name)
        if dimension.status == "failed":
            failed_checks = [
                check.name for check in dimension.checks if check.status == "failed"
            ]
            reasons.append(f"{name}=failed checks={failed_checks!r}")
    for name in item.gate.require_measured_dimensions:
        if getattr(report, name).status == "not_run":
            reasons.append(f"{name}=not_run")
    if report.quality.score < item.gate.min_validator_quality:
        reasons.append(
            f"validator_quality={report.quality.score} < {item.gate.min_validator_quality}"
        )
    return not reasons, reasons, view, report


def _review_index(
    suite: OfficeContentSuite,
    reviews: ReviewFile | None,
) -> dict[tuple[str, str], ReviewAnnotation]:
    if reviews is None:
        return {}
    items = {item.id: item for item in suite.items}
    result: dict[tuple[str, str], ReviewAnnotation] = {}
    for annotation in reviews.reviews:
        item = items.get(annotation.item_id)
        if item is None:
            raise OfficeContentSuiteError(f"复核引用未知 item：{annotation.item_id}")
        criterion = next(
            (
                value
                for value in item.review_criteria
                if value.id == annotation.criterion_id
            ),
            None,
        )
        if criterion is None:
            raise OfficeContentSuiteError(
                f"复核引用未知 criterion：{annotation.item_id}/{annotation.criterion_id}"
            )
        if annotation.score > criterion.max_score:
            raise OfficeContentSuiteError(
                f"{annotation.item_id}/{annotation.criterion_id} score 超过 {criterion.max_score}"
            )
        result[(annotation.item_id, annotation.criterion_id)] = annotation
    return result


def _score_item(
    suite: OfficeContentSuite,
    item: OfficeContentItem,
    submission_root: Path,
    review_index: dict[tuple[str, str], ReviewAnnotation],
) -> tuple[dict[str, Any], dict[str, Any]]:
    path = submission_root / item.id / "submission" / item.output_file
    gate_passed, gate_reasons, view, validation = _gate(item, path)
    sha256 = _artifact_sha256(path) if path.is_file() else None
    dimension_totals = {name: 0.0 for name in _DIMENSIONS}
    dimension_earned = {name: 0.0 for name in _DIMENSIONS}
    rubric_results: list[dict[str, Any]] = []
    critical_failures: list[str] = []
    if view is not None and validation is not None:
        for rubric in item.rubric:
            outcome = evaluate_check(rubric.check, view, validation)
            dimension_totals[rubric.dimension] += rubric.weight
            dimension_earned[rubric.dimension] += rubric.weight * int(outcome.passed)
            if rubric.critical and not outcome.passed:
                critical_failures.append(rubric.id)
            rubric_results.append(
                {
                    "id": rubric.id,
                    "dimension": rubric.dimension,
                    "description": rubric.description,
                    "source_refs": rubric.source_refs,
                    "weight": rubric.weight,
                    "critical": rubric.critical,
                    "passed": outcome.passed,
                    "detail": outcome.detail,
                    "evidence": outcome.evidence,
                    "actual": _json_safe(outcome.actual),
                }
            )
    else:
        for rubric in item.rubric:
            dimension_totals[rubric.dimension] += rubric.weight
            if rubric.critical:
                critical_failures.append(rubric.id)
            rubric_results.append(
                {
                    "id": rubric.id,
                    "dimension": rubric.dimension,
                    "description": rubric.description,
                    "source_refs": rubric.source_refs,
                    "weight": rubric.weight,
                    "critical": rubric.critical,
                    "passed": False,
                    "detail": "gate 未得到可解析文件，检查未运行",
                    "evidence": [],
                    "actual": None,
                }
            )
    dimension_scores = {
        name: round(
            100.0 * dimension_earned[name] / dimension_totals[name]
            if dimension_totals[name]
            else 0.0,
            2,
        )
        for name in _DIMENSIONS
    }
    raw_automatic = sum(
        dimension_scores[name] * suite.dimension_weights[name] for name in _DIMENSIONS
    )
    penalty_results: list[dict[str, Any]] = []
    penalty_points = 0.0
    blocking_penalties: list[str] = []
    if view is not None and validation is not None:
        for penalty in item.penalties:
            outcome = evaluate_check(penalty.trigger, view, validation)
            hit = outcome.passed
            if hit:
                penalty_points += penalty.points
                if penalty.blocking:
                    blocking_penalties.append(penalty.id)
            penalty_results.append(
                {
                    "id": penalty.id,
                    "description": penalty.description,
                    "source_refs": penalty.source_refs,
                    "points": penalty.points,
                    "blocking": penalty.blocking,
                    "hit": hit,
                    "detail": outcome.detail,
                    "evidence": outcome.evidence,
                }
            )
    automatic_score = (
        round(max(0.0, raw_automatic - penalty_points), 2) if gate_passed else 0.0
    )
    automatic_pass = bool(
        gate_passed
        and automatic_score >= item.pass_threshold
        and not critical_failures
        and not blocking_penalties
    )

    review_results: list[dict[str, Any]] = []
    review_earned = 0
    review_total = 0
    review_dimension_earned = {name: 0 for name in _REVIEW_DIMENSIONS}
    review_dimension_total = {name: 0 for name in _REVIEW_DIMENSIONS}
    review_annotations: list[ReviewAnnotation] = []
    for criterion in item.review_criteria:
        annotation = review_index.get((item.id, criterion.id))
        stale = bool(annotation is not None and annotation.artifact_sha256 != sha256)
        if annotation is not None and not stale:
            review_annotations.append(annotation)
            review_earned += annotation.score
            review_total += criterion.max_score
            review_dimension_earned[criterion.dimension] += annotation.score
            review_dimension_total[criterion.dimension] += criterion.max_score
        review_results.append(
            {
                "id": criterion.id,
                "dimension": criterion.dimension,
                "description": criterion.description,
                "anchors": criterion.anchors,
                "max_score": criterion.max_score,
                "minimum_score": criterion.minimum_score,
                "status": "stale" if stale else "scored" if annotation else "pending",
                "score": annotation.score
                if annotation is not None and not stale
                else None,
                "evidence": annotation.evidence
                if annotation is not None and not stale
                else None,
                "reviewer": annotation.reviewer
                if annotation is not None and not stale
                else None,
                "source": annotation.source
                if annotation is not None and not stale
                else None,
                "model_provenance": (
                    {
                        "provider": annotation.provider,
                        "model": annotation.model,
                        "prompt_fingerprint": annotation.prompt_fingerprint,
                        "authorization_note_fingerprint": (
                            annotation.authorization_note_fingerprint
                        ),
                        "calibration_status": annotation.calibration_status,
                        "render_mode": annotation.render_mode,
                    }
                    if annotation is not None
                    and not stale
                    and annotation.source == "model"
                    else None
                ),
            }
        )
    review_complete = bool(review_results) and all(
        result["status"] == "scored" for result in review_results
    )
    review_score = (
        round(100.0 * review_earned / review_total, 2) if review_complete else None
    )
    review_dimension_scores: dict[str, float] | None = (
        {
            name: round(
                100.0 * review_dimension_earned[name] / review_dimension_total[name],
                2,
            )
            for name in _REVIEW_DIMENSIONS
        }
        if review_complete
        else None
    )
    review_failures = (
        [
            str(result["id"])
            for result in review_results
            if int(result["score"]) < int(result["minimum_score"])
        ]
        if review_complete
        else []
    )
    review_pass = not review_failures if review_complete else None
    review_eligibility_failures = (
        [
            f"{annotation.criterion_id}:uncalibrated_model_review"
            for annotation in review_annotations
            if annotation.source == "model"
        ]
        if review_complete
        else []
    )
    review_eligible_for_benchmark = (
        not review_eligibility_failures if review_complete else None
    )
    final_score = (
        round(
            suite.automatic_weight * automatic_score
            + suite.review_weight * float(review_score),
            2,
        )
        if review_score is not None
        else None
    )
    engineering_pass = (
        automatic_pass and review_pass is True and final_score >= item.pass_threshold
        if final_score is not None and review_pass is not None
        else None
    )
    benchmark_pass = (
        engineering_pass is True and review_eligible_for_benchmark is True
        if engineering_pass is not None and review_eligible_for_benchmark is not None
        else None
    )
    validation_payload = (
        validation.model_dump(mode="json") if validation is not None else None
    )
    result = {
        "id": item.id,
        "split": item.split,
        "artifact_type": item.artifact_type,
        "task_type": item.task_type,
        "category": item.category,
        "difficulty": item.difficulty,
        "artifact": str(path),
        "artifact_sha256": sha256,
        "gate": {
            "passed": gate_passed,
            "reasons": gate_reasons,
            "validation": validation_payload,
        },
        "rubric": rubric_results,
        "dimension_scores": dimension_scores,
        "raw_automatic_score": round(raw_automatic, 2),
        "penalties": penalty_results,
        "penalty_points": round(penalty_points, 2),
        "critical_failures": critical_failures,
        "blocking_penalties": blocking_penalties,
        "automatic_score": automatic_score,
        "automatic_pass": automatic_pass,
        "review": review_results,
        "review_complete": review_complete,
        "review_score": review_score,
        "review_dimension_scores": review_dimension_scores,
        "review_failures": review_failures,
        "review_pass": review_pass,
        "review_sources": sorted(
            {annotation.source for annotation in review_annotations}
        ),
        "review_eligibility_failures": review_eligibility_failures,
        "review_eligible_for_benchmark": review_eligible_for_benchmark,
        "final_score": final_score,
        "engineering_pass": engineering_pass,
        "benchmark_pass": benchmark_pass,
    }
    packet = {
        "item_id": item.id,
        "artifact": str(path),
        "artifact_sha256": sha256,
        "prompt": item.prompt,
        "automatic_score": automatic_score,
        "failed_automatic_items": [
            {
                "id": rubric["id"],
                "description": rubric["description"],
                "detail": rubric["detail"],
            }
            for rubric in rubric_results
            if not rubric["passed"]
        ],
        "criteria": [
            {
                "criterion_id": criterion.id,
                "dimension": criterion.dimension,
                "description": criterion.description,
                "anchors": criterion.anchors,
                "minimum_score": criterion.minimum_score,
            }
            for criterion in item.review_criteria
        ],
    }
    return result, packet


def _select_items(
    suite: OfficeContentSuite,
    split: Literal["dev", "test", "all"],
    *,
    include_test: bool,
    test_access_note: str | None,
) -> list[OfficeContentItem]:
    if split in {"test", "all"}:
        if not include_test:
            raise OfficeContentSuiteError("访问 test 必须显式传入 include_test")
        if not test_access_note or not test_access_note.strip():
            raise OfficeContentSuiteError("访问 test 必须记录非空 test_access_note")
    selected = [item for item in suite.items if split == "all" or item.split == split]
    if not selected:
        raise OfficeContentSuiteError(f"suite 中没有 split={split} 的任务")
    return selected


def prepare_suite(
    suite: OfficeContentSuite,
    workspace_root: Path,
    *,
    split: Literal["dev", "test", "all"] = "dev",
    include_test: bool = False,
    test_access_note: str | None = None,
) -> dict[str, Any]:
    if workspace_root.exists():
        raise OfficeContentSuiteError("workspace_root 已存在；prepare 不覆盖既有目录")
    items = _select_items(
        suite,
        split,
        include_test=include_test,
        test_access_note=test_access_note,
    )
    workspace_root.mkdir(parents=True)
    prepared: list[dict[str, Any]] = []
    for item in items:
        case_root = workspace_root / item.id
        input_root = case_root / "inputs"
        submission_root = case_root / "submission"
        input_root.mkdir(parents=True)
        submission_root.mkdir()
        for fixture in item.fixtures:
            path = input_root.joinpath(*PurePosixPath(fixture.path).parts)
            path.parent.mkdir(parents=True, exist_ok=True)
            path.write_text(fixture.content, encoding="utf-8")
        task_payload = {
            "id": item.id,
            "artifact_type": item.artifact_type,
            "prompt": item.prompt,
            "inputs": [f"inputs/{fixture.path}" for fixture in item.fixtures],
            "required_output": f"submission/{item.output_file}",
        }
        (case_root / "task.json").write_text(
            json.dumps(task_payload, ensure_ascii=False, indent=2, sort_keys=True)
            + "\n",
            encoding="utf-8",
        )
        (case_root / "TASK.md").write_text(
            f"# {item.id}\n\n{item.prompt}\n\n"
            f"输入资料位于 `inputs/`；最终文件必须保存为 "
            f"`submission/{item.output_file}`。\n",
            encoding="utf-8",
        )
        prepared.append(task_payload)
    manifest = {
        "schema_version": SCHEMA_VERSION,
        "suite": suite_summary(suite),
        "suite_sha256": _suite_sha256(suite),
        "split": split,
        "test_access_note": test_access_note.strip() if test_access_note else None,
        "items": prepared,
    }
    (workspace_root / "manifest.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2, sort_keys=True) + "\n",
        encoding="utf-8",
    )
    return manifest


def _group_summary(
    results: list[dict[str, Any]],
    key: Literal["artifact_type", "task_type", "difficulty"],
) -> dict[str, dict[str, Any]]:
    totals = Counter(str(item[key]) for item in results)
    return {
        name: {
            "items": count,
            "gate_passed": sum(
                bool(item["gate"]["passed"]) for item in results if item[key] == name
            ),
            "automatic_passed": sum(
                bool(item["automatic_pass"]) for item in results if item[key] == name
            ),
            "mean_automatic_score": round(
                sum(
                    float(item["automatic_score"])
                    for item in results
                    if item[key] == name
                )
                / count,
                2,
            ),
        }
        for name, count in sorted(totals.items())
    }


def evaluate_suite(
    suite: OfficeContentSuite,
    submission_root: Path,
    output_root: Path,
    *,
    split: Literal["dev", "test", "all"] = "dev",
    include_test: bool = False,
    test_access_note: str | None = None,
    reviews: ReviewFile | None = None,
) -> dict[str, Any]:
    if output_root.exists():
        raise OfficeContentSuiteError("output_root 已存在；score 不覆盖既有结果")
    items = _select_items(
        suite,
        split,
        include_test=include_test,
        test_access_note=test_access_note,
    )
    review_index = _review_index(suite, reviews)
    output_root.mkdir(parents=True)
    scored = [_score_item(suite, item, submission_root, review_index) for item in items]
    results = [item[0] for item in scored]
    packets = [item[1] for item in scored]
    gate_passes = sum(bool(item["gate"]["passed"]) for item in results)
    automatic_passes = sum(bool(item["automatic_pass"]) for item in results)
    adjudicated = [item for item in results if item["review_complete"]]
    review_passes = sum(item["review_pass"] is True for item in adjudicated)
    engineering_passes = sum(item["engineering_pass"] is True for item in adjudicated)
    benchmark_eligible = sum(
        item["review_eligible_for_benchmark"] is True for item in adjudicated
    )
    benchmark_passes = sum(bool(item["benchmark_pass"]) for item in adjudicated)
    summary = {
        "items": len(results),
        "gate_passed": gate_passes,
        "gate_pass_rate": gate_passes / len(results) if results else 0.0,
        "automatic_passed": automatic_passes,
        "automatic_pass_rate": automatic_passes / len(results) if results else 0.0,
        "mean_automatic_score": round(
            sum(float(item["automatic_score"]) for item in results) / len(results), 2
        )
        if results
        else 0.0,
        "mean_dimension_scores": {
            dimension: round(
                sum(float(item["dimension_scores"][dimension]) for item in results)
                / len(results),
                2,
            )
            if results
            else 0.0
            for dimension in _DIMENSIONS
        },
        "review_complete_items": len(adjudicated),
        "review_coverage": len(adjudicated) / len(results) if results else 0.0,
        "review_passed": review_passes,
        "review_pass_rate": review_passes / len(adjudicated) if adjudicated else None,
        "engineering_passed": engineering_passes,
        "engineering_pass_rate": engineering_passes / len(adjudicated)
        if adjudicated
        else None,
        "benchmark_eligible_items": benchmark_eligible,
        "benchmark_eligibility_rate": benchmark_eligible / len(adjudicated)
        if adjudicated
        else None,
        "mean_review_dimension_scores": {
            dimension: round(
                sum(
                    float(item["review_dimension_scores"][dimension])
                    for item in adjudicated
                )
                / len(adjudicated),
                2,
            )
            for dimension in _REVIEW_DIMENSIONS
        }
        if adjudicated
        else None,
        "mean_final_score": round(
            sum(float(item["final_score"]) for item in adjudicated) / len(adjudicated),
            2,
        )
        if adjudicated
        else None,
        "benchmark_passed": benchmark_passes,
        "benchmark_pass_rate": benchmark_passes / len(adjudicated)
        if adjudicated
        else None,
        "by_format": _group_summary(results, "artifact_type"),
        "by_task_type": _group_summary(results, "task_type"),
        "by_difficulty": _group_summary(results, "difficulty"),
    }
    payload = {
        "schema_version": REPORT_SCHEMA_VERSION,
        "generated_at": datetime.now(UTC).isoformat(),
        "suite": suite_summary(suite),
        "suite_sha256": _suite_sha256(suite),
        "scorer_fingerprint": _scorer_fingerprint(),
        "split": split,
        "test_access_note": test_access_note.strip() if test_access_note else None,
        "scoring": {
            "dimension_weights": suite.dimension_weights,
            "automatic_weight": suite.automatic_weight,
            "review_weight": suite.review_weight,
            "gate_failure_score": 0,
            "final_score_requires_complete_review": True,
            "review_minimum_scores_enforced": True,
            "uncalibrated_model_reviews_are_benchmark_ineligible": True,
        },
        "summary": summary,
        "results": results,
    }
    (output_root / "report.json").write_text(
        json.dumps(payload, ensure_ascii=False, indent=2, sort_keys=True) + "\n",
        encoding="utf-8",
    )
    with (output_root / "review-packet.jsonl").open("w", encoding="utf-8") as handle:
        for packet in packets:
            handle.write(json.dumps(packet, ensure_ascii=False, sort_keys=True) + "\n")
    return payload


def _main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(
        description="Validate, prepare, or score WorkPilot office-content tasks"
    )
    subparsers = parser.add_subparsers(dest="command", required=True)
    validate_parser = subparsers.add_parser("validate")
    validate_parser.add_argument("--suite", type=Path, default=DEFAULT_SUITE)
    prepare_parser = subparsers.add_parser("prepare")
    prepare_parser.add_argument("--suite", type=Path, default=DEFAULT_SUITE)
    prepare_parser.add_argument("--workspace-root", type=Path, required=True)
    score_parser = subparsers.add_parser("score")
    score_parser.add_argument("--suite", type=Path, default=DEFAULT_SUITE)
    score_parser.add_argument("--submission-root", type=Path, required=True)
    score_parser.add_argument("--output-dir", type=Path, required=True)
    score_parser.add_argument("--reviews", type=Path)
    score_parser.add_argument(
        "--require-complete-reviews",
        action="store_true",
        help="发布门禁：所有选中任务必须完成当前文件哈希绑定的复核并通过最低分",
    )
    for command_parser in (prepare_parser, score_parser):
        command_parser.add_argument(
            "--split", choices=("dev", "test", "all"), default="dev"
        )
        command_parser.add_argument("--include-test", action="store_true")
        command_parser.add_argument("--test-access-note")
    args = parser.parse_args(argv)
    try:
        suite = load_suite(args.suite)
        if args.command == "validate":
            print(
                json.dumps(
                    suite_summary(suite), ensure_ascii=False, indent=2, sort_keys=True
                )
            )
            return 0
        common = {
            "split": args.split,
            "include_test": args.include_test,
            "test_access_note": args.test_access_note,
        }
        if args.command == "prepare":
            result = prepare_suite(suite, args.workspace_root, **common)
            print(
                json.dumps(
                    {
                        "prepared": len(result["items"]),
                        "workspace_root": str(args.workspace_root),
                    },
                    ensure_ascii=False,
                    indent=2,
                    sort_keys=True,
                )
            )
            return 0
        report = evaluate_suite(
            suite,
            args.submission_root,
            args.output_dir,
            reviews=load_reviews(args.reviews),
            **common,
        )
    except OfficeContentSuiteError as error:
        print(str(error), file=sys.stderr)
        return 2
    print(json.dumps(report["summary"], ensure_ascii=False, indent=2, sort_keys=True))
    if args.require_complete_reviews:
        return (
            0
            if report["summary"]["review_complete_items"] == report["summary"]["items"]
            and report["summary"]["benchmark_passed"] == report["summary"]["items"]
            else 1
        )
    return (
        0 if report["summary"]["automatic_passed"] == report["summary"]["items"] else 1
    )


if __name__ == "__main__":  # pragma: no cover
    raise SystemExit(_main())


__all__ = [
    "DEFAULT_SUITE",
    "REVIEW_SCHEMA_VERSION",
    "ArtifactView",
    "CheckOutcome",
    "OfficeContentItem",
    "OfficeContentSuite",
    "OfficeContentSuiteError",
    "PptxChartDataCheck",
    "ReviewAnnotation",
    "ReviewFile",
    "ScopeSpec",
    "TextClaimContainsCheck",
    "TextClaimNotContainsCheck",
    "TextContainsCheck",
    "TextLabeledValuesCheck",
    "XlsxCellsValueCheck",
    "XlsxFormulaValueCheck",
    "evaluate_check",
    "evaluate_suite",
    "extract_artifact",
    "load_reviews",
    "load_suite",
    "prepare_suite",
    "suite_summary",
]
